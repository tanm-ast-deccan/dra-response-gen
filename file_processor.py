"""
file_processor.py — Multi-format document ingestion for the Claude adapter.

Handles: PDF, DOCX, XLSX, PPTX, CSV, TXT, MD, HTML, JPEG, PNG, GIF, WEBP

Architecture:
    ┌──────────────────────────────────────────────────────────┐
    │                    FileProcessor                         │
    │                                                          │
    │  Input: list of file paths (any supported format)        │
    │                                                          │
    │  For each file:                                          │
    │    1. Detect format from extension                       │
    │    2. Validate (exists, size < limit, supported)         │
    │    3. Extract/convert to Claude API content block:       │
    │       - PDF    → {"type": "document", base64}  (native) │
    │       - Images → {"type": "image", base64}     (native) │
    │       - Others → {"type": "text", extracted}   (parsed) │
    │    4. Apply truncation if content exceeds token budget   │
    │                                                          │
    │  Output: list of Claude API content blocks               │
    │          ready to append to the user message             │
    └──────────────────────────────────────────────────────────┘

Why three strategies?
    - PDF: Claude reads PDFs natively via the Documents API. Sending as
      base64 preserves tables, charts, and layout that text extraction
      would lose. This is the gold standard for fidelity.
    - Images: Claude's vision capability reads images natively. Sending
      as base64 lets the model see charts, diagrams, screenshots, etc.
    - Text extraction: For DOCX/XLSX/PPTX/CSV, there's no native API
      support, so we extract text content. This loses formatting but
      preserves the actual data. For evaluation purposes, this is
      sufficient since the agent needs the *information*, not the layout.

Token budget management:
    Claude's context window is 200K tokens (standard). With 10 files at
    5MB each, naive inclusion would blow past this limit. The processor
    applies per-file and total token budgets, truncating with clear
    markers so the agent knows content was cut.
"""

from __future__ import annotations

import os
import io
import csv
import base64
import mimetypes
from dataclasses import dataclass, field
from pathlib import Path
from typing import Optional


# ─── Configuration ───────────────────────────────────────────────────────

MAX_FILE_SIZE_BYTES = 5 * 1024 * 1024       # 5 MB per file
MAX_TOTAL_SIZE_BYTES = 50 * 1024 * 1024     # 50 MB total across all files
MAX_FILES = 10

# Token budgets for text-extracted content
# Rule of thumb: 1 token ≈ 4 characters for English text
CHARS_PER_TOKEN = 4
DEFAULT_PER_FILE_TOKEN_BUDGET = 15_000      # ~60K chars, enough for most docs
DEFAULT_TOTAL_TOKEN_BUDGET = 100_000        # ~400K chars across all files

# For native formats (PDF, images), we send the raw binary.
# These have their own API-side limits but are generally more efficient
# than text extraction since Claude processes them natively.
MAX_IMAGE_DIMENSION = 2048  # resize images larger than this
JPEG_QUALITY = 85           # quality for JPEG re-encoding

# Supported formats grouped by processing strategy
NATIVE_DOCUMENT_EXTS = {".pdf"}
NATIVE_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".webp"}
TEXT_EXTRACT_EXTS = {".docx", ".xlsx", ".xls", ".pptx", ".csv", ".tsv",
                     ".txt", ".md", ".html", ".htm", ".json", ".xml"}

ALL_SUPPORTED_EXTS = NATIVE_DOCUMENT_EXTS | NATIVE_IMAGE_EXTS | TEXT_EXTRACT_EXTS


# ─── Data structures ────────────────────────────────────────────────────

@dataclass
class ProcessedFile:
    """
    Result of processing a single file.
    
    Contains the Claude API content block(s) and metadata for logging.
    """
    filename: str
    filepath: str
    format: str                              # "pdf", "image", "text"
    content_blocks: list[dict]               # Claude API content blocks
    original_size_bytes: int
    extracted_chars: int = 0                  # for text-extracted files
    truncated: bool = False
    error: Optional[str] = None


@dataclass
class ProcessingResult:
    """Result of processing all files for a task."""
    files: list[ProcessedFile] = field(default_factory=list)
    content_blocks: list[dict] = field(default_factory=list)  # all blocks, ready for API
    total_original_bytes: int = 0
    total_extracted_chars: int = 0
    errors: list[str] = field(default_factory=list)
    skipped: list[str] = field(default_factory=list)


# ─── The processor ───────────────────────────────────────────────────────

class FileProcessor:
    """
    Processes multiple files into Claude API content blocks.
    
    Usage:
        processor = FileProcessor()
        result = processor.process_files([
            "/path/to/report.pdf",
            "/path/to/data.xlsx",
            "/path/to/chart.png",
        ])
        
        # result.content_blocks is a list ready to append to the
        # user message content array in the Claude Messages API.
    """
    
    def __init__(
        self,
        per_file_token_budget: int = DEFAULT_PER_FILE_TOKEN_BUDGET,
        total_token_budget: int = DEFAULT_TOTAL_TOKEN_BUDGET,
        max_file_size: int = MAX_FILE_SIZE_BYTES,
        max_files: int = MAX_FILES,
    ):
        self.per_file_token_budget = per_file_token_budget
        self.total_token_budget = total_token_budget
        self.max_file_size = max_file_size
        self.max_files = max_files
        self._total_chars_used = 0
    
    def process_files(self, file_paths: list[str]) -> ProcessingResult:
        """
        Process all files and return Claude API content blocks.
        
        Files are processed in order. If the total token budget is
        exhausted, remaining files are truncated more aggressively
        or skipped entirely.
        """
        result = ProcessingResult()
        self._total_chars_used = 0
        
        # ── Validate count ───────────────────────────────────────
        if len(file_paths) > self.max_files:
            result.errors.append(
                f"Too many files: {len(file_paths)} provided, max is {self.max_files}. "
                f"Processing first {self.max_files} only."
            )
            file_paths = file_paths[:self.max_files]
        
        # ── Check total size upfront ─────────────────────────────
        total_size = 0
        for fpath in file_paths:
            if os.path.exists(fpath):
                total_size += os.path.getsize(fpath)
        
        if total_size > MAX_TOTAL_SIZE_BYTES:
            result.errors.append(
                f"Total file size ({total_size / 1024 / 1024:.1f} MB) exceeds "
                f"limit ({MAX_TOTAL_SIZE_BYTES / 1024 / 1024:.0f} MB)."
            )
        
        result.total_original_bytes = total_size
        
        # ── Process each file ────────────────────────────────────
        for fpath in file_paths:
            processed = self._process_single_file(fpath)
            result.files.append(processed)
            
            if processed.error:
                result.errors.append(f"{processed.filename}: {processed.error}")
                continue
            
            result.content_blocks.extend(processed.content_blocks)
            result.total_extracted_chars += processed.extracted_chars
        
        # ── Add file manifest ────────────────────────────────────
        # Prepend a text block listing all files so Claude knows
        # what documents are available and can reference them.
        manifest = self._build_manifest(result)
        result.content_blocks.insert(0, {
            "type": "text",
            "text": manifest,
        })
        
        return result
    
    def _process_single_file(self, filepath: str) -> ProcessedFile:
        """Route a file to the appropriate processor based on extension."""
        
        filename = os.path.basename(filepath)
        ext = Path(filepath).suffix.lower()
        
        # ── Validate ─────────────────────────────────────────────
        if not os.path.exists(filepath):
            return ProcessedFile(
                filename=filename, filepath=filepath, format="unknown",
                content_blocks=[], original_size_bytes=0,
                error=f"File not found: {filepath}",
            )
        
        file_size = os.path.getsize(filepath)
        
        if file_size > self.max_file_size:
            return ProcessedFile(
                filename=filename, filepath=filepath, format="unknown",
                content_blocks=[], original_size_bytes=file_size,
                error=f"File too large: {file_size / 1024 / 1024:.1f} MB "
                      f"(max {self.max_file_size / 1024 / 1024:.0f} MB)",
            )
        
        if ext not in ALL_SUPPORTED_EXTS:
            return ProcessedFile(
                filename=filename, filepath=filepath, format="unknown",
                content_blocks=[], original_size_bytes=file_size,
                error=f"Unsupported format: {ext}",
            )
        
        # ── Route to processor ───────────────────────────────────
        try:
            if ext in NATIVE_DOCUMENT_EXTS:
                return self._process_pdf(filepath, filename, file_size)
            elif ext in NATIVE_IMAGE_EXTS:
                return self._process_image(filepath, filename, file_size, ext)
            elif ext == ".docx":
                return self._process_docx(filepath, filename, file_size)
            elif ext in (".xlsx", ".xls"):
                return self._process_xlsx(filepath, filename, file_size)
            elif ext == ".pptx":
                return self._process_pptx(filepath, filename, file_size)
            elif ext in (".csv", ".tsv"):
                return self._process_csv(filepath, filename, file_size, ext)
            else:
                # Plain text formats: txt, md, html, json, xml
                return self._process_text(filepath, filename, file_size)
        except Exception as e:
            return ProcessedFile(
                filename=filename, filepath=filepath, format="unknown",
                content_blocks=[], original_size_bytes=file_size,
                error=f"Processing error: {type(e).__name__}: {e}",
            )
    
    # ── Native PDF ───────────────────────────────────────────────────
    
    def _process_pdf(self, filepath, filename, file_size) -> ProcessedFile:
        """
        Send PDF as a native document block.
        
        Claude reads PDFs natively — tables, charts, headers, footnotes,
        and layout are all preserved. This is strictly better than text
        extraction for any PDF with non-trivial formatting.
        
        API format:
            {
                "type": "document",
                "source": {
                    "type": "base64",
                    "media_type": "application/pdf",
                    "data": "<base64 string>"
                }
            }
        """
        with open(filepath, "rb") as f:
            raw = f.read()
        
        b64 = base64.standard_b64encode(raw).decode("ascii")
        
        content_block = {
            "type": "document",
            "source": {
                "type": "base64",
                "media_type": "application/pdf",
                "data": b64,
            },
            # Cache control: PDFs are static, so marking as ephemeral
            # is appropriate. For very large PDFs, prompt caching helps.
            "cache_control": {"type": "ephemeral"},
        }
        
        # Add a label so Claude can reference this file by name
        label_block = {
            "type": "text",
            "text": f"\n[Document: {filename} — PDF, {file_size / 1024:.0f} KB]\n",
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="pdf",
            content_blocks=[label_block, content_block],
            original_size_bytes=file_size,
        )
    
    # ── Native images ────────────────────────────────────────────────
    
    def _process_image(self, filepath, filename, file_size, ext) -> ProcessedFile:
        """
        Send image as a native image block.
        
        Claude's vision reads images directly. We resize if necessary
        to stay within reasonable token costs (larger images = more tokens).
        
        API format:
            {
                "type": "image",
                "source": {
                    "type": "base64",
                    "media_type": "image/jpeg",
                    "data": "<base64 string>"
                }
            }
        """
        from PIL import Image
        
        # Map extensions to MIME types
        mime_map = {
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".png": "image/png",
            ".gif": "image/gif",
            ".webp": "image/webp",
        }
        media_type = mime_map.get(ext, "image/jpeg")
        
        # Open and optionally resize
        img = Image.open(filepath)
        original_dims = img.size
        resized = False
        
        if max(img.size) > MAX_IMAGE_DIMENSION:
            img.thumbnail((MAX_IMAGE_DIMENSION, MAX_IMAGE_DIMENSION), Image.LANCZOS)
            resized = True
        
        # Encode to bytes
        buf = io.BytesIO()
        if ext in (".jpg", ".jpeg"):
            # Convert RGBA to RGB for JPEG
            if img.mode in ("RGBA", "P"):
                img = img.convert("RGB")
            img.save(buf, format="JPEG", quality=JPEG_QUALITY)
        elif ext == ".png":
            img.save(buf, format="PNG")
        elif ext == ".gif":
            img.save(buf, format="GIF")
        elif ext == ".webp":
            img.save(buf, format="WEBP")
        else:
            img.save(buf, format="JPEG", quality=JPEG_QUALITY)
            media_type = "image/jpeg"
        
        b64 = base64.standard_b64encode(buf.getvalue()).decode("ascii")
        
        content_block = {
            "type": "image",
            "source": {
                "type": "base64",
                "media_type": media_type,
                "data": b64,
            },
        }
        
        resize_note = f", resized from {original_dims}" if resized else ""
        label_block = {
            "type": "text",
            "text": (f"\n[Image: {filename} — {img.size[0]}x{img.size[1]}"
                     f"{resize_note}]\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="image",
            content_blocks=[label_block, content_block],
            original_size_bytes=file_size,
        )
    
    # ── DOCX extraction ──────────────────────────────────────────────
    
    def _process_docx(self, filepath, filename, file_size) -> ProcessedFile:
        """
        Extract text from Word documents.
        
        Extracts:
          - All paragraph text (preserving newlines)
          - Table content (formatted as pipe-delimited)
          - Header/footer text (noted separately)
        
        Does NOT extract:
          - Embedded images (would need separate handling)
          - Track changes (gets the current text state)
          - Comments (could be added if needed)
        """
        import docx
        
        doc = docx.Document(filepath)
        parts = []
        
        for element in doc.element.body:
            tag = element.tag.split("}")[-1]  # strip namespace
            
            if tag == "p":
                # Paragraph
                para = docx.text.paragraph.Paragraph(element, doc)
                text = para.text.strip()
                if text:
                    parts.append(text)
            
            elif tag == "tbl":
                # Table — format as pipe-delimited
                table = docx.table.Table(element, doc)
                table_lines = []
                for row in table.rows:
                    cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    table_lines.append("| " + " | ".join(cells) + " |")
                
                if table_lines:
                    # Add separator after header row
                    if len(table_lines) > 1:
                        n_cols = len(table.rows[0].cells)
                        separator = "| " + " | ".join(["---"] * n_cols) + " |"
                        table_lines.insert(1, separator)
                    parts.append("\n".join(table_lines))
        
        full_text = "\n\n".join(parts)
        
        # Apply token budget
        full_text, truncated = self._apply_token_budget(full_text, filename)
        
        content_block = {
            "type": "text",
            "text": (f"\n--- Document: {filename} (Word) ---\n"
                     f"{full_text}\n"
                     f"--- End: {filename} ---\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="text",
            content_blocks=[content_block],
            original_size_bytes=file_size,
            extracted_chars=len(full_text),
            truncated=truncated,
        )
    
    # ── XLSX extraction ──────────────────────────────────────────────
    
    def _process_xlsx(self, filepath, filename, file_size) -> ProcessedFile:
        """
        Extract data from Excel workbooks.
        
        Strategy:
          - Process each sheet separately
          - Format as pipe-delimited tables (Markdown-style)
          - Include sheet names as headers
          - For large sheets, truncate rows (keep headers + first N rows)
        
        This captures the DATA, not the formatting. Formulas are evaluated
        to their current values (via data_only=True).
        """
        import openpyxl
        
        wb = openpyxl.load_workbook(filepath, data_only=True, read_only=True)
        parts = []
        
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows_data = []
            
            for row in ws.iter_rows(values_only=True):
                # Convert each cell to string, handle None
                cells = [str(c).strip() if c is not None else "" for c in row]
                # Skip entirely empty rows
                if any(cells):
                    rows_data.append(cells)
            
            if not rows_data:
                continue
            
            # Build table
            sheet_lines = [f"### Sheet: {sheet_name}"]
            
            # Normalize column count (some rows may have different lengths)
            max_cols = max(len(r) for r in rows_data)
            for row in rows_data:
                while len(row) < max_cols:
                    row.append("")
            
            # Format as pipe-delimited table
            for i, row in enumerate(rows_data):
                sheet_lines.append("| " + " | ".join(row) + " |")
                if i == 0:
                    sheet_lines.append("| " + " | ".join(["---"] * max_cols) + " |")
            
            sheet_lines.append(f"({len(rows_data) - 1} data rows)")
            parts.append("\n".join(sheet_lines))
        
        wb.close()
        
        full_text = "\n\n".join(parts)
        full_text, truncated = self._apply_token_budget(full_text, filename)
        
        content_block = {
            "type": "text",
            "text": (f"\n--- Document: {filename} (Excel) ---\n"
                     f"{full_text}\n"
                     f"--- End: {filename} ---\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="text",
            content_blocks=[content_block],
            original_size_bytes=file_size,
            extracted_chars=len(full_text),
            truncated=truncated,
        )
    
    # ── PPTX extraction ──────────────────────────────────────────────
    
    def _process_pptx(self, filepath, filename, file_size) -> ProcessedFile:
        """
        Extract text from PowerPoint presentations.
        
        Extracts:
          - Slide titles and body text
          - Table content within slides
          - Speaker notes (valuable context often missed)
        
        Format: Each slide labeled with number + title.
        """
        from pptx import Presentation
        from pptx.util import Inches
        
        prs = Presentation(filepath)
        parts = []
        
        for i, slide in enumerate(prs.slides, 1):
            slide_parts = [f"### Slide {i}"]
            
            for shape in slide.shapes:
                # Text frames (titles, body text, text boxes)
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_parts.append(text)
                
                # Tables within slides
                if shape.has_table:
                    table = shape.table
                    table_lines = []
                    for row in table.rows:
                        cells = [cell.text.strip().replace("\n", " ")
                                 for cell in row.cells]
                        table_lines.append("| " + " | ".join(cells) + " |")
                    
                    if table_lines and len(table.rows) > 1:
                        n_cols = len(table.rows[0].cells)
                        separator = "| " + " | ".join(["---"] * n_cols) + " |"
                        table_lines.insert(1, separator)
                    
                    slide_parts.append("\n".join(table_lines))
            
            # Speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_parts.append(f"[Speaker notes: {notes_text}]")
            
            # Only include slide if it has content beyond just the header
            if len(slide_parts) > 1:
                parts.append("\n".join(slide_parts))
        
        full_text = "\n\n".join(parts)
        full_text, truncated = self._apply_token_budget(full_text, filename)
        
        content_block = {
            "type": "text",
            "text": (f"\n--- Document: {filename} (PowerPoint, "
                     f"{len(prs.slides)} slides) ---\n"
                     f"{full_text}\n"
                     f"--- End: {filename} ---\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="text",
            content_blocks=[content_block],
            original_size_bytes=file_size,
            extracted_chars=len(full_text),
            truncated=truncated,
        )
    
    # ── CSV/TSV extraction ───────────────────────────────────────────
    
    def _process_csv(self, filepath, filename, file_size, ext) -> ProcessedFile:
        """
        Read CSV/TSV files as pipe-delimited tables.
        
        For large CSVs (common in evaluation tasks), we:
          - Always preserve the header row
          - Show first N rows that fit within token budget
          - Report total row count so the agent knows what was truncated
        """
        delimiter = "\t" if ext == ".tsv" else ","
        
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            # Use csv.reader for robust parsing (handles quoted fields, etc.)
            reader = csv.reader(f, delimiter=delimiter)
            rows = []
            for row in reader:
                rows.append(row)
                # Safety limit: don't read more than 50K rows into memory
                if len(rows) > 50_000:
                    break
        
        total_rows = len(rows)
        
        if not rows:
            return ProcessedFile(
                filename=filename, filepath=filepath, format="text",
                content_blocks=[{
                    "type": "text",
                    "text": f"\n[File: {filename} — empty CSV/TSV]\n",
                }],
                original_size_bytes=file_size,
            )
        
        # Normalize column count
        max_cols = max(len(r) for r in rows)
        for row in rows:
            while len(row) < max_cols:
                row.append("")
        
        # Build pipe-delimited table
        lines = []
        for i, row in enumerate(rows):
            cells = [c.strip().replace("\n", " ").replace("|", "\\|") for c in row]
            lines.append("| " + " | ".join(cells) + " |")
            if i == 0:
                lines.append("| " + " | ".join(["---"] * max_cols) + " |")
        
        full_text = "\n".join(lines)
        full_text, truncated = self._apply_token_budget(full_text, filename)
        
        row_note = f" ({total_rows - 1} data rows)"
        if total_rows > 50_000:
            row_note = f" (50,000+ rows, showing subset)"
        
        content_block = {
            "type": "text",
            "text": (f"\n--- Document: {filename} (CSV{row_note}) ---\n"
                     f"{full_text}\n"
                     f"--- End: {filename} ---\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="text",
            content_blocks=[content_block],
            original_size_bytes=file_size,
            extracted_chars=len(full_text),
            truncated=truncated,
        )
    
    # ── Plain text extraction ────────────────────────────────────────
    
    def _process_text(self, filepath, filename, file_size) -> ProcessedFile:
        """
        Read plain text files (TXT, MD, HTML, JSON, XML).
        
        These are sent as-is (with truncation if needed). The file
        extension is noted so Claude can interpret the format.
        """
        ext = Path(filepath).suffix.lower()
        format_names = {
            ".txt": "Text", ".md": "Markdown", ".html": "HTML",
            ".htm": "HTML", ".json": "JSON", ".xml": "XML",
        }
        format_name = format_names.get(ext, "Text")
        
        with open(filepath, "r", encoding="utf-8", errors="replace") as f:
            full_text = f.read()
        
        full_text, truncated = self._apply_token_budget(full_text, filename)
        
        content_block = {
            "type": "text",
            "text": (f"\n--- Document: {filename} ({format_name}) ---\n"
                     f"{full_text}\n"
                     f"--- End: {filename} ---\n"),
        }
        
        return ProcessedFile(
            filename=filename, filepath=filepath, format="text",
            content_blocks=[content_block],
            original_size_bytes=file_size,
            extracted_chars=len(full_text),
            truncated=truncated,
        )
    
    # ── Token budget management ──────────────────────────────────────
    
    def _apply_token_budget(self, text: str, filename: str) -> tuple[str, bool]:
        """
        Truncate text to fit within per-file and total token budgets.
        
        Strategy: Keep the beginning and end of the document (head + tail),
        with a clear truncation marker in the middle. This preserves:
          - Headers, table of contents, executive summary (usually at top)
          - Conclusions, references, appendices (usually at bottom)
          - A clear signal to Claude that content was omitted
        
        Returns: (possibly_truncated_text, was_truncated)
        """
        # Calculate remaining budget from total pool
        remaining_total = self.total_token_budget - (self._total_chars_used // CHARS_PER_TOKEN)
        effective_budget = min(self.per_file_token_budget, max(remaining_total, 1000))
        
        max_chars = effective_budget * CHARS_PER_TOKEN
        
        if len(text) <= max_chars:
            self._total_chars_used += len(text)
            return text, False
        
        # Truncate: keep 70% from head, 30% from tail
        head_chars = int(max_chars * 0.7)
        tail_chars = int(max_chars * 0.3) - 200  # reserve space for marker
        
        truncation_marker = (
            f"\n\n[... TRUNCATED: {filename} — showing first ~{head_chars // 1000}K "
            f"and last ~{tail_chars // 1000}K characters of "
            f"{len(text) // 1000}K total. "
            f"{(len(text) - head_chars - tail_chars) // 1000}K characters omitted. ...]\n\n"
        )
        
        truncated_text = text[:head_chars] + truncation_marker + text[-tail_chars:]
        
        self._total_chars_used += len(truncated_text)
        return truncated_text, True
    
    # ── Manifest ─────────────────────────────────────────────────────
    
    def _build_manifest(self, result: ProcessingResult) -> str:
        """
        Build a text manifest listing all provided files.
        
        This goes at the top of the content blocks so Claude
        immediately knows what documents are available and can
        plan its research strategy accordingly.
        """
        lines = [
            "The following documents have been provided for this research task:",
            "",
        ]
        
        for f in result.files:
            if f.error:
                lines.append(f"  ❌ {f.filename} — ERROR: {f.error}")
            else:
                size_kb = f.original_size_bytes / 1024
                trunc = " [TRUNCATED]" if f.truncated else ""
                lines.append(f"  • {f.filename} ({f.format}, {size_kb:.0f} KB){trunc}")
        
        if result.errors:
            lines.append("")
            lines.append(f"Note: {len(result.errors)} file(s) had processing issues.")
        
        lines.append("")
        lines.append(
            "Reference these files using [File: filename] in your citations. "
            "Treat file contents as primary source material."
        )
        
        return "\n".join(lines)


# ─── Convenience function ────────────────────────────────────────────────

def process_files(
    file_paths: list[str],
    per_file_token_budget: int = DEFAULT_PER_FILE_TOKEN_BUDGET,
    total_token_budget: int = DEFAULT_TOTAL_TOKEN_BUDGET,
) -> ProcessingResult:
    """
    One-shot convenience function.
    
    Usage:
        from file_processor import process_files
        result = process_files(["report.pdf", "data.xlsx", "chart.png"])
        content_blocks = result.content_blocks
    """
    processor = FileProcessor(
        per_file_token_budget=per_file_token_budget,
        total_token_budget=total_token_budget,
    )
    return processor.process_files(file_paths)
