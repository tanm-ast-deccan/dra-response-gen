# src/document_parser.py
import os
import warnings
from pathlib import Path

# Suppress PyMuPDF's informational suggestion about pymupdf_layout
warnings.filterwarnings("ignore", message=".*pymupdf_layout.*")


def read_document(file_path: str) -> str:
    """
    Reads and extracts text from a document.
    Supports .pdf, .docx, .xlsx, .xls, .csv, and .txt/.md files.

    Args:
        file_path: Path to the document.

    Returns:
        Extracted text content as a string.

    Raises:
        FileNotFoundError: If the file doesn't exist.
        NotImplementedError: If the file type is not supported.
    """
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(f"Document not found: {file_path}")

    extension = path.suffix.lower()

    if extension == ".txt" or extension == ".md":
        return _read_txt(path)
    elif extension == ".pdf":
        return _read_pdf(path)
    elif extension == ".docx":
        return _read_docx(path)
    elif extension in (".xlsx", ".xls", ".xlsm"):
        return _read_xlsx(path)
    elif extension in (".csv", ".tsv"):
        return _read_csv(path)
    elif extension == ".pptx":
        return _read_pptx(path)
    elif extension == ".json":
        return _read_json(path)
    else:
        raise NotImplementedError(
            f"File type '{extension}' is not supported. "
            f"Use .pdf, .docx, .xlsx, .xls, .csv, .txt, or .md"
        )

def _read_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"### Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs)
                    if t.strip(): parts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    parts.append("| " + " | ".join(c.text.strip() for c in row.cells) + " |")
    return "\n".join(parts)

def _read_json(path):
    import json
    with open(path, encoding="utf-8", errors="replace") as f:
        return json.dumps(json.load(f), indent=2, ensure_ascii=False)
    
def _read_txt(path: Path) -> str:
    """Read plain text file."""
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        return f.read()


def _read_pdf(path: Path) -> str:
    """Read PDF file using PyMuPDF (fitz) with table extraction."""
    try:
        import fitz  # PyMuPDF
    except ImportError:
        raise ImportError(
            "PyMuPDF is required to read PDF files. Install with: pip install PyMuPDF"
        )

    doc = fitz.open(path)
    text_parts = []

    for page in doc:
        try:
            tables = page.find_tables()
            if tables and len(tables.tables) > 0:
                page_text_parts = []
                blocks = page.get_text("blocks")
                table_rects = [t.bbox for t in tables.tables]

                for block in blocks:
                    block_rect = fitz.Rect(block[:4])
                    in_table = any(
                        block_rect.intersects(fitz.Rect(tr)) for tr in table_rects
                    )
                    if not in_table and block[4].strip():
                        page_text_parts.append(block[4].strip())

                for table in tables.tables:
                    md_table = _convert_pdf_table_to_markdown(table)
                    if md_table:
                        page_text_parts.append(md_table)

                text_parts.append("\n\n".join(page_text_parts))
            else:
                text_parts.append(page.get_text())
        except AttributeError:
            text_parts.append(page.get_text())

    doc.close()
    return "\n\n".join(text_parts)


def _convert_pdf_table_to_markdown(table) -> str:
    """Convert PyMuPDF table to markdown format."""
    try:
        data = table.extract()
        if not data or len(data) == 0:
            return ""

        rows_text = []
        for i, row in enumerate(data):
            cells = [
                str(cell).strip().replace("\n", " ") if cell else "" for cell in row
            ]
            row_text = " | ".join(cells)
            rows_text.append(f"| {row_text} |")
            if i == 0:
                separator = "|" + "|".join(["---"] * len(cells)) + "|"
                rows_text.append(separator)

        return "\n".join(rows_text)
    except Exception:
        return ""


def _read_docx(path: Path) -> str:
    """Read Word document using python-docx, including tables and text boxes."""
    try:
        from docx import Document
        from docx.oxml.ns import qn
    except ImportError:
        raise ImportError(
            "python-docx is required to read DOCX files. Install with: pip install python-docx"
        )

    doc = Document(path)
    text_parts = []

    try:
        for element in doc.element.body:
            if element.tag.endswith("p"):
                for para in doc.paragraphs:
                    if para._element is element:
                        if para.text.strip():
                            text_parts.append(para.text)
                        textbox_content = _extract_textboxes_from_element(element)
                        if textbox_content:
                            text_parts.append(textbox_content)
                        break
            elif element.tag.endswith("tbl"):
                for table in doc.tables:
                    if table._tbl is element:
                        table_text = _extract_table_as_text(table)
                        if table_text.strip():
                            text_parts.append(table_text)
                        break
    except Exception as e:
        print(f"  Warning: Structured parsing failed ({e}), using fallback...")
        text_parts = []
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        for table in doc.tables:
            table_text = _extract_table_as_text(table)
            if table_text.strip():
                text_parts.append(table_text)

    return "\n\n".join(text_parts)


def _extract_textboxes_from_element(element) -> str:
    """Extract text from text boxes within an element."""
    try:
        from docx.oxml.ns import qn
        textbox_texts = []
        for txbx in element.iter(qn("w:txbxContent")):
            for p in txbx.iter(qn("w:p")):
                text = "".join(node.text for node in p.iter(qn("w:t")) if node.text)
                if text.strip():
                    textbox_texts.append(text.strip())
        return "\n".join(textbox_texts)
    except Exception:
        return ""


def _extract_table_as_text(table) -> str:
    """Extract table content as markdown-style text."""
    try:
        rows_text = []
        prev_row_cells = None

        for i, row in enumerate(table.rows):
            cells = []
            for j, cell in enumerate(row.cells):
                cell_text = cell.text.strip().replace("\n", " ")
                if prev_row_cells and j < len(prev_row_cells):
                    if cell._tc is prev_row_cells[j]._tc:
                        cell_text = ""
                cells.append(cell_text)

            row_text = " | ".join(cells)
            rows_text.append(f"| {row_text} |")
            if i == 0:
                separator = "|" + "|".join(["---"] * len(cells)) + "|"
                rows_text.append(separator)
            prev_row_cells = row.cells

        return "\n".join(rows_text)
    except Exception:
        try:
            rows_text = []
            for row in table.rows:
                cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                rows_text.append("| " + " | ".join(cells) + " |")
            return "\n".join(rows_text)
        except Exception:
            return ""


def _read_xlsx(path: Path) -> str:
    """
    Read Excel file (.xlsx/.xls/.xlsm) using openpyxl.

    Extracts all sheets as markdown tables, prefixed with the sheet name.
    Each sheet is separated by a clear header so the LLM can identify
    which sheet data came from.
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError(
            "openpyxl is required to read Excel files. Install with: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(path, data_only=True)
    sheet_parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]

        # Skip completely empty sheets
        if ws.max_row == 0 or ws.max_column == 0:
            continue

        rows_text = []
        max_col = ws.max_column

        for i, row in enumerate(ws.iter_rows(max_row=ws.max_row, values_only=True)):
            # Skip rows that are entirely None/empty
            if all(cell is None or str(cell).strip() == "" for cell in row):
                continue

            cells = [
                str(cell).strip().replace("\n", " ") if cell is not None else ""
                for cell in row
            ]
            # Pad to max_col
            while len(cells) < max_col:
                cells.append("")

            rows_text.append("| " + " | ".join(cells) + " |")

            # Separator after first (header) row
            if i == 0:
                rows_text.append("|" + "|".join(["---"] * max_col) + "|")

        if rows_text:
            sheet_parts.append(
                f"### Sheet: {sheet_name}\n" + "\n".join(rows_text)
            )

    wb.close()

    if not sheet_parts:
        return "(Excel file appears to be empty)"

    return "\n\n".join(sheet_parts)


def _read_csv(path: Path) -> str:
    """Read CSV/TSV file as a markdown table."""
    import csv

    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f, delimiter=delimiter)
        rows = list(reader)

    if not rows:
        return "(CSV file is empty)"

    max_cols = max(len(r) for r in rows)
    lines = []
    for i, row in enumerate(rows):
        while len(row) < max_cols:
            row.append("")
        cells = [c.strip().replace("\n", " ") for c in row]
        lines.append("| " + " | ".join(cells) + " |")
        if i == 0:
            lines.append("|" + "|".join(["---"] * max_cols) + "|")

    return "\n".join(lines)