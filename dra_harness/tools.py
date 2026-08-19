"""
tools.py — Full-featured local tool registry for the agentic loop.

Provides APEX-Agents-equivalent tool capabilities:
  - python_execute: run Python code in a subprocess
  - bash_execute: run shell commands
  - read_file: read file contents (text, CSV, XLSX, PDF)
  - write_file: write content to a file
  - list_directory: list files in a directory
  - web_fetch: fetch content from a URL
  - calculator: safe arithmetic evaluation

Web search is handled natively by OpenRouter (see GenParams.web_search /
provider.py) and is NOT a local tool.

Enable tools per run via GenParams.enabled_tools = ["all"] for everything,
or list specific names. Register your own with @register_tool.
"""

from __future__ import annotations

import json
import os
import logging
from dataclasses import dataclass
from typing import Callable


#: Default ceiling on a single read_file result.
#:
#: This was 0, meaning "lossless, never truncate", with the docstring telling the
#: model to bound only a "pathological" input. A 21.8 MB annual report is not
#: pathological — it is the normal case in finance, and it extracts to roughly
#: 3,000,000 characters (~750,000 tokens), which no model accepts in one tool
#: result. The call would fail or cost a fortune, and the model had no hint that
#: another route existed.
#:
#: The ceiling is deliberately generous: most inputs fit under it untouched, and
#: when one does not the message names the alternative rather than just stopping.
#: Reading a 200-page report selectively — search, then read the pages that
#: matter — is the skill the benchmark exists to measure, so the cap must push
#: towards that and not remove it.
DEFAULT_READ_CHARS = 100_000

logger = logging.getLogger("dra.tools")


@dataclass
class Tool:
    name: str
    description: str
    parameters: dict          # JSON schema for the function arguments
    func: Callable            # func(**arguments) -> str (or JSON-serializable)

    def schema(self) -> dict:
        """OpenAI-compatible tool schema."""
        return {
            "type": "function",
            "function": {
                "name": self.name,
                "description": self.description,
                "parameters": self.parameters,
            },
        }

    def run(self, arguments: dict) -> str:
        result = self.func(**(arguments or {}))
        if isinstance(result, str):
            return result
        try:
            return json.dumps(result, ensure_ascii=False, default=str)
        except (TypeError, ValueError):
            return str(result)


TOOL_REGISTRY: dict[str, Tool] = {}


def register_tool(name: str, description: str, parameters: dict):
    """Decorator to register a local tool callable in TOOL_REGISTRY."""
    def deco(func: Callable) -> Callable:
        TOOL_REGISTRY[name] = Tool(name, description, parameters, func)
        return func
    return deco


def schemas_for(names: list[str]) -> list[dict]:
    """Return OpenAI tool schemas for the given enabled tool names.
    Pass ["all"] to enable every registered tool."""
    if names == ["all"]:
        return [t.schema() for t in TOOL_REGISTRY.values()]
    schemas = []
    for n in (names or []):
        tool = TOOL_REGISTRY.get(n)
        if tool is None:
            raise KeyError(
                f"Tool '{n}' is not registered. Available: {list(TOOL_REGISTRY)}"
            )
        schemas.append(tool.schema())
    return schemas


def execute(name: str, arguments: dict) -> str:
    """Execute a registered tool by name, returning a string result."""
    tool = TOOL_REGISTRY.get(name)
    if tool is None:
        return f"[tool error] unknown tool: {name}"
    try:
        return tool.run(arguments)
    except Exception as e:  # noqa: BLE001 — tool errors must not crash the loop
        return f"[tool error] {name} failed: {e}"


# ═══════════════════════════════════════════════════════════════════════════════
#  REGISTERED TOOLS
# ═══════════════════════════════════════════════════════════════════════════════


# ─── Python execution ─────────────────────────────────────────────────────────

@register_tool(
    name="python_execute",
    description=(
        "Execute Python code and return stdout/stderr. Use for calculations, "
        "data analysis, file processing, chart generation, or any programmatic "
        "task. Common libraries available: pandas, numpy, scipy, openpyxl, "
        "matplotlib, json, csv, os. The working directory is the task staging "
        "folder — input files are available there."
    ),
    parameters={
        "type": "object",
        "properties": {
            "code": {
                "type": "string",
                "description": "Python code to execute. Print results to stdout.",
            }
        },
        "required": ["code"],
    },
)
def _python_execute(code: str) -> str:
    import subprocess
    import tempfile

    staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())

    with tempfile.NamedTemporaryFile(
        mode="w", suffix=".py", dir=staging, delete=False
    ) as f:
        f.write(code)
        tmp_path = f.name

    try:
        result = subprocess.run(
            ["python", tmp_path],
            capture_output=True,
            text=True,
            timeout=120,
            cwd=staging,
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += f"\n[stderr]\n{result.stderr}"
        return (output.strip() or "[no output]")[:15000]
    except subprocess.TimeoutExpired:
        return "[error] execution timed out (120s)"
    except Exception as e:
        return f"[error] {e}"
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


# ─── Bash / shell execution ──────────────────────────────────────────────────

@register_tool(
    name="bash_execute",
    description=(
        "Execute a shell command and return stdout/stderr. Use for file "
        "operations, installing packages, running scripts, or system commands. "
        "The working directory is the task staging folder."
    ),
    parameters={
        "type": "object",
        "properties": {
            "command": {
                "type": "string",
                "description": "Shell command to execute.",
            }
        },
        "required": ["command"],
    },
)
def _bash_execute(command: str) -> str:
    import subprocess

    staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())

    try:
        result = subprocess.run(
            command,
            shell=True,
            capture_output=True,
            text=True,
            timeout=120,
            cwd=staging,
        )
        output = ""
        if result.stdout:
            output += result.stdout
        if result.stderr:
            output += f"\n[stderr]\n{result.stderr}"
        return (output.strip() or "[no output]")[:15000]
    except subprocess.TimeoutExpired:
        return "[error] execution timed out (120s)"
    except Exception as e:
        return f"[error] {e}"


# ─── Read file ────────────────────────────────────────────────────────────────

@register_tool(
    name="read_file",
    description=(
        "Read the contents of a file. Supports .txt, .csv, .json, .md, .py, "
        ".xlsx (all sheets), .pdf (text extraction), .docx. "
        f"Text content is capped at {DEFAULT_READ_CHARS:,} characters by default. "
        "When a file exceeds that, the result states how many characters remain "
        "and how to reach them - use bash_execute or python_execute to search the "
        "file and read only the region you need. Pass max_chars to raise the cap, "
        "or max_chars=0 for a lossless read of a file you know is small."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the file to read (absolute or relative to staging dir).",
            }
        },
        "required": ["path"],
    },
)



def _read_file(path: str, max_chars: int = -1) -> str:
    """Read a file to text for the model.

    Capped at DEFAULT_READ_CHARS by default. Pass max_chars=0 for a genuinely
    lossless read (only sensible on a file you already know is small), or a
    positive value to set your own ceiling.

    When the cap bites, the result says how much remains and how to reach it:
    python_execute with pdfplumber or openpyxl lets you search the file and read
    only the part you need.

    Supported: text formats, .json (parsed + pretty-printed), .md, .xlsx/.xls,
    .pdf, .docx, .pptx (slides + tables + speaker notes, notes clearly labeled).
    """
    import os

    if max_chars < 0:
        max_chars = DEFAULT_READ_CHARS

    staging = os.environ.get("INDRAYUDH_STAGING_DIR", os.getcwd())
    if not os.path.isabs(path):
        path = os.path.join(staging, path)
    if not os.path.isfile(path):
        return f"[error] file not found: {path}"

    ext = os.path.splitext(path)[1].lower()

    def cap(s: str) -> str:
        # max_chars <= 0 means lossless (no truncation).
        if not max_chars or len(s) <= max_chars:
            return s
        remaining = len(s) - max_chars
        hint = ("python_execute with openpyxl" if ext in (".xlsx", ".xls")
                else "python_execute with pdfplumber" if ext == ".pdf"
                else "python_execute" if ext in (".docx", ".pptx")
                else "bash_execute with grep, or python_execute")
        return (s[:max_chars]
                + f"\n\n[... TRUNCATED. This file extracts to {len(s):,} characters"
                  f" and you have been shown the first {max_chars:,};"
                  f" {remaining:,} characters remain, INCLUDING most of the"
                  f" document body.\n"
                  f" Do NOT assume a figure is absent because it is not above."
                  f" To reach the rest, use {hint} to search the file and read"
                  f" only the part you need — for example locate the page or"
                  f" sheet containing a label, then extract that region."
                  f" You may also call read_file again with a larger max_chars,"
                  f" but a whole file of this size will not fit in context.]")

    try:
        # ---- JSON: parse + pretty-print (structured, not a raw wall of text) ----
        if ext == ".json":
            import json
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                raw = f.read()
            try:
                obj = json.loads(raw)
                return cap(json.dumps(obj, indent=2, ensure_ascii=False))
            except json.JSONDecodeError:
                # not valid JSON — return raw so nothing is lost
                return cap(raw)

        # ---- Plain-text formats (now includes .md explicitly) ----
        if ext in (".txt", ".csv", ".md", ".markdown", ".py", ".js", ".html",
                   ".xml", ".yaml", ".yml", ".toml", ".log", ".tsv", ".sql", ".rst", ".ini"):
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                return cap(f.read())

        # ---- Excel ----
        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
            parts = []
            for sn in wb.sheetnames:
                ws = wb[sn]
                parts.append(f"=== Sheet: {sn} ===")
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join(str(c) if c is not None else "" for c in row))
            wb.close()
            return cap("\n".join(parts))

        # ---- PDF ----
        if ext == ".pdf":
            try:
                import pdfplumber
            except ImportError:
                return "[error] pdfplumber not installed — pip install pdfplumber"
            text_parts = []
            with pdfplumber.open(path) as pdf:
                for i, page in enumerate(pdf.pages, 1):
                    t = page.extract_text()
                    if t:
                        text_parts.append(f"=== Page {i} ===\n{t}")
                    # also pull tables, which extract_text can garble
                    for ti, table in enumerate(page.extract_tables() or [], 1):
                        rows = ["\t".join((c or "") for c in row) for row in table]
                        if rows:
                            text_parts.append(f"=== Page {i} Table {ti} ===\n" + "\n".join(rows))
            return cap("\n\n".join(text_parts)) or "[no extractable text in PDF]"

        # ---- Word ----
        if ext == ".docx":
            try:
                from docx import Document
            except ImportError:
                return "[error] python-docx not installed — pip install python-docx"
            doc = Document(path)
            parts = [p.text for p in doc.paragraphs if p.text.strip()]
            for table in doc.tables:
                for row in table.rows:
                    parts.append("\t".join(c.text.strip() for c in row.cells))
            return cap("\n".join(parts)) or "[empty document]"

        # ---- PowerPoint: slides + tables + speaker notes (notes labeled) ----
        if ext in (".pptx", ".pptm"):
            try:
                from pptx import Presentation
            except ImportError:
                return "[error] python-pptx not installed — pip install python-pptx"
            prs = Presentation(path)
            parts = []
            for si, slide in enumerate(prs.slides, 1):
                parts.append(f"=== Slide {si} ===")
                for shape in slide.shapes:
                    # text frames (titles, bodies, text boxes)
                    if shape.has_text_frame:
                        txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
                        if txt.strip():
                            parts.append(txt)
                    # tables
                    if shape.has_table:
                        tbl = shape.table
                        for row in tbl.rows:
                            parts.append("\t".join(cell.text.strip() for cell in row.cells))
                # speaker notes — clearly labeled so a reviewer can see (and audit) them
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text
                    if notes and notes.strip():
                        parts.append(f"--- Speaker notes (slide {si}) ---\n{notes.strip()}")
            return cap("\n".join(parts)) or "[no extractable content in presentation]"

        # ---- Fallback: read as text ----
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            return cap(f.read())

    except Exception as e:
        return f"[error] reading {path}: {e}"
    
# def _read_file(path: str) -> str:
#     staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())

#     # Resolve relative paths against staging dir
#     if not os.path.isabs(path):
#         path = os.path.join(staging, path)

#     if not os.path.isfile(path):
#         return f"[error] file not found: {path}"

#     ext = os.path.splitext(path)[1].lower()
#     max_chars = 50000

#     try:
#         # Plain text formats
#         if ext in (".txt", ".csv", ".json", ".md", ".py", ".js", ".html",
#                    ".xml", ".yaml", ".yml", ".toml", ".log", ".tsv", ".sql"):
#             with open(path, "r", encoding="utf-8", errors="replace") as f:
#                 return f.read()[:max_chars]

#         # Excel
#         if ext in (".xlsx", ".xls"):
#             import openpyxl
#             wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
#             parts = []
#             for sn in wb.sheetnames:
#                 ws = wb[sn]
#                 parts.append(f"=== Sheet: {sn} ===")
#                 for row in ws.iter_rows(values_only=True):
#                     parts.append("\t".join(str(c) if c is not None else "" for c in row))
#             wb.close()
#             return "\n".join(parts)[:max_chars]

#         # PDF
#         if ext == ".pdf":
#             try:
#                 import pdfplumber
#                 text_parts = []
#                 with pdfplumber.open(path) as pdf:
#                     for page in pdf.pages:
#                         t = page.extract_text()
#                         if t:
#                             text_parts.append(t)
#                 return "\n\n".join(text_parts)[:max_chars] or "[no extractable text in PDF]"
#             except ImportError:
#                 return "[error] pdfplumber not installed — pip install pdfplumber"

#         # Word
#         if ext == ".docx":
#             try:
#                 from docx import Document
#                 doc = Document(path)
#                 text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
#                 # Also extract tables
#                 for table in doc.tables:
#                     for row in table.rows:
#                         text += "\n" + "\t".join(c.text.strip() for c in row.cells)
#                 return text[:max_chars] or "[empty document]"
#             except ImportError:
#                 return "[error] python-docx not installed — pip install python-docx"

#         # Fallback: try reading as text
#         with open(path, "r", encoding="utf-8", errors="replace") as f:
#             return f.read()[:max_chars]

#     except Exception as e:
#         return f"[error] reading {path}: {e}"


# ─── Write file ───────────────────────────────────────────────────────────────



@register_tool(
    name="write_file",
    description=(
        "Write text content to a file. Use for creating reports, saving "
        "analysis results, generating CSV/JSON output, or writing code files. "
        "The file is created in the task staging directory."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Filename or relative path for the output file.",
            },
            "content": {
                "type": "string",
                "description": "Content to write to the file.",
            }
        },
        "required": ["path", "content"],
    },
)
def _write_file(path: str, content: str) -> str:
    staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())

    if not os.path.isabs(path):
        path = os.path.join(staging, path)

    try:
        os.makedirs(os.path.dirname(path) or ".", exist_ok=True)
        with open(path, "w", encoding="utf-8") as f:
            f.write(content)
        size = os.path.getsize(path)
        return f"File written: {path} ({size:,} bytes)"
    except Exception as e:
        return f"[error] writing {path}: {e}"


# ─── List directory ───────────────────────────────────────────────────────────

@register_tool(
    name="list_directory",
    description=(
        "List files and subdirectories in a directory. Defaults to the "
        "task staging directory if no path is given."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Directory path to list. Defaults to staging dir.",
                "default": ".",
            }
        },
    },
)
def _list_directory(path: str = ".") -> str:
    staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())

    if not os.path.isabs(path):
        path = os.path.join(staging, path)

    if not os.path.isdir(path):
        return f"[error] not a directory: {path}"

    try:
        entries = []
        for name in sorted(os.listdir(path)):
            full = os.path.join(path, name)
            if os.path.isdir(full):
                entries.append(f"  {name}/")
            else:
                size = os.path.getsize(full)
                entries.append(f"  {name}  ({size:,} bytes)")
        return f"Directory: {path}\n" + "\n".join(entries) if entries else f"[empty directory: {path}]"
    except Exception as e:
        return f"[error] listing {path}: {e}"


# ─── Web fetch ────────────────────────────────────────────────────────────────

@register_tool(
    name="web_fetch",
    description=(
        "Fetch the text content of a web page at a given URL. Use for "
        "retrieving data from APIs, reading web pages, or downloading "
        "publicly available documents. Returns text content truncated "
        "to 30,000 characters."
    ),
    parameters={
        "type": "object",
        "properties": {
            "url": {
                "type": "string",
                "description": "The URL to fetch.",
            }
        },
        "required": ["url"],
    },
)
def _web_fetch(url: str) -> str:
    import urllib.request
    import urllib.error

    try:
        req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=30) as resp:
            content_type = resp.headers.get("Content-Type", "")
            raw = resp.read(500_000)  # cap at 500KB

            # Detect encoding
            encoding = "utf-8"
            if "charset=" in content_type:
                encoding = content_type.split("charset=")[-1].split(";")[0].strip()

            text = raw.decode(encoding, errors="replace")

            # Strip HTML tags for readability
            if "html" in content_type.lower():
                import re
                text = re.sub(r"<script[^>]*>.*?</script>", "", text, flags=re.DOTALL)
                text = re.sub(r"<style[^>]*>.*?</style>", "", text, flags=re.DOTALL)
                text = re.sub(r"<[^>]+>", " ", text)
                text = re.sub(r"\s+", " ", text).strip()

            return text[:30000] or "[empty response]"

    except urllib.error.HTTPError as e:
        return f"[error] HTTP {e.code}: {e.reason}"
    except urllib.error.URLError as e:
        return f"[error] URL error: {e.reason}"
    except Exception as e:
        return f"[error] fetching {url}: {e}"


# ─── Calculator (safe arithmetic) ─────────────────────────────────────────────

@register_tool(
    name="calculator",
    description="Evaluate a basic arithmetic expression and return the result.",
    parameters={
        "type": "object",
        "properties": {
            "expression": {
                "type": "string",
                "description": "Arithmetic expression, e.g. '2 * (3 + 4) / 7'",
            }
        },
        "required": ["expression"],
    },
)
def _calculator(expression: str) -> str:
    import ast
    import operator as op

    ops = {
        ast.Add: op.add, ast.Sub: op.sub, ast.Mult: op.mul,
        ast.Div: op.truediv, ast.Pow: op.pow, ast.Mod: op.mod,
        ast.USub: op.neg, ast.UAdd: op.pos, ast.FloorDiv: op.floordiv,
    }

    def _eval(node):
        if isinstance(node, ast.Constant):
            if isinstance(node.value, (int, float)):
                return node.value
            raise ValueError("only numeric constants allowed")
        if isinstance(node, ast.BinOp) and type(node.op) in ops:
            return ops[type(node.op)](_eval(node.left), _eval(node.right))
        if isinstance(node, ast.UnaryOp) and type(node.op) in ops:
            return ops[type(node.op)](_eval(node.operand))
        raise ValueError("unsupported expression")

    try:
        tree = ast.parse(expression, mode="eval")
        return str(_eval(tree.body))
    except Exception as e:  # noqa: BLE001
        return f"[calculator error] {e}"
    

# ─── Search in file ───────────────────────────────────────────────────────────

@register_tool(
    name="search_in_file",
    description=(
        "Search for a keyword or phrase within a file and return matching "
        "lines with surrounding context. Use when a file is large and you "
        "need to find specific data points without reading the entire file."
    ),
    parameters={
        "type": "object",
        "properties": {
            "path": {
                "type": "string",
                "description": "Path to the file to search.",
            },
            "query": {
                "type": "string",
                "description": "Keyword or phrase to search for (case-insensitive).",
            },
            "context_lines": {
                "type": "integer",
                "description": "Number of lines of context above and below each match. Default 3.",
                "default": 3,
            }
        },
        "required": ["path", "query"],
    },
)
def _search_in_file(path: str, query: str, context_lines: int = 3) -> str:
    staging = os.environ.get("DRA_AGENT_WORKDIR", os.getcwd())
    if not os.path.isabs(path):
        path = os.path.join(staging, path)

    if not os.path.isfile(path):
        return f"[error] file not found: {path}"

    # Extract text first (reuse read_file logic for non-text formats).
    # CRITICAL: pass max_chars=0 for a LOSSLESS extraction. read_file's default
    # 100K-char cap is meant to protect the MODEL's context on a full read — but
    # here we are SEARCHING, and the search result is separately capped at 15K
    # below. If we inherited the 100K cap, a match past the first ~30-50 pages of
    # a large PDF would be silently truncated away before the search ran, and the
    # tool would report "no matches" for data that is actually in the file — a
    # false negative the model cannot detect. Search the whole file; cap the
    # OUTPUT, not the input.
    ext = os.path.splitext(path)[1].lower()
    try:
        if ext in (".xlsx", ".xls"):
            text = _read_file(path, max_chars=0)
        elif ext == ".pdf":
            text = _read_file(path, max_chars=0)
        elif ext == ".docx":
            text = _read_file(path, max_chars=0)
        else:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
    except Exception as e:
        return f"[error] reading {path}: {e}"

    lines = text.split("\n")
    query_lower = query.lower()
    matches = []

    for i, line in enumerate(lines):
        if query_lower in line.lower():
            start = max(0, i - context_lines)
            end = min(len(lines), i + context_lines + 1)
            snippet = "\n".join(
                f"{'>>>' if j == i else '   '} {lines[j]}"
                for j in range(start, end)
            )
            matches.append(f"[Match at line {i+1}]\n{snippet}")

    if not matches:
        return f"No matches for '{query}' in {os.path.basename(path)}"

    total = len(matches)
    header = (f"Found {total} match(es) for '{query}' in "
              f"{os.path.basename(path)}:\n\n")
    body = header + "\n\n".join(matches)
    if len(body) <= 15000:
        return body
    # Too many/large matches to return in full. Rather than cut mid-match with no
    # signal, keep whole matches up to the cap and tell the model exactly how many
    # were omitted, so it can narrow the query or page through instead of assuming
    # it saw everything.
    kept, acc = [], len(header)
    for m in matches:
        if acc + len(m) + 2 > 15000:
            break
        kept.append(m)
        acc += len(m) + 2
    omitted = total - len(kept)
    note = (f"\n\n[... {omitted} more match(es) not shown (output capped). "
            f"Narrow the query or use read_file with max_chars to target a "
            f"specific region.]")
    return header + "\n\n".join(kept) + note


# ─── Web search (Serper/Google primary, DuckDuckGo fallback) ─────────────────

@register_tool(
    name="web_search",
    description=(
        "Search the web and return top results (titles, URLs, snippets).\n\n"
        "Uses Google (via Serper) when SERPER_API_KEY is set — far better on "
        "institutional/government/data sources — and falls back to DuckDuckGo "
        "otherwise or if Serper returns nothing. Use for finding current data, "
        "verifying facts, or retrieving external information not in the files."
    ),
    parameters={
        "type": "object",
        "properties": {
            "query": {
                "type": "string",
                "description": "Search query.",
            },
            "max_results": {
                "type": "integer",
                "description": "Number of results to return. Default 5.",
                "default": 5,
            }
        },
        "required": ["query"],
    },
)
def _web_search(query: str, max_results: int = 5) -> str:
    """Serper (Google) primary, DuckDuckGo fallback.

    Route is decided by SERPER_API_KEY: when set, Serper is tried first; on a
    missing key, empty result set, or ANY Serper error we fall through to DDG.
    Every fallback is LOGGED (logger 'dra.tools') so a Serper outage or key
    problem is visible in the run logs rather than silently degrading search
    quality — DDG returns far weaker results on institutional/government/data
    queries, which is exactly where DRA tasks need the web.
    """
    serper_key = os.environ.get("SERPER_API_KEY")
    if serper_key:
        try:
            out = _web_search_serper(query, max_results, serper_key)
            if out:
                return out
            logger.warning("web_search: Serper returned no results for %r — "
                           "falling back to DuckDuckGo", query)
        except Exception as e:                                  # noqa: BLE001
            logger.warning("web_search: Serper failed (%s) for %r — falling "
                           "back to DuckDuckGo", e, query)
    else:
        logger.info("web_search: SERPER_API_KEY not set — using DuckDuckGo")
    return _web_search_ddg(query, max_results)


def _web_search_serper(query: str, max_results: int, api_key: str) -> str:
    """Google results via Serper (https://google.serper.dev/search). Returns a
    formatted string, or "" when Serper yields no organic results (so the caller
    falls back to DDG). Raises on transport/HTTP errors (caller logs + falls
    back)."""
    import requests

    resp = requests.post(
        "https://google.serper.dev/search",
        headers={"X-API-KEY": api_key, "Content-Type": "application/json"},
        json={"q": query, "num": max_results},
        timeout=20,
    )
    resp.raise_for_status()
    data = resp.json()
    organic = data.get("organic", []) or []
    if not organic:
        return ""
    parts = []
    for i, r in enumerate(organic[:max_results], 1):
        parts.append(
            f"[{i}] {r.get('title', 'No title')}\n"
            f"    URL: {r.get('link', 'N/A')}\n"
            f"    {r.get('snippet', 'No snippet')}"
        )
    return "\n\n".join(parts)


def _web_search_ddg(query: str, max_results: int = 5) -> str:
    """DuckDuckGo fallback (no API key needed)."""
    try:
        from duckduckgo_search import DDGS
    except ImportError:
        return "[error] duckduckgo-search not installed — pip install duckduckgo-search"

    try:
        with DDGS() as ddgs:
            results = list(ddgs.text(query, max_results=max_results))

        if not results:
            return f"No results found for '{query}'"

        parts = []
        for i, r in enumerate(results, 1):
            parts.append(
                f"[{i}] {r.get('title', 'No title')}\n"
                f"    URL: {r.get('href', 'N/A')}\n"
                f"    {r.get('body', 'No snippet')}"
            )
        return "\n\n".join(parts)

    except Exception as e:
        return f"[error] web search failed: {e}"