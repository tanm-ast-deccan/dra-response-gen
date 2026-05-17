"""
corpus_tools.py — Shared file access functions for Tier 2 corpus serving.

These functions are the core of the Tier 2 architecture. They provide
structured access to a corpus of evaluation files, used by:
  - The MCP server (corpus_server.py) for OpenAI's native MCP client
  - Claude's custom tools in the agent loop
  - Gemini's fallback (if File Search is unavailable)

The corpus is a directory of files that the SME provides as part of
a PromptPackage. The tools let agents navigate and read the corpus
on demand, rather than stuffing everything into context upfront.

This follows the APEX-Agents / Archipelago pattern:
  - list_documents() → what's in the corpus?
  - search_corpus()  → which files are relevant to my query?
  - fetch_document()  → give me the full content of this file

Design principles:
  - File content is extracted once and cached in memory
  - Search is keyword-based (not vector / embedding-based) for
    simplicity and reproducibility
  - Binary files (PDF, DOCX, XLSX) are extracted to text on load
  - Images return metadata only (agents can't "see" images via MCP)
"""

from __future__ import annotations

import os
import re
import hashlib
from pathlib import Path
from dataclasses import dataclass, field
from typing import Optional


# ─── Data structures ─────────────────────────────────────────────────────

@dataclass
class CorpusDocument:
    """A single document in the corpus, with extracted text."""
    doc_id: str                # stable hash-based ID
    filename: str              # original filename
    path: str                  # full path on disk
    extension: str             # .pdf, .docx, etc.
    size_bytes: int
    text_content: str          # extracted text (empty for unsupported)
    num_tokens_est: int        # rough token estimate (chars / 4)
    extractable: bool          # whether text was successfully extracted
    mime_type: str = ""


@dataclass
class SearchResult:
    """A single search hit."""
    doc_id: str
    filename: str
    snippet: str               # context around the match
    relevance_score: float     # 0.0 - 1.0
    match_count: int           # number of keyword matches


# ─── Text extraction ─────────────────────────────────────────────────────

def _extract_text(filepath: str) -> tuple[str, bool]:
    """
    Extract readable text from a file.

    Returns (text, success). For text-based formats, reads directly.
    For binary formats (PDF, DOCX, XLSX, PPTX), uses extraction
    libraries if available.
    """
    ext = Path(filepath).suffix.lower()

    # ── Plain text formats ────────────────────────────────────────
    if ext in (".txt", ".md", ".csv", ".tsv", ".html", ".htm",
               ".json", ".xml", ".yaml", ".yml", ".log"):
        try:
            with open(filepath, "r", encoding="utf-8", errors="replace") as f:
                return f.read(), True
        except Exception as e:
            return f"[Read error: {e}]", False

    # ── PDF ───────────────────────────────────────────────────────
    if ext == ".pdf":
        try:
            import pdfplumber
            pages = []
            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"[Page {i+1}]\n{text}")
            return "\n\n".join(pages), bool(pages)
        except ImportError:
            return "[PDF extraction requires: pip install pdfplumber]", False
        except Exception as e:
            return f"[PDF extraction error: {e}]", False

    # ── DOCX ──────────────────────────────────────────────────────
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(filepath)
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            return "\n\n".join(paragraphs), bool(paragraphs)
        except ImportError:
            return "[DOCX extraction requires: pip install python-docx]", False
        except Exception as e:
            return f"[DOCX extraction error: {e}]", False

    # ── XLSX ──────────────────────────────────────────────────────
    if ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"=== Sheet: {sheet_name} ===")
                rows_read = 0
                for row in ws.iter_rows(max_row=1000, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    if any(cells):  # skip fully empty rows
                        parts.append("\t".join(cells))
                        rows_read += 1
                parts.append(f"[{rows_read} rows]")
            wb.close()
            text = "\n".join(parts)
            return text, bool(parts)
        except ImportError:
            return "[XLSX extraction requires: pip install openpyxl]", False
        except Exception as e:
            return f"[XLSX extraction error: {e}]", False

    # ── PPTX ──────────────────────────────────────────────────────
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text)
                if slide_texts:
                    parts.append(
                        f"--- Slide {i+1} ---\n" + "\n".join(slide_texts)
                    )
            return "\n\n".join(parts), bool(parts)
        except ImportError:
            return "[PPTX extraction requires: pip install python-pptx]", False
        except Exception as e:
            return f"[PPTX extraction error: {e}]", False

    # ── Images — metadata only ────────────────────────────────────
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff"):
        size = os.path.getsize(filepath)
        return (
            f"[Image file: {os.path.basename(filepath)}, "
            f"size={size:,} bytes, format={ext}. "
            f"Image content not available via text extraction.]"
        ), False

    # ── Unsupported ───────────────────────────────────────────────
    return f"[Unsupported format: {ext}]", False


def _make_doc_id(filepath: str) -> str:
    """Generate a stable document ID from the filepath."""
    # Use filename + size for a simple stable ID
    filename = os.path.basename(filepath)
    size = os.path.getsize(filepath) if os.path.exists(filepath) else 0
    raw = f"{filename}:{size}"
    return hashlib.sha256(raw.encode()).hexdigest()[:12]


def _guess_mime(ext: str) -> str:
    """Guess MIME type from extension."""
    mime_map = {
        ".pdf": "application/pdf",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".csv": "text/csv",
        ".tsv": "text/tab-separated-values",
        ".txt": "text/plain",
        ".md": "text/markdown",
        ".html": "text/html",
        ".htm": "text/html",
        ".json": "application/json",
        ".xml": "application/xml",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".png": "image/png",
    }
    return mime_map.get(ext, "application/octet-stream")


# ─── The Corpus ──────────────────────────────────────────────────────────

class Corpus:
    """
    In-memory index of a file corpus for Tier 2 access.

    Load a directory of files, extract text, and provide structured
    access via list/search/fetch operations.

    Usage:
        corpus = Corpus("/path/to/eval/files")
        corpus.load()

        # List all documents
        docs = corpus.list_documents()

        # Search for relevant files
        hits = corpus.search("EBITDA restructuring")

        # Read a specific file
        content = corpus.fetch_document("abc123def456")
    """

    def __init__(self, corpus_dir: str):
        self.corpus_dir = os.path.abspath(corpus_dir)
        self.documents: dict[str, CorpusDocument] = {}  # doc_id → doc
        self._loaded = False

    def load(self) -> int:
        """
        Load and index all files in the corpus directory.

        Returns the number of documents loaded. Files are extracted
        to text on load and cached in memory.
        """
        if not os.path.isdir(self.corpus_dir):
            raise FileNotFoundError(
                f"Corpus directory not found: {self.corpus_dir}"
            )

        self.documents.clear()

        for entry in sorted(os.listdir(self.corpus_dir)):
            filepath = os.path.join(self.corpus_dir, entry)
            if not os.path.isfile(filepath):
                continue

            # Skip hidden files and common junk
            if entry.startswith(".") or entry.startswith("~"):
                continue

            ext = Path(entry).suffix.lower()
            doc_id = _make_doc_id(filepath)
            text_content, extractable = _extract_text(filepath)

            doc = CorpusDocument(
                doc_id=doc_id,
                filename=entry,
                path=filepath,
                extension=ext,
                size_bytes=os.path.getsize(filepath),
                text_content=text_content,
                num_tokens_est=len(text_content) // 4,
                extractable=extractable,
                mime_type=_guess_mime(ext),
            )
            self.documents[doc_id] = doc

        self._loaded = True
        return len(self.documents)

    def load_from_paths(self, file_paths: list[str]) -> int:
        """
        Load specific files (not a directory scan).

        Useful when the PromptPackage specifies explicit file paths
        rather than a directory.
        """
        self.documents.clear()

        for filepath in file_paths:
            if not os.path.isfile(filepath):
                continue

            entry = os.path.basename(filepath)
            ext = Path(entry).suffix.lower()
            doc_id = _make_doc_id(filepath)
            text_content, extractable = _extract_text(filepath)

            doc = CorpusDocument(
                doc_id=doc_id,
                filename=entry,
                path=filepath,
                extension=ext,
                size_bytes=os.path.getsize(filepath),
                text_content=text_content,
                num_tokens_est=len(text_content) // 4,
                extractable=extractable,
                mime_type=_guess_mime(ext),
            )
            self.documents[doc_id] = doc

        self._loaded = True
        return len(self.documents)

    # ─── Tool implementations ─────────────────────────────────────

    def list_documents(self) -> list[dict]:
        """
        List all documents in the corpus.

        Returns metadata only (not full content) — the agent decides
        which files to read based on filenames, sizes, and types.

        MCP tool: list_documents
        Claude tool: list_files
        """
        docs = []
        for doc in sorted(self.documents.values(), key=lambda d: d.filename):
            docs.append({
                "id": doc.doc_id,
                "filename": doc.filename,
                "type": doc.extension,
                "size_bytes": doc.size_bytes,
                "tokens_est": doc.num_tokens_est,
                "extractable": doc.extractable,
            })
        return docs

    def search_corpus(
        self,
        query: str,
        max_results: int = 5,
    ) -> list[dict]:
        """
        Search across all documents for keyword matches.

        Simple TF-based keyword search. Each query term is searched
        independently; documents matching more terms rank higher.
        Snippets show context around the first match.

        MCP tool: search
        Claude tool: search_files
        """
        # Tokenize query into lowercase terms
        terms = [t.lower() for t in re.split(r'\W+', query) if len(t) > 2]
        if not terms:
            return []

        scored: list[tuple[float, int, CorpusDocument, str]] = []

        for doc in self.documents.values():
            if not doc.extractable:
                continue

            text_lower = doc.text_content.lower()
            total_matches = 0
            terms_found = 0
            first_snippet = ""

            for term in terms:
                count = text_lower.count(term)
                if count > 0:
                    total_matches += count
                    terms_found += 1

                    # Extract snippet around first occurrence
                    if not first_snippet:
                        idx = text_lower.index(term)
                        start = max(0, idx - 100)
                        end = min(len(doc.text_content), idx + len(term) + 100)
                        first_snippet = doc.text_content[start:end].strip()

            if terms_found > 0:
                # Score: fraction of query terms found × log(match count)
                import math
                relevance = (terms_found / len(terms)) * (
                    1 + math.log(1 + total_matches)
                )
                scored.append((relevance, total_matches, doc, first_snippet))

        # Sort by relevance descending
        scored.sort(key=lambda x: x[0], reverse=True)

        results = []
        for relevance, match_count, doc, snippet in scored[:max_results]:
            results.append({
                "id": doc.doc_id,
                "filename": doc.filename,
                "snippet": snippet[:300],
                "relevance_score": round(relevance, 3),
                "match_count": match_count,
            })

        return results

    def fetch_document(
        self,
        doc_id: str,
        max_chars: Optional[int] = None,
    ) -> dict:
        """
        Fetch the full text content of a document by ID.

        Returns the complete extracted text. For large documents,
        max_chars can limit the returned content (agent can call
        again with an offset in a more sophisticated implementation).

        MCP tool: fetch
        Claude tool: read_file
        """
        doc = self.documents.get(doc_id)
        if doc is None:
            # Try matching by filename as fallback
            for d in self.documents.values():
                if d.filename == doc_id or d.filename.lower() == doc_id.lower():
                    doc = d
                    break

        if doc is None:
            return {
                "error": f"Document not found: {doc_id}",
                "available_ids": [d.doc_id for d in self.documents.values()],
            }

        content = doc.text_content
        truncated = False
        if max_chars and len(content) > max_chars:
            content = content[:max_chars]
            truncated = True

        return {
            "id": doc.doc_id,
            "filename": doc.filename,
            "type": doc.extension,
            "size_bytes": doc.size_bytes,
            "content": content,
            "truncated": truncated,
            "total_chars": len(doc.text_content),
        }

    # ─── Stats ────────────────────────────────────────────────────

    @property
    def total_tokens_est(self) -> int:
        """Estimated total tokens across all extractable documents."""
        return sum(
            d.num_tokens_est for d in self.documents.values()
            if d.extractable
        )

    @property
    def stats(self) -> dict:
        """Summary statistics about the corpus."""
        extractable = [d for d in self.documents.values() if d.extractable]
        return {
            "total_documents": len(self.documents),
            "extractable_documents": len(extractable),
            "total_size_bytes": sum(d.size_bytes for d in self.documents.values()),
            "total_tokens_est": self.total_tokens_est,
            "types": dict(sorted(
                {ext: 0 for ext in set(d.extension for d in self.documents.values())}.items()
            )),
        }

    def __repr__(self) -> str:
        return (
            f"Corpus(dir={self.corpus_dir!r}, "
            f"docs={len(self.documents)}, "
            f"tokens≈{self.total_tokens_est:,})"
        )
