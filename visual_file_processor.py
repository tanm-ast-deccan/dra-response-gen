"""
visual_file_processor.py — Rasterize documents to images and extract visual content.

Used by the Qwen adapter's read_file_visual tool to give a text-only reasoning
model (Qwen 3.7 Max) full visual understanding of charts, diagrams, tables,
and layouts via a separate vision model call (Qwen 3.7 Plus or any multimodal).

Architecture:
    PDF/PPTX/DOCX/image → rasterize to PNG pages → base64 encode
    → call vision model with "describe this page in detail" → return text

Pipeline:
    ┌──────────┐    ┌─────────────┐    ┌──────────────┐    ┌──────────┐
    │ Any file │───▶│ Rasterize   │───▶│ base64 PNGs  │───▶│ Vision   │───▶ text
    │ on disk  │    │ to pages    │    │ per page     │    │ model    │
    └──────────┘    └─────────────┘    └──────────────┘    └──────────┘

Rasterization methods (tried in order):
    1. pymupdf (fitz) — fastest, best quality, handles most PDFs
    2. pdf2image + poppler — fallback for PDFs
    3. Pillow — direct for image files
    4. LibreOffice headless — converts DOCX/PPTX/XLSX to PDF first

Vision model options (configured via env):
    - Qwen 3.7 Plus (DashScope) — default, same vendor
    - Qwen-VL series (open-weight, local)
    - GPT-4o / Claude (cross-vendor, if API keys available)
"""

from __future__ import annotations

import os
import re
import sys
import json
import base64
import asyncio
import logging
import shutil
import tempfile
from pathlib import Path
from typing import Optional

_log = logging.getLogger("dra.visual")

# ── Configuration ─────────────────────────────────────────────────────────────

# Vision model for describing rasterized pages
VISION_MODEL     = os.environ.get("QWEN_VISION_MODEL", "qwen3.7-plus-preview")
VISION_BASE_URL  = os.environ.get("QWEN_VISION_BASE_URL",
                                   os.environ.get("QWEN_BASE_URL",
                                                  "https://dashscope.aliyuncs.com/compatible-mode/v1"))
VISION_API_KEY   = os.environ.get("QWEN_VISION_API_KEY",
                                   os.environ.get("QWEN_API_KEY", ""))

# Rasterization settings
DPI             = 150           # resolution for PDF→image
MAX_PAGES       = 20            # max pages to rasterize per document
MAX_IMAGE_DIM   = 2048          # max width or height in pixels
JPEG_QUALITY    = 85            # quality for JPEG compression (smaller than PNG)
CONVERT_TIMEOUT = 120           # seconds for LibreOffice conversion


# ── Rasterization ─────────────────────────────────────────────────────────────

def rasterize_file(filepath: str) -> list[dict]:
    """
    Rasterize a file into a list of page images.

    Returns:
        List of {"page": int, "base64": str, "mime": str, "width": int, "height": int}
        Empty list if the file cannot be rasterized.
    """
    ext = os.path.splitext(filepath)[1].lower()

    if ext == ".pdf":
        return _rasterize_pdf(filepath)
    elif ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
        return _rasterize_image(filepath)
    elif ext in (".pptx", ".docx", ".xlsx", ".odt", ".ods", ".odp"):
        return _rasterize_via_libreoffice(filepath)
    else:
        return []


def _rasterize_pdf(filepath: str) -> list[dict]:
    """Rasterize PDF pages to images. Tries pymupdf first, then pdf2image."""
    # Method 1: pymupdf (fitz) — fastest, best quality
    try:
        import fitz
        doc = fitz.open(filepath)
        pages = []
        for i, page in enumerate(doc):
            if i >= MAX_PAGES:
                break
            # Render at target DPI
            zoom = DPI / 72.0
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat, alpha=False)
            # Resize if too large
            img_bytes = pix.tobytes("png")
            b64, w, h = _resize_and_encode(img_bytes, "png")
            pages.append({
                "page": i + 1,
                "base64": b64,
                "mime": "image/png",
                "width": w,
                "height": h,
            })
        doc.close()
        return pages
    except ImportError:
        pass

    # Method 2: pdf2image + poppler
    try:
        from pdf2image import convert_from_path
        images = convert_from_path(
            filepath, dpi=DPI, first_page=1,
            last_page=MAX_PAGES, fmt="png",
        )
        pages = []
        for i, img in enumerate(images):
            b64, w, h = _pil_to_base64(img)
            pages.append({
                "page": i + 1,
                "base64": b64,
                "mime": "image/png",
                "width": w,
                "height": h,
            })
        return pages
    except ImportError:
        pass

    _log.warning("PDF rasterization unavailable: install pymupdf or pdf2image+poppler")
    return []


def _rasterize_image(filepath: str) -> list[dict]:
    """Load an image file directly."""
    try:
        from PIL import Image
        img = Image.open(filepath)
        if img.mode == "RGBA":
            img = img.convert("RGB")
        b64, w, h = _pil_to_base64(img)
        return [{
            "page": 1,
            "base64": b64,
            "mime": "image/png",
            "width": w,
            "height": h,
        }]
    except ImportError:
        _log.warning("Image rasterization unavailable: install Pillow")
        return []
    except Exception as e:
        _log.warning(f"Image rasterization failed for {filepath}: {e}")
        return []


def _rasterize_via_libreoffice(filepath: str) -> list[dict]:
    """Convert DOCX/PPTX/XLSX to PDF via LibreOffice, then rasterize the PDF."""
    lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo_bin:
        _log.warning("LibreOffice not found — cannot rasterize office files")
        return []

    with tempfile.TemporaryDirectory(prefix="visual_lo_") as tmpdir:
        try:
            result = subprocess.run(
                [lo_bin, "--headless", "--convert-to", "pdf",
                 "--outdir", tmpdir, filepath],
                capture_output=True, timeout=CONVERT_TIMEOUT,
            )
            if result.returncode != 0:
                _log.warning(f"LibreOffice conversion failed: {result.stderr[:200]}")
                return []

            # Find the generated PDF
            stem = os.path.splitext(os.path.basename(filepath))[0]
            pdf_path = os.path.join(tmpdir, f"{stem}.pdf")
            if not os.path.exists(pdf_path):
                # LibreOffice might use slightly different naming
                pdfs = [f for f in os.listdir(tmpdir) if f.endswith(".pdf")]
                if pdfs:
                    pdf_path = os.path.join(tmpdir, pdfs[0])
                else:
                    return []

            return _rasterize_pdf(pdf_path)

        except subprocess.TimeoutExpired:
            _log.warning(f"LibreOffice conversion timed out for {filepath}")
            return []
        except Exception as e:
            _log.warning(f"LibreOffice conversion error: {e}")
            return []


# Need this import at module level for _rasterize_via_libreoffice
import subprocess


def _pil_to_base64(img) -> tuple[str, int, int]:
    """Convert a PIL Image to base64 PNG, resizing if needed."""
    from PIL import Image
    import io

    w, h = img.size
    if w > MAX_IMAGE_DIM or h > MAX_IMAGE_DIM:
        ratio = min(MAX_IMAGE_DIM / w, MAX_IMAGE_DIM / h)
        new_w, new_h = int(w * ratio), int(h * ratio)
        img = img.resize((new_w, new_h), Image.LANCZOS)
        w, h = new_w, new_h

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    b64 = base64.b64encode(buf.getvalue()).decode("ascii")
    return b64, w, h


def _resize_and_encode(img_bytes: bytes, fmt: str) -> tuple[str, int, int]:
    """Resize raw image bytes if needed and return base64."""
    try:
        from PIL import Image
        import io
        img = Image.open(io.BytesIO(img_bytes))
        return _pil_to_base64(img)
    except ImportError:
        # No Pillow — just encode as-is
        return base64.b64encode(img_bytes).decode("ascii"), 0, 0


# ── Vision model call ─────────────────────────────────────────────────────────

async def describe_visual_content(
    pages: list[dict],
    filename: str,
    context_hint: str = "",
    model: Optional[str] = None,
    base_url: Optional[str] = None,
    api_key: Optional[str] = None,
) -> str:
    """
    Send rasterized pages to a vision model and get a detailed description.

    Args:
        pages: Output from rasterize_file()
        filename: Original filename (for context)
        context_hint: Optional hint about what to look for
        model: Vision model to use (default: env QWEN_VISION_MODEL)
        base_url: API endpoint (default: env QWEN_VISION_BASE_URL)
        api_key: API key (default: env QWEN_VISION_API_KEY)

    Returns:
        Detailed text description of all visual content in the file.
    """
    if not pages:
        return "[No visual content could be extracted from this file]"

    vision_model = model or VISION_MODEL
    vision_url = base_url or VISION_BASE_URL
    vision_key = api_key or VISION_API_KEY

    if not vision_key:
        return (
            "[Vision model not configured. Set QWEN_VISION_API_KEY and "
            "QWEN_VISION_MODEL in .env to enable visual file reading.]"
        )

    try:
        from openai import AsyncOpenAI
        client = AsyncOpenAI(api_key=vision_key, base_url=vision_url, timeout=120)
    except ImportError:
        return "[openai package required for vision model calls]"

    # Build multimodal message with all pages
    content_parts = [
        {
            "type": "text",
            "text": (
                f"You are analyzing the file '{filename}' which has been "
                f"rasterized into {len(pages)} page image(s). "
                f"Provide a DETAILED description of EVERY visual element: "
                f"charts (type, axes, data points, trends), tables (all rows "
                f"and columns with values), diagrams (structure, labels, "
                f"connections), images, formatting, and any text visible in "
                f"the images. Be exhaustive — numbers matter, labels matter, "
                f"colors and legends matter. "
                f"{'Additional context: ' + context_hint if context_hint else ''}"
            ),
        },
    ]

    for page_info in pages:
        content_parts.append({
            "type": "image_url",
            "image_url": {
                "url": f"data:{page_info['mime']};base64,{page_info['base64']}",
            },
        })
        content_parts.append({
            "type": "text",
            "text": f"[Page {page_info['page']}]",
        })

    try:
        response = await client.chat.completions.create(
            model=vision_model,
            messages=[{"role": "user", "content": content_parts}],
            max_tokens=4096,
            temperature=0.2,  # low temperature for factual description
        )

        description = response.choices[0].message.content or ""
        token_info = ""
        if response.usage:
            token_info = (
                f"\n[Vision model: {vision_model}, "
                f"tokens: {response.usage.prompt_tokens}in/"
                f"{response.usage.completion_tokens}out]"
            )

        return description + token_info

    except Exception as e:
        _log.error(f"Vision model call failed: {e}")
        return f"[Vision model error: {e}]"


# ── Combined file reader (text + visual) ──────────────────────────────────────

def has_visual_content(filepath: str) -> bool:
    """Check if a file likely contains visual elements worth rasterizing."""
    ext = os.path.splitext(filepath)[1].lower()

    # Image files: always visual
    if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff"):
        return True

    # PPTX: almost always has visual layout
    if ext == ".pptx":
        try:
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        return True
            return False
        except Exception:
            return True

    # PDF: check if it has images or is scanned
    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(filepath)
            for page in doc:
                images = page.get_images()
                if images:
                    doc.close()
                    return True
                # Also check if page has very little text (likely scanned)
                text = page.get_text().strip()
                if len(text) < 50:
                    doc.close()
                    return True
            doc.close()
            return False
        except ImportError:
            # Can't check — assume visual if PDF is > 100KB
            return os.path.getsize(filepath) > 100_000
        except Exception:
            return False

    # XLSX: check for chart sheets (openpyxl)
    if ext == ".xlsx":
        try:
            from openpyxl import load_workbook
            wb = load_workbook(filepath, read_only=True, data_only=True)
            for ws in wb.worksheets:
                if hasattr(ws, '_charts') and ws._charts:
                    wb.close()
                    return True
            wb.close()
            return False
        except Exception:
            return False

    # DOCX: check for images
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(filepath)
            for rel in doc.part.rels.values():
                if "image" in rel.reltype:
                    return True
            return False
        except Exception:
            return False

    return False


async def read_file_complete(
    filepath: str,
    text_reader_func,
    file_paths: list[str],
) -> str:
    """
    Read a file with both text extraction AND visual description.

    This is the upgraded read_file that provides full file understanding.
    Calls the text reader first, then checks for visual content and
    adds a vision model description if present.

    Args:
        filepath: Path to the file
        text_reader_func: The existing _exec_read_file function
        file_paths: List of all task file paths (for text reader)

    Returns:
        Combined text: extracted text + visual description (if applicable)
    """
    filename = os.path.basename(filepath)

    # Always get text extraction first
    text_content = text_reader_func(filename, file_paths)

    # Check if visual description would add value
    if not has_visual_content(filepath):
        return text_content

    # Rasterize and describe
    pages = rasterize_file(filepath)
    if not pages:
        return text_content + "\n\n[Note: File may contain visual elements " \
               "but rasterization was not available.]"

    visual_desc = await describe_visual_content(pages, filename)

    return (
        f"{text_content}\n\n"
        f"══ VISUAL CONTENT DESCRIPTION ══\n"
        f"The following visual elements were detected in {filename}:\n\n"
        f"{visual_desc}"
    )