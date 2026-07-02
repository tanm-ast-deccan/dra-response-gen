"""
file_gen.py — Unified, provider-agnostic file creation via local code execution.

OpenRouter's chat API has no server-side sandbox, so file deliverables
(xlsx/docx/pptx) are produced the same way for EVERY provider:

  1. The prompt asks the model to emit a sentinel-wrapped Python block that
     writes `output.<fmt>` (instructions built by file_gen_instructions()).
  2. We extract that code, run it locally in the staging dir (where the input
     files live on disk), and capture the produced file.
  3. On failure, we send the traceback back to the SAME model and ask for a
     fix, retrying up to params.file_fix_attempts times.

Also provides extract_file_text() used to inline file content into prompts.
"""

from __future__ import annotations

import os
import re
import sys
import shutil
import asyncio
import logging
import subprocess
from pathlib import Path
from typing import Optional

logger = logging.getLogger("indrayudh.file_gen")

FILE_OUTPUT_START = "# IPL_FILE_OUTPUT_START"
FILE_OUTPUT_END = "# IPL_FILE_OUTPUT_END"

_LIB_HINT = {
    "xlsx": "openpyxl or pandas (with the openpyxl engine)",
    "docx": "python-docx (from docx import Document)",
    "pptx": "python-pptx (from pptx import Presentation)",
}


# ─── Prompt instructions ──────────────────────────────────────────────────────

def file_gen_instructions(fmt: str, file_list: list[str]) -> str:
    """Build the file-generation instruction appended to the user prompt."""
    files = "\n".join(f"  - {os.path.basename(p)}" for p in (file_list or [])) or "  (none)"
    lib = _LIB_HINT.get(fmt, fmt)
    return (
        f"\n\n{'=' * 60}\nFILE GENERATION REQUIREMENT\n{'=' * 60}\n\n"
        f"After your analysis you MUST append ONE Python code block that writes "
        f"the required {fmt.upper()} deliverable.\n\n"
        f"Input files available in the working directory:\n{files}\n\n"
        f"Rules:\n"
        f"1. Use {lib}. Standard file-handling libraries only — no sklearn, "
        f"statsmodels, pulp, scipy.\n"
        f"2. Save the file as exactly: output.{fmt} (current working directory).\n"
        f"3. Base content on your analysis; use exact input filenames if you read them.\n"
        f"4. Wrap the code with EXACTLY these sentinel lines as the first and "
        f"last lines inside the block:\n\n"
        f"```python\n{FILE_OUTPUT_START}\n# ... your code ...\n{FILE_OUTPUT_END}\n```\n\n"
        f"A report without the required {fmt.upper()} file is incomplete."
    )


# ─── Public entry point ───────────────────────────────────────────────────────

async def generate(
    report_text: str,
    task,
    params,
    driver,
    model_slug: str,
) -> tuple[list[str], dict, float]:
    """
    Generate every requested output format for a task.

    Returns (output_files, output_file_errors, fix_cost_usd).
    """
    output_files: list[str] = []
    errors: dict = {}
    total_fix_cost = 0.0

    for fmt in (task.output_formats or []):
        files, err, cost = await _generate_one(
            report_text, task, params, driver, model_slug, fmt
        )
        output_files.extend(files)
        if err:
            errors[fmt] = err
        total_fix_cost += cost

    return output_files, errors, round(total_fix_cost, 6)


async def _generate_one(
    report_text, task, params, driver, model_slug, fmt
) -> tuple[list[str], str, float]:
    fix_cost = 0.0
    max_attempts = max(1, params.file_fix_attempts)

    # Staging dir = where input files live (so generated code can read them).
    if task.file_paths and os.path.exists(task.file_paths[0]):
        staging_dir = os.path.dirname(os.path.abspath(task.file_paths[0]))
    else:
        staging_dir = task.output_dir or os.getcwd()
    out_dir = task.output_dir or staging_dir
    os.makedirs(out_dir, exist_ok=True)
    os.makedirs(staging_dir, exist_ok=True)

    script_path = os.path.join(staging_dir, f"_ipl_gen_{task.run_id}_{fmt}.py")
    chdir_header = (
        "import os as _os\n"
        "_os.chdir(_os.path.dirname(_os.path.abspath(__file__)))\n\n"
    )

    code = _extract_code(report_text, fmt)

    def _move_output() -> tuple[list[str], str]:
        produced = os.path.join(staging_dir, f"output.{fmt}")
        if not os.path.exists(produced):
            produced = _find_output_file(staging_dir, fmt, script_path)
        if produced and os.path.exists(produced):
            dest = os.path.join(out_dir, f"{task.provider}_{task.run_id}.{fmt}")
            shutil.move(produced, dest)
            logger.info("file saved → %s (%d bytes)", dest, os.path.getsize(dest))
            return [dest], ""
        return [], f"Code ran (rc=0) but output.{fmt} was not found."

    try:
        for attempt in range(1, max_attempts + 1):
            # Need code? Either none was emitted, or a prior attempt failed.
            if code is None:
                code, c = await _request_code(
                    report_text, task, params, driver, model_slug, fmt, prior_error=None
                )
                fix_cost += c
                if code is None:
                    return [], f"Model did not emit a {fmt} code block.", fix_cost

            _write_script(script_path, chdir_header + code)
            rc, stderr = _run_script(script_path, staging_dir, params.code_exec_timeout)

            if rc == 0:
                files, err = _move_output()
                return files, err, fix_cost

            logger.warning("[%s] %s gen attempt %d/%d failed: %s",
                           task.run_id, fmt, attempt, max_attempts, stderr[:200])
            if attempt == max_attempts:
                return [], f"Code failed after {max_attempts} attempts: {stderr[:500]}", fix_cost

            # Ask the model to fix it.
            code, c = await _request_code(
                report_text, task, params, driver, model_slug, fmt, prior_error=(code, stderr)
            )
            fix_cost += c
            if code is None:
                return [], f"Model did not return fixed {fmt} code.", fix_cost
    finally:
        if os.path.exists(script_path):
            try:
                os.remove(script_path)
            except OSError:
                pass

    return [], f"Unexpected exit of {fmt} generation loop.", fix_cost


# ─── Model code requests ──────────────────────────────────────────────────────

async def _request_code(
    report_text, task, params, driver, model_slug, fmt, prior_error
) -> tuple[Optional[str], float]:
    """Ask the model to produce (or fix) the file-generation code block."""
    lib = _LIB_HINT.get(fmt, fmt)
    if prior_error is None:
        prompt = (
            f"You produced this research report but no Python code block to "
            f"generate the required {fmt.upper()} file:\n\n{report_text}\n\n"
            f"Using the findings above, write ONE complete Python code block using "
            f"{lib} that saves the file as exactly output.{fmt} in the current "
            f"working directory. Wrap it with EXACTLY these sentinels:\n\n"
            f"```python\n{FILE_OUTPUT_START}\n# ... code ...\n{FILE_OUTPUT_END}\n```\n\n"
            f"Return ONLY the code block."
        )
    else:
        broken_code, stderr = prior_error
        prompt = (
            f"The Python code to generate a .{fmt} file failed with:\n\n"
            f"```\n{stderr}\n```\n\nHere is the code that failed:\n\n"
            f"```python\n{broken_code}\n```\n\n"
            f"Fix it so it runs without errors and saves output.{fmt} in the current "
            f"working directory. Wrap the fixed code with EXACTLY these sentinels:\n\n"
            f"```python\n{FILE_OUTPUT_START}\n# ... fixed code ...\n{FILE_OUTPUT_END}\n```\n\n"
            f"Return ONLY the code block."
        )

    # Codegen should be deterministic-ish, no web, no tools.
    codegen_params = params.merged({"web_search": False, "enabled_tools": []})
    try:
        resp = await driver.chat(
            model_slug=model_slug,
            messages=[{"role": "user", "content": prompt}],
            params=codegen_params,
            tools_schemas=None,
        )
    except Exception as e:  # noqa: BLE001
        logger.error("code request failed: %s", e)
        return None, 0.0

    code = _extract_code(resp.text, fmt)
    return code, resp.cost_usd


# ─── Local execution ──────────────────────────────────────────────────────────

def _write_script(path: str, content: str) -> None:
    with open(path, "w", encoding="utf-8") as f:
        f.write(content)


def _run_script(path: str, cwd: str, timeout: int) -> tuple[int, str]:
    try:
        proc = subprocess.run(
            [sys.executable, path],
            cwd=cwd, capture_output=True, text=True, timeout=timeout,
        )
        return proc.returncode, (proc.stderr or "").strip()
    except subprocess.TimeoutExpired:
        return 1, f"Execution timed out after {timeout}s"
    except Exception as e:  # noqa: BLE001
        return 1, f"Execution error: {e}"


# ─── Code extraction ──────────────────────────────────────────────────────────

_FENCE_RE = re.compile(
    r"`{3,}[ \t]*(?:python|py)?[ \t]*\n(.*?)\n?`{3,}",
    re.DOTALL | re.IGNORECASE,
)

_LIB_KEYWORDS = {
    "xlsx": ["openpyxl", "pandas", "xlsxwriter"],
    "docx": ["docx", "Document"],
    "pptx": ["pptx", "Presentation"],
}


def _extract_code(text: str, fmt: str) -> Optional[str]:
    if not text:
        return None

    # 1) Sentinel block (most reliable)
    s = text.find(FILE_OUTPUT_START)
    e = text.find(FILE_OUTPUT_END)
    if s != -1 and e != -1 and e > s:
        raw = text[s:e + len(FILE_OUTPUT_END)]
        lines = [l for l in raw.splitlines() if not re.match(r"^`{3,}", l.strip())]
        return _sanitize_code("\n".join(lines).strip())

    # 2) Any fenced block that references the right library
    hints = _LIB_KEYWORDS.get(fmt, [])
    last = None
    for m in _FENCE_RE.finditer(text):
        body = m.group(1)
        if any(h in body for h in hints):
            last = body
    if last is not None:
        return _sanitize_code(last.strip())

    return None


def _sanitize_code(code: str) -> str:
    """Strip inline citation markers that some models leak into code."""
    code = re.sub(r"(['\"][^'\"]*['\"]) *\[cite:\s*[\d,\s]+\](\s*:)", r"\1\2", code)
    code = re.sub(r"(:\s*)\[cite:\s*[\d,\s]+\]", r"\1[]", code)
    code = re.sub(r"\[cite:\s*[\d,\s]+\]", "", code)
    return code


def _find_output_file(staging_dir: str, fmt: str, exclude: str) -> Optional[str]:
    candidates = []
    for fname in os.listdir(staging_dir):
        fpath = os.path.join(staging_dir, fname)
        if fname.endswith(f".{fmt}") and fpath != exclude and os.path.isfile(fpath):
            candidates.append(fpath)
    return max(candidates, key=os.path.getmtime) if candidates else None


# ─── File text extraction (for inlining into prompts) ─────────────────────────

def extract_file_text(fpath: str, max_chars: int = 40_000) -> str:
    """Best-effort plain-text extraction for inlining a file into a prompt."""
    ext = Path(fpath).suffix.lower()

    def _truncate(t: str) -> str:
        return t if len(t) <= max_chars else t[:max_chars] + "\n\n[... truncated ...]"

    try:
        if ext in (".txt", ".md", ".csv", ".tsv", ".json", ".xml", ".html", ".htm", ".py"):
            with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                return _truncate(f.read())

        if ext == ".docx":
            from docx import Document
            doc = Document(fpath)
            return _truncate("\n".join(p.text for p in doc.paragraphs if p.text.strip()))

        if ext in (".xlsx", ".xls"):
            from openpyxl import load_workbook
            wb = load_workbook(fpath, read_only=True, data_only=True)
            parts = []
            for name in wb.sheetnames:
                ws = wb[name]
                parts.append(f"=== Sheet: {name} ===")
                for row in ws.iter_rows(max_row=500, values_only=True):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            wb.close()
            return _truncate("\n".join(parts))

        if ext == ".pdf":
            try:
                import fitz  # PyMuPDF
                doc = fitz.open(fpath)
                text = "\n\n".join(page.get_text() for page in doc)
                doc.close()
                return _truncate(text)
            except ImportError:
                from pdfminer.high_level import extract_text
                return _truncate(extract_text(fpath))

        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(fpath)
            parts = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [sh.text for sh in slide.shapes
                         if hasattr(sh, "text") and sh.text.strip()]
                if texts:
                    parts.append(f"=== Slide {i} ===\n" + "\n".join(texts))
            return _truncate("\n\n".join(parts))

        if ext in (".png", ".jpg", ".jpeg", ".gif", ".webp"):
            return f"[Image file {os.path.basename(fpath)} — not inlined as text]"

    except ImportError as e:
        return f"[{os.path.basename(fpath)} — extraction needs a missing library: {e}]"
    except Exception as e:  # noqa: BLE001
        return f"[Error extracting {os.path.basename(fpath)}: {e}]"

    return f"[{os.path.basename(fpath)} — unsupported format {ext}]"
