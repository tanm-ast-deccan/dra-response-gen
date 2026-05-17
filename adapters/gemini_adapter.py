"""
gemini_adapter.py — Gemini Deep Research adapter (rewrite).

Architecture
------------
Stage 1 (always):
    1. _prepare_files() — convert non-native formats and upload everything
       via the Gemini Files API so the DRA sees all data, no truncation.
    2. interactions.create(input=[prompt_part + file_parts],
                           agent=DRA, background=True)
    3. Poll until completed / failed / stale.
    4. Extract report text + citations from interaction.outputs[-1].

Stage 2 (only when output_formats is set):
    5. Extract the Python code block the DRA embedded in its report.
    6. Execute it locally in the staging dir (where input files live on disk).
    7. Move the output file to task.output_files_dir.

File format strategy — NO truncation, ever
-------------------------------------------
    PDF                → Files API upload, native understanding (application/pdf)
    CSV                → Files API upload (text/csv)
    TXT / MD / PY /
    JSON / XML         → Files API upload (text/plain)
    XLSX (multi-tab)   → convert ALL sheets to labelled TSV text
                         → Files API upload (text/plain)
                         Original .xlsx stays on disk for local code execution.
    DOCX               → extract full text via python-docx
                         → Files API upload (text/plain)
                         Original .docx stays on disk for local code execution.
    PPTX               → extract full text via python-pptx
                         → Files API upload (text/plain)
                         Original .pptx stays on disk for local code execution.

Why originals stay on disk:
    The DRA sees the text content of every file.  When it writes Python code
    to generate the output file it will reference the original filenames
    (e.g. openpyxl.load_workbook('network_flow_model.xlsx')).  Those files
    ARE present in the staging dir, so local execution works correctly.

Comparison with Claude and o3
------------------------------
    Aspect               Claude          o3              Gemini (this)
    ─────────────────── ─────────────── ─────────────── ──────────────────────
    Code written by      Model           Model           Model (in report)
    Input files in exec  Yes (sandbox)   Yes (container) Yes (local disk)
    Stage 2 API call     No              No              Yes (on error, up to 3 fix calls)
    File content quality Real computed   Real computed   Real computed
    Exec environment     Anthropic sbx   OpenAI ctr      DGX Spark local

Required packages in conda env for local code execution
---------------------------------------------------------
    openpyxl      — reading + writing .xlsx (multi-tab)
    pandas        — data manipulation, CSV I/O, XLSX I/O via openpyxl engine
    python-docx   — reading + writing .docx   (pip: python-docx, import: docx)
    python-pptx   — reading + writing .pptx   (pip: python-pptx, import: pptx)
    reportlab     — PDF generation if needed
    matplotlib    — charts / plots embedded in files
    numpy         — numerical computation
    scipy         — scientific computation

    Verify on Spark:
        conda run -n adobe python -c "
            import openpyxl, pandas, docx, pptx, reportlab, matplotlib, numpy"
"""

from __future__ import annotations

import os
import re
import io
import time
import shutil
import asyncio
import logging
import subprocess
import tempfile
from datetime import datetime, timezone
from typing import Optional

from models import ResearchTask, AgentResult

_log = logging.getLogger("dra.gemini")


# ─── Constants ────────────────────────────────────────────────────────────────

DEEP_RESEARCH_AGENT       = "deep-research-pro-preview-12-2025"
POLL_INTERVAL_SECONDS     = 10
STALE_THRESHOLD_SECONDS   = 1800
CODE_EXEC_TIMEOUT_SECONDS = 120    # local subprocess timeout

# Hard limit on total inline file content before the API call.
# 1M token context = ~4M chars total. We reserve ~250K tokens for the
# prompt text and model output, leaving ~750K tokens = ~3M chars for files.
# Truncation is NEVER allowed — hitting this raises immediately so the
# failure is explicit in AgentResult and visible to the benchmark scorer.
MAX_INLINE_CHARS = 3_000_000

# Sentinels the DRA must use to wrap its file-generation code.
# Using explicit start/end markers instead of markdown fences makes
# extraction robust against language-tag variations (```python vs ```Python
# vs ```) and non-code content mixed around the block.
FILE_OUTPUT_MARKER = "GEMINI_FILE_OUTPUT"
FILE_OUTPUT_START  = "# GEMINI_FILE_OUTPUT_START"
FILE_OUTPUT_END    = "# GEMINI_FILE_OUTPUT_END"


def _pdf_to_text(fpath: str) -> str:
    """
    Extract text from a PDF for inline embedding.

    For large PDFs (annual reports, legal documents), applies smart section
    filtering to drop non-analytical content (statutory reports, corporate
    governance, CSR, secretarial audit, directors report boilerplate) while
    preserving all financially relevant content (MD&A, financial statements,
    notes to accounts).

    This keeps token counts within Gemini DRA's 131K input limit for typical
    Indian annual reports without losing any analytically relevant content.

    Drop logic is keyword-based and page-range aware — no hardcoded page numbers.
    Falls back to full extraction if PyMuPDF is unavailable.
    """
    try:
        import fitz  # PyMuPDF

        doc = fitz.open(fpath)
        total_pages = len(doc)

        # ── Smart section filtering for large PDFs ────────────────────────
        # Only apply if document is large enough to risk token limit
        # Threshold: ~400K chars (~100K tokens) — safe margin below 131K
        LARGE_PDF_THRESHOLD = 400_000

        # Quick size estimate from first 10 pages
        sample_chars = sum(len(doc[i].get_text())
                          for i in range(min(10, total_pages)))
        estimated_total = sample_chars * (total_pages / min(10, total_pages))

        if estimated_total > LARGE_PDF_THRESHOLD:
            pages_text = _pdf_smart_extract(doc, fpath)
        else:
            pages_text = [page.get_text() for page in doc]

        doc.close()
        return "\n\n".join(p for p in pages_text if p.strip()) \
               or "[PDF: no extractable text]"

    except ImportError:
        pass
    except Exception as e:
        print(f"[Gemini] PyMuPDF error on {os.path.basename(fpath)}: {e}")

    try:
        from pdfminer.high_level import extract_text as _pdfminer_extract
        text = _pdfminer_extract(fpath)
        return text.strip() or "[PDF: no extractable text]"
    except ImportError:
        pass
    except Exception as e:
        print(f"[Gemini] pdfminer error on {os.path.basename(fpath)}: {e}")

    size = os.path.getsize(fpath)
    return (
        f"[PDF: {os.path.basename(fpath)}, {size} bytes — text extraction "
        f"unavailable. Install: pip install pymupdf]"
    )


def _pdf_smart_extract(doc, fpath: str) -> list[str]:
    """
    Smart PDF extraction that drops non-analytical statutory sections.

    Strategy:
    1. Scan all pages and classify each as KEEP or DROP based on keywords
    2. DROP sections: Corporate Governance, CSR, Secretarial Audit,
       Directors Report boilerplate, Board committee reports
    3. KEEP sections: Business overview, MD&A, Financial Statements,
       Notes to Accounts, Key financial highlights
    4. If dropping sections still leaves too much text, drop auditor
       report boilerplate (standard opinion text) as well

    Works for Indian annual reports (SEBI format) without hardcoded page numbers.
    """
    import os

    # Keywords that indicate DROP sections (statutory boilerplate)
    DROP_KEYWORDS = [
        'secretarial audit',
        'form no. mr-',
        'annexure-a',
        'board of directors\nthe following',
        'corporate governance report',
        'certificate on corporate governance',
        'csr committee',
        'annual report on csr',
        'related party transactions',
        'dividend distribution policy',
    ]

    # Keywords that indicate KEEP sections (financial content)
    KEEP_KEYWORDS = [
        'financial statements',
        'balance sheet',
        'profit and loss',
        'cash flow',
        'notes to standalone',
        'notes to consolidated',
        'statement of changes in equity',
        'management discussion',
        'independent auditor',
        'key audit matters',
    ]

    # Classify each page
    page_classes = []
    for i in range(len(doc)):
        text_lower = doc[i].get_text().lower()[:500]

        is_drop = any(kw in text_lower for kw in DROP_KEYWORDS)
        is_keep = any(kw in text_lower for kw in KEEP_KEYWORDS)

        if is_keep:
            page_classes.append('KEEP')
        elif is_drop:
            page_classes.append('DROP')
        else:
            page_classes.append('NEUTRAL')

    # Propagate DROP classification — if a page is in a DROP section,
    # neutral pages following it are also dropped until a KEEP page appears
    final_classes = list(page_classes)
    in_drop_section = False
    for i in range(len(final_classes)):
        if final_classes[i] == 'DROP':
            in_drop_section = True
        elif final_classes[i] == 'KEEP':
            in_drop_section = False
        if in_drop_section and final_classes[i] == 'NEUTRAL':
            final_classes[i] = 'DROP'

    kept = sum(1 for c in final_classes if c != 'DROP')
    dropped = sum(1 for c in final_classes if c == 'DROP')
    kept_chars = sum(len(doc[i].get_text())
                    for i in range(len(doc)) if final_classes[i] != 'DROP')

    print(f"[Gemini] Smart PDF extraction: keeping {kept} pages, "
          f"dropping {dropped} pages, "
          f"~{kept_chars:,} chars (~{kept_chars//4:,} tokens) "
          f"[{os.path.basename(fpath)}]")

    return [doc[i].get_text() if final_classes[i] != 'DROP' else ''
            for i in range(len(doc))]


class StaleInteractionError(RuntimeError):
    pass


# ─── Pricing ──────────────────────────────────────────────────────────────────

PRICING = {
    "deep-research-pro-preview-12-2025": {
        "input_per_mtok":        2.00,
        "output_per_mtok":      12.00,
        "cached_input_per_mtok": 0.50,
    },
    "gemini-2.0-flash-deep-research": {
        "input_per_mtok":        0.10,
        "output_per_mtok":       0.40,
        "cached_input_per_mtok": 0.025,
    },
}


def estimate_cost(
    model: str,
    input_tokens: int,
    output_tokens: int,
    cached_tokens: int = 0,
) -> float:
    prices = None
    for key in PRICING:
        if model.startswith(key) or key.startswith(model):
            prices = PRICING[key]
            break
    if prices is None:
        prices = PRICING[DEEP_RESEARCH_AGENT]
    uncached = max(0, input_tokens - cached_tokens)
    return round(
        (uncached / 1_000_000)      * prices["input_per_mtok"]
        + (cached_tokens / 1_000_000) * prices["cached_input_per_mtok"]
        + (output_tokens / 1_000_000) * prices["output_per_mtok"],
        6,
    )


# ─── File conversion helpers ──────────────────────────────────────────────────

def _libreoffice_evaluate(fpath: str) -> str:
    """
    Use LibreOffice headless to evaluate all formulas in an xlsx file
    and return the path to a new xlsx with cached computed values.

    This solves the openpyxl data_only=True problem where formula cells
    return None if the file was never saved with cached values in Excel.
    LibreOffice opens, computes all formulas, and re-saves with values.

    Returns the path to the evaluated file (in a temp dir).
    Falls back to original path if LibreOffice is unavailable or fails.
    """
    import subprocess
    import shutil
    import tempfile

    libreoffice_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not libreoffice_bin:
        return fpath  # fallback: no LibreOffice available

    try:
        tmp_dir = tempfile.mkdtemp(prefix="lo_eval_")
        result = subprocess.run(
            [
                libreoffice_bin,
                "--headless",
                "--convert-to", "xlsx",
                "--outdir", tmp_dir,
                fpath,
            ],
            capture_output=True,
            text=True,
            timeout=60,
        )
        if result.returncode != 0:
            print(f"[Gemini] LibreOffice eval failed for {os.path.basename(fpath)}: "
                  f"{result.stderr.strip()[:200]}")
            return fpath  # fallback to original

        # LibreOffice saves as <original_stem>.xlsx in tmp_dir
        stem = os.path.splitext(os.path.basename(fpath))[0]
        evaluated = os.path.join(tmp_dir, f"{stem}.xlsx")
        if os.path.exists(evaluated):
            print(f"[Gemini] LibreOffice evaluated formulas: {os.path.basename(fpath)}")
            return evaluated
        else:
            return fpath  # fallback

    except Exception as e:
        print(f"[Gemini] LibreOffice eval error for {os.path.basename(fpath)}: {e}")
        return fpath  # fallback to original


def _xlsx_to_text(fpath: str) -> str:
    """
    Convert ALL sheets of an xlsx workbook to labelled TSV text.

    Multi-tab workbooks are fully preserved — each sheet gets a clear
    header so the DRA knows which data belongs to which sheet.
    Empty rows are dropped; cell values are coerced to strings.

    Formula evaluation: LibreOffice headless is used to pre-evaluate
    all formulas before extraction, ensuring computed values are present
    even for workbooks saved without cached values (e.g. template files).
    """
    try:
        import openpyxl
    except ImportError:
        raise ImportError("pip install openpyxl  (needed for XLSX conversion)")

    # Pre-evaluate formulas via LibreOffice so data_only=True returns values
    evaluated_path = _libreoffice_evaluate(fpath)

    wb = openpyxl.load_workbook(evaluated_path, data_only=True)
    sections = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        all_rows = list(ws.iter_rows(values_only=True))

        # Detect actual header row (0-based) — first row with ≥2 non-null cells
        header_idx = 0
        for i, row in enumerate(all_rows):
            if sum(1 for c in row if c is not None) >= 2:
                header_idx = i
                break

        # Keep only rows from header onwards, drop fully empty rows
        rows = []
        for row in all_rows[header_idx:]:
            if any(cell is not None for cell in row):
                # Strip whitespace from all cell values to avoid silent mismatches
                rows.append(
                    "\t".join("" if cell is None else str(cell).strip() for cell in row)
                )

        if rows:
            sections.append(
                f"=== Sheet: {sheet_name} ({len(rows)} rows, "
                f"use header={header_idx} when reading with pandas) ===\n"
                + "\n".join(rows)
            )
    return "\n\n".join(sections) if sections else "[Empty workbook]"


def _docx_to_text(fpath: str) -> str:
    """Extract full text from a .docx file, preserving paragraph structure."""
    try:
        from docx import Document
    except ImportError:
        raise ImportError("pip install python-docx  (needed for DOCX conversion)")
    doc = Document(fpath)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract table content
    for table in doc.tables:
        for row in table.rows:
            cell_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
            if cell_texts:
                paragraphs.append("\t".join(cell_texts))
    return "\n".join(paragraphs) if paragraphs else "[Empty document]"


def _pptx_to_text(fpath: str) -> str:
    """Extract full text from a .pptx file, slide by slide."""
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx  (needed for PPTX conversion)")
    prs = Presentation(fpath)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"=== Slide {i} ===\n" + "\n".join(texts))
    return "\n\n".join(slides) if slides else "[Empty presentation]"


# ─── The adapter ──────────────────────────────────────────────────────────────

class GeminiAdapter:
    """
    Runs a deep research task using the Gemini Deep Research agent
    via the Interactions API (google.genai SDK >= 1.55.0).
    """

    def __init__(
        self,
        api_key: Optional[str] = None,
        model: str = DEEP_RESEARCH_AGENT,
        dry_run: bool = False,
        use_context_cache: bool = False,
        max_research_passes: Optional[int] = 4,
    ):
        self.api_key              = api_key or os.environ.get("GOOGLE_API_KEY")
        self.model                = model
        self.dry_run              = dry_run
        self.use_context_cache    = use_context_cache
        self.max_research_passes  = max_research_passes  # None = unlimited

        if not dry_run:
            try:
                from google import genai
                from google.genai import types
                self._client = genai.Client(api_key=self.api_key)
                self._types  = types
            except ImportError:
                raise ImportError(
                    "pip install -U google-genai  (>= 1.55.0 required)"
                )

    # ─── Public entry point ───────────────────────────────────────────

    async def run(self, task: ResearchTask) -> AgentResult:
        started_at = datetime.now(timezone.utc)
        print(f"[Gemini] Starting {task.task_id} "
              f"(model={self.model}, dry_run={self.dry_run})")

        if self.dry_run:
            return self._dry_run_result(task, started_at)

        try:
            # Stage 1 — Multi-pass Research
            uploaded_files = await self._prepare_files(task)

            n_native   = sum(1 for f in uploaded_files if f.get("uri"))
            n_fallback = sum(1 for f in uploaded_files if f.get("inline_text"))
            print(f"[Gemini] Files: {n_native} native uploads, "
                  f"{n_fallback} inline fallbacks")

            total_input_tokens  = 0
            total_output_tokens = 0
            total_cached_tokens = 0
            previous_id  = None
            interaction  = None

            # pass1_report: research text from Pass 1. This is ALWAYS what gets
            #               saved as result.response_text — never the code.
            # pass2_text:   output from Pass 2 (code generation). Used only as
            #               the source for code extraction in Stage 2.
            # pass2_ran:    True if Pass 2 actually executed.
            pass1_report = ""
            pass2_text   = ""
            pass2_ran    = False

            # Multi-pass only when file output is required.
            # Pass 2's sole purpose is code generation — text-only tasks
            # never need it since Pass 1 already produces the full response.
            effective_passes = (
                (self.max_research_passes or 999)
                if task.output_formats else 1
            )

            for pass_num in range(1, effective_passes + 1):

                # Budget guard
                current_cost = estimate_cost(
                    self.model, total_input_tokens,
                    total_output_tokens, total_cached_tokens
                )
                if pass_num > 1 and current_cost > task.max_cost_usd:
                    print(f"[Gemini] Budget ${task.max_cost_usd} exceeded "
                          f"(${current_cost:.2f}) — stopping at pass {pass_num - 1}")
                    break

                # Timeout guard
                elapsed = (datetime.now(timezone.utc) - started_at).total_seconds()
                if pass_num > 1 and elapsed > task.timeout_seconds:
                    print(f"[Gemini] Timeout {task.timeout_seconds}s exceeded "
                          f"— stopping at pass {pass_num - 1}")
                    break

                # Skip Pass 2 if Pass 1 already produced a code block.
                # No need for a second API call if the code is already there.
                if pass_num > 1 and task.output_formats:
                    fmt_check = task.output_formats[0]
                    has_code = (
                        _extract_marked_code_block(pass1_report, FILE_OUTPUT_MARKER)
                        is not None
                        or _extract_last_file_code_block(pass1_report, fmt_check)
                        is not None
                    )
                    if has_code:
                        print(f"[Gemini] Pass 1 already contains code block — "
                              f"skipping Pass {pass_num}")
                        break

                # Build input
                if pass_num == 1:
                    input_payload = self._build_input(task, uploaded_files)
                    print(f"[Gemini] Research pass 1 "
                          f"({'closed' if task.is_closed else 'open'} corpus, "
                          f"{len(uploaded_files)} file(s))...")
                else:
                    input_payload = [{"type": "text",
                                      "text": self._continuation_prompt(
                                          pass_num, task, pass1_report
                                      )}]
                    print(f"[Gemini] Code generation pass {pass_num} "
                          f"(pass1 report: {len(pass1_report):,} chars)...")

                create_kwargs = dict(
                    input=input_payload,
                    agent=self.model,
                    background=True,
                    store=True,
                )
                if previous_id is not None:
                    create_kwargs["previous_interaction_id"] = previous_id

                interaction = await asyncio.to_thread(
                    self._client.interactions.create,
                    **create_kwargs,
                )
                print(f"[Gemini] Pass {pass_num} started: {interaction.id}")
                interaction = await self._poll(interaction, task)

                # Accumulate token counts across all passes
                usage = getattr(interaction, "usage", None)
                if usage:
                    total_input_tokens  += getattr(usage, "total_input_tokens",  0) or 0
                    total_output_tokens += getattr(usage, "total_output_tokens", 0) or 0
                    total_cached_tokens += getattr(usage, "total_cached_tokens", 0) or 0

                # Capture text from each pass separately.
                # Pass 1 → pass1_report (research, always the saved response).
                # Pass 2 → pass2_text  (code, used only for code extraction).
                pass_outputs = getattr(interaction, "outputs", [])
                pass_text = (
                    getattr(pass_outputs[-1], "text", "") or ""
                    if pass_outputs else ""
                )
                if pass_num == 1:
                    pass1_report = pass_text
                else:
                    pass2_text = pass_text
                    pass2_ran  = True

                previous_id = interaction.id

            # Extract metadata (tokens, cost, timestamps) from the final interaction.
            # _extract_result may set response_text to code (if Pass 2 ran last) —
            # we unconditionally override it with pass1_report so the research
            # report is always what gets saved, never the code.
            result = self._extract_result(
                interaction, task, started_at,
                total_input_tokens, total_output_tokens, total_cached_tokens,
            )
            # Strip the code block from pass1_report before saving as response_text.
            # When Pass 1 includes code inline, pass1_report = research + code block.
            # The code block is extracted separately by _run_code_block — it must
            # not appear in the saved research response.
            clean_report = (
                _strip_file_output_block(pass1_report)
                if task.output_formats else pass1_report
            )
            result.response_text = clean_report
            result.completed     = bool(clean_report)
            result.error         = None if clean_report else "Empty response from Pass 1"

            # Stage 2 — File generation (local execution with self-correction).
            # code_source: text to scan for the code block.
            #   Pass 2 ran  → use pass2_text (dedicated code-gen output)
            #   No Pass 2   → use pass1_report (Gemini included code inline in Pass 1)
            if (task.output_formats
                    and task.output_files_dir
                    and result.completed):
                code_source = pass2_text if pass2_ran else pass1_report
                output_files, output_file_errors, fix_cost = \
                    await self._run_code_block(
                        code_source, task, previous_id, pass1_report
                    )
                result.output_files       = output_files
                result.output_file_errors = output_file_errors
                result.total_cost_usd     = round(
                    result.total_cost_usd + fix_cost, 6
                )
                if output_file_errors:
                    result.completed = False
                    result.error = (
                        f"File generation failed for: "
                        f"{list(output_file_errors.keys())}"
                    )

            return result

        except Exception as e:
            _log.exception("[%s] Gemini adapter error", task.task_id)
            print(f"[Gemini] Error: {e}")
            return self._error_result(task, started_at, str(e))

    # ─── Continuation prompt ──────────────────────────────────────────

    def _continuation_prompt(
        self,
        pass_num: int,
        task: "ResearchTask",
        pass1_report: str = "",
    ) -> str:
        """
        Build a continuation prompt for research pass N (N >= 2).

        Final pass + output_formats set: skip synthesis prose, go straight to
        code generation. Includes the Pass 1 research report inline so Gemini
        writes code from actual findings — previous_interaction_id alone does
        NOT carry research context across interactions.

        Final pass + no output_formats: synthesize research into final report.

        Intermediate pass: identify gaps and conduct targeted research.

        Unlimited mode (max_research_passes=None): always intermediate.
        """
        is_final = (
            self.max_research_passes is not None
            and pass_num == self.max_research_passes
        )
        if is_final and task.output_formats:
            fmt = task.output_formats[0]
            lib = {
                "docx": "python-docx (from docx import Document)",
                "xlsx": "openpyxl or pandas with openpyxl engine",
                "pptx": "python-pptx (from pptx import Presentation)",
            }.get(fmt, fmt)
            report_section = (
                f"\n\nHere is the research report you produced:\n\n"
                f"{pass1_report}\n\n"
                if pass1_report else ""
            )
            return (
                f"Your research is complete. Do not conduct any new searches "
                f"or add any more analysis.{report_section}"
                f"Using the research findings above, write ONLY a Python code "
                f"block using {lib} that generates output.{fmt}. "
                f"The code must save the file as exactly: output.{fmt} "
                f"in the current working directory.\n\n"
                f"Wrap the code with EXACTLY these sentinel lines:\n\n"
                f"```python\n"
                f"# GEMINI_FILE_OUTPUT_START\n"
                f"# ... your code ...\n"
                f"# GEMINI_FILE_OUTPUT_END\n"
                f"```\n\n"
                f"Return nothing else — no text before or after the code block. "
                f"Do NOT include [cite: N] markers inside the code."
            )
        if is_final:
            return (
                "Stop all research now. Do not conduct any new web searches. "
                "Using only the information you have already gathered in this "
                "conversation, write your final synthesized answer directly. "
                "Do not describe how to write a report — write the actual "
                "answer to the original research question. Resolve any "
                "contradictions in your findings, fill remaining gaps where "
                "possible from existing gathered data, and ensure all claims "
                "are cited."
            )
        return (
            f"You have completed research pass {pass_num - 1}. "
            "Review your findings so far and identify the most critical gaps: "
            "missing quantitative data, unverified claims, or areas needing "
            "deeper analysis. Conduct targeted additional research to address "
            "these specific gaps and update your report with the new findings."
        )

    # ─── File preparation ─────────────────────────────────────────────

    # MIME type mapping for Files API native uploads
    _MIME_MAP = {
        ".pdf":  "application/pdf",
        ".csv":  "text/csv",
        ".txt":  "text/plain",
        ".md":   "text/plain",
        ".json": "application/json",
        ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ".xls":  "application/vnd.ms-excel",
        ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ".doc":  "application/msword",
        ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ".ppt":  "application/vnd.ms-powerpoint",
    }

    async def _prepare_files(
        self, task: ResearchTask
    ) -> list[dict]:
        """
        Upload every input file to the Gemini Files API and return a list
        of file metadata dicts: {basename, uri, mime_type, fpath}.

        Using native Files API uploads instead of inline text extraction:
        - PDFs are understood natively by Gemini (no text extraction loss)
        - XLSXs are read with formula evaluation by Gemini's backend
        - No token limit issues from bloated text conversions
        - No data loss from openpyxl data_only=True on formula-only workbooks

        Original files remain on disk in the staging dir for local code
        execution during Stage 2 file generation.
        """
        uploaded: list[dict] = []

        for fpath in task.file_paths:
            if not os.path.exists(fpath):
                print(f"[Gemini] File not found, skipping: {fpath}")
                continue

            basename = os.path.basename(fpath)
            ext      = os.path.splitext(fpath)[1].lower()
            mime     = self._MIME_MAP.get(ext, "application/octet-stream")

            try:
                if ext == ".pdf":
                    # PDFs: upload natively — Interactions API supports application/pdf
                    file_obj = await asyncio.to_thread(
                        self._client.files.upload,
                        file=fpath,
                        config={"mime_type": "application/pdf",
                                "display_name": basename},
                    )
                    uploaded.append({
                        "basename":  basename,
                        "uri":       file_obj.uri,
                        "mime_type": "application/pdf",
                        "fpath":     fpath,
                    })
                    size_kb = os.path.getsize(fpath) / 1024
                    print(f"[Gemini] {basename} → Files API upload "
                          f"({size_kb:.1f} KB, native PDF)")
                else:
                    # Non-PDF: inline text extraction
                    # (xlsx, docx, pptx, csv — not supported as native document type)
                    text = self._extract_file_text(fpath, ext, basename)
                    uploaded.append({
                        "basename":    basename,
                        "uri":         None,
                        "mime_type":   None,
                        "fpath":       fpath,
                        "inline_text": text,
                    })
                    print(f"[Gemini] {basename} → inline text "
                          f"({len(text):,} chars)")

            except Exception as e:
                print(f"[Gemini] Failed to process {basename}: {e}")

        return uploaded

    def _extract_file_text(self, fpath: str, ext: str, basename: str) -> str:
        """Fallback inline text extraction when Files API upload fails."""
        if ext in (".xlsx", ".xls"):
            return f"[SPREADSHEET: {basename}]\n" + _xlsx_to_text(fpath)
        elif ext in (".docx", ".doc"):
            return f"[DOCUMENT: {basename}]\n" + _docx_to_text(fpath)
        elif ext in (".pptx", ".ppt"):
            return f"[PRESENTATION: {basename}]\n" + _pptx_to_text(fpath)
        elif ext == ".pdf":
            return f"[PDF: {basename}]\n" + _pdf_to_text(fpath)
        else:
            with open(fpath, "r", encoding="utf-8", errors="replace") as f:
                return f"[FILE: {basename}]\n" + f.read()

    # ─── Prompt construction ──────────────────────────────────────────

    def _build_input(
        self,
        task: ResearchTask,
        uploaded_files: list[dict],
    ) -> list:
        """
        Build the interactions.create input list using native Files API parts.

        Structure:
          [
            Part(text=<prompt + constraints + filename list>),
            Part(file_data=FileData(file_uri=..., mime_type=...)),  # per uploaded file
            Part(text=<inline fallback text>),                       # if upload failed
            Part(text=<file generation instructions>),               # if output_formats set
          ]

        For IAT-1: prepend closed-corpus constraint in text part.
        For output_formats: append file generation instructions as text part.
        """
        parts = []

        # ── Build main text part (prompt + constraints + filename list) ──
        sections = []

        if task.is_closed:
            sections.append(
                "[CONSTRAINT — CLOSED CORPUS]\n"
                "This is an IAT-1 closed-book evaluation task.\n"
                "You MUST NOT use web search or any external sources.\n"
                "Base your entire analysis ONLY on the provided files.\n"
            )
            print(f"[Gemini] IAT-1 (Closed): web search disabled via instruction")

        sections.append(task.prompt)

        if task.file_paths:
            file_list = "\n".join(
                f"  - {os.path.basename(p)}" for p in task.file_paths
            )
            sections.append(
                "\n[AVAILABLE CORPUS FILES]\n"
                "The following files are available for your analysis "
                "and in the working directory for any code you write:\n"
                + file_list
            )

        parts.append({"type": "text", "text": "\n".join(sections)})

        # ── Add file parts (native upload or inline fallback) ──
        for f in uploaded_files:
            if f.get("uri"):
                # Native document part — Interactions API only supports PDF natively
                # Other formats (xlsx, docx, pptx, csv) fall back to inline text
                parts.append({
                    "type":      "document",
                    "uri":       f["uri"],
                    "mime_type": "application/pdf",
                })
            elif f.get("inline_text"):
                # Inline text for non-PDF files or failed uploads
                parts.append({
                    "type": "text",
                    "text": f"\n--- {f['basename']} ---\n{f['inline_text']}\n",
                })

        # ── File generation instructions ──
        if task.output_formats and task.file_paths:
            file_list = "\n".join(
                f"  - {os.path.basename(p)}" for p in task.file_paths
            )
            fmt = task.output_formats[0]

            if fmt in ('pptx', 'docx'):
                approach_instruction = (
                    f"IMPORTANT — HOW TO WRITE YOUR CODE:\n"
                    f"Your research above already contains all computed values, "
                    f"findings, and recommendations. Your code must write those "
                    f"findings DIRECTLY into the output file — do NOT re-read or "
                    f"re-process the input files in your code. Do NOT use data "
                    f"science libraries (sklearn, statsmodels, pulp, scipy) — the "
                    f"analysis is already done. Simply format and write your "
                    f"conclusions into the {fmt.upper()} file.\n"
                )
            else:  # xlsx
                approach_instruction = (
                    f"IMPORTANT — HOW TO WRITE YOUR CODE:\n"
                    f"You may read from the input files listed above using pandas "
                    f"or openpyxl only. Do NOT use sklearn, statsmodels, pulp, "
                    f"scipy or other analysis libraries — use only standard "
                    f"file-handling libraries. Use the exact filenames listed above, "
                    f"including the correct header= parameter when reading with pandas.\n"
                )

            file_gen_text = (
                f"\n\n{'='*60}\n"
                f"FILE GENERATION REQUIREMENT\n"
                f"{'='*60}\n\n"
                f"After completing your full analysis you MUST append a "
                f"Python code block at the END of your report.\n\n"
                f"The following input files are available in your working directory:\n"
                f"{file_list}\n\n"
                f"Required output format: {fmt}\n\n"
                f"{approach_instruction}\n"
                f"Your code MUST be wrapped with EXACTLY these sentinel lines:\n\n"
                f"```python\n"
                f"# GEMINI_FILE_OUTPUT_START\n"
                f"import openpyxl  # or relevant library\n"
                f"# ... your complete, self-contained code ...\n"
                f"# save as output.{fmt}\n"
                f"# GEMINI_FILE_OUTPUT_END\n"
                f"```\n\n"
                f"Rules:\n"
                f"1. The sentinels # GEMINI_FILE_OUTPUT_START and "
                f"# GEMINI_FILE_OUTPUT_END must appear as the first and last "
                f"lines of code inside the block.\n"
                f"2. Save the output file as exactly: output.{fmt}\n"
                f"3. Use standard libraries only:\n"
                f"   - xlsx: openpyxl or pandas (with openpyxl engine)\n"
                f"   - docx: from docx import Document\n"
                f"   - pptx: from pptx import Presentation\n"
                f"4. The code block must contain ONLY valid Python. Do NOT "
                f"include citation markers like [cite: N] — they cause "
                f"SyntaxError and will fail execution.\n\n"
                f"Do NOT skip this step. A report without the required file "
                f"is considered incomplete."
            )
            parts.append({"type": "text", "text": file_gen_text})



        return parts

    # ─── Polling ──────────────────────────────────────────────────────

    async def _poll(self, interaction, task: ResearchTask):
        """Poll the interaction until completed, failed, or stale."""
        last_status          = None
        last_status_changed  = time.monotonic()

        while True:
            interaction = await asyncio.to_thread(
                self._client.interactions.get, interaction.id
            )
            status = getattr(interaction, "status", "unknown")

            if status != last_status:
                last_status         = status
                last_status_changed = time.monotonic()
                _log.info("[%s] Gemini status → %s", task.task_id, status)
                print(f"[Gemini] Status: {status}")

            if status == "completed":
                print(f"[Gemini] Completed: {interaction.id}")
                return interaction

            if status == "failed":
                raise RuntimeError(
                    f"Gemini Deep Research failed: "
                    f"{getattr(interaction, 'error', 'unknown error')}"
                )

            stale_for = time.monotonic() - last_status_changed
            if stale_for >= max(STALE_THRESHOLD_SECONDS, task.timeout_seconds):
                raise StaleInteractionError(
                    f"Interaction stuck in '{status}' for "
                    f"{stale_for / 60:.1f}m — aborting"
                )

            await asyncio.sleep(POLL_INTERVAL_SECONDS)

    # ─── Result extraction ────────────────────────────────────────────

    def _extract_result(self, interaction, task: ResearchTask,
                        started_at,
                        total_input_tokens: int = 0,
                        total_output_tokens: int = 0,
                        total_cached_tokens: int = 0) -> AgentResult:
        completed_at = datetime.now(timezone.utc)

        report_text = ""
        try:
            outputs = getattr(interaction, "outputs", [])
            if outputs:
                report_text = getattr(outputs[-1], "text", "") or ""
        except Exception as e:
            print(f"[Gemini] Error extracting text: {e}")

        citations = self._extract_citations(interaction)

        # Use accumulated tokens if provided, else read from interaction
        if total_input_tokens == 0 and total_output_tokens == 0:
            try:
                usage = getattr(interaction, "usage", None)
                if usage:
                    total_input_tokens  = getattr(usage, "total_input_tokens",  0) or 0
                    total_output_tokens = getattr(usage, "total_output_tokens", 0) or 0
                    total_cached_tokens = getattr(usage, "total_cached_tokens", 0) or 0
            except Exception:
                pass

        total_cost = estimate_cost(
            self.model, total_input_tokens, total_output_tokens, total_cached_tokens
        )
        completed = bool(report_text)
        error     = None if completed else "Empty response from Interactions API"

        print(f"[Gemini] Done. "
              f"tokens=({total_input_tokens}in [{total_cached_tokens} cached], "
              f"{total_output_tokens}out), "
              f"cost=${total_cost:.4f}, citations={len(citations)}")

        return AgentResult(
            task_id=task.task_id,
            agent="gemini",
            model=self.model,
            response_text=report_text,
            citations=citations,
            tool_call_log=[],
            input_tokens=total_input_tokens,
            output_tokens=total_output_tokens,
            total_cost_usd=total_cost,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=1,
            completed=completed,
            forced_stop=False,
            error=error,
            output_files=[],
            output_file_errors={},
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    def _extract_citations(self, interaction) -> list[dict]:
        """
        Extract citations from the completed interaction.

        The Deep Research agent embeds citations inline in the report text
        as numbered references:
            [1] Title — https://example.com/...
        or as markdown links in the text body.

        The Interactions API does not return grounding_metadata on the
        interaction object for DRA responses — citations live in the text.
        """
        citations  = []
        seen_urls  = set()
        try:
            report_text = ""
            outputs = getattr(interaction, "outputs", [])
            if outputs:
                report_text = getattr(outputs[-1], "text", "") or ""

            if not report_text:
                return citations

            import re
            # Pattern 1: [N] Title — URL  (reference list format)
            ref_pattern = re.compile(
                r'\[(\d+)\]\s+(.+?)\s+[—–-]+\s+(https?://\S+)',
                re.MULTILINE
            )
            for m in ref_pattern.finditer(report_text):
                url = m.group(3).rstrip('.,)')
                if url not in seen_urls:
                    citations.append({
                        "url":     url,
                        "title":   m.group(2).strip(),
                        "snippet": "",
                    })
                    seen_urls.add(url)

            # Pattern 2: [Title](URL) markdown links
            md_pattern = re.compile(r'\[([^\]]+)\]\((https?://[^\)]+)\)')
            for m in md_pattern.finditer(report_text):
                url = m.group(2)
                if url not in seen_urls:
                    citations.append({
                        "url":     url,
                        "title":   m.group(1).strip(),
                        "snippet": "",
                    })
                    seen_urls.add(url)

            # Pattern 3: bare URLs
            if not citations:
                url_pattern = re.compile(r'https?://\S{10,}')
                for m in url_pattern.finditer(report_text):
                    url = m.group(0).rstrip('.,)')
                    if url not in seen_urls:
                        citations.append({"url": url, "title": "", "snippet": ""})
                        seen_urls.add(url)

        except Exception as e:
            print(f"[Gemini] Error extracting citations: {e}")
        return citations

    # ─── Local code execution (Stage 2) ──────────────────────────────

    async def _run_code_block(
        self,
        report_text: str,
        task: ResearchTask,
        previous_id: Optional[str] = None,
        pass1_report: str = "",
    ) -> tuple[list[str], dict, float]:
        """
        Extract the Python code block the DRA wrote in its report and
        execute it locally in the staging directory where input files live.

        pass1_report: the research text from Pass 1. Pasted directly into the
        generation-from-report prompt so Gemini writes code from actual findings.
        previous_interaction_id alone does NOT carry research context across
        interactions — Gemini treats each interaction independently.

        previous_id: kept for execution-failure fix requests where we pass the
        broken code and stderr — no research context needed there.

        On failure, sends the error back to Gemini and asks for a fix,
        then re-executes. Repeats up to MAX_FIX_ATTEMPTS times total.
        Uses background=True + _poll() for fix calls — same as Stage 1.

        Returns (output_files, output_file_errors, fix_cost_usd).
        fix_cost_usd is the accumulated cost of any fix API calls made.
        """
        import sys
        MAX_FIX_ATTEMPTS = 3
        fmt = task.output_formats[0]
        fix_cost_usd = 0.0

        # ── Extract code block ────────────────────────────────────────
        code = _extract_marked_code_block(report_text, FILE_OUTPUT_MARKER)
        if code is None:
            code = _extract_last_file_code_block(report_text, fmt)

        lib_hint = {
            "docx": "python-docx (from docx import Document)",
            "xlsx": "openpyxl or pandas with openpyxl engine",
            "pptx": "python-pptx (from pptx import Presentation)",
        }.get(fmt, fmt)

        if code is None:
            no_initial_code = True
            print(f"[Gemini] No code block in report — will request "
                  f"code generation from research context")
        else:
            no_initial_code = False
            code = _sanitize_code(code)
            print(f"[Gemini] Extracted code block ({len(code)} chars) — executing locally")

        # ── Staging dir = where input files live ──────────────────────
        if task.file_paths:
            staging_dir = os.path.dirname(os.path.abspath(task.file_paths[0]))
        else:
            staging_dir = task.output_files_dir

        os.makedirs(task.output_files_dir, exist_ok=True)

        script_path = os.path.join(
            staging_dir, f"_gemini_gen_{task.task_id}.py"
        )

        chdir_header = (
            "import os as _os\n"
            "_os.chdir(_os.path.dirname(_os.path.abspath(__file__)))\n\n"
        )

        def _write_script(c: str):
            with open(script_path, "w", encoding="utf-8") as sf:
                sf.write(chdir_header + c)

        def _move_output() -> tuple[list[str], dict]:
            """Look for output.{fmt} and move to output_files_dir."""
            expected_path = os.path.join(staging_dir, f"output.{fmt}")
            if os.path.exists(expected_path):
                dest = os.path.join(
                    task.output_files_dir,
                    f"gemini_{task.task_id}.{fmt}"
                )
                shutil.move(expected_path, dest)
                size = os.path.getsize(dest)
                print(f"[Gemini] ✓ File saved → {dest} ({size:,} bytes)")
                return [dest], {}
            found = _find_output_file(staging_dir, fmt, script_path)
            if found:
                dest = os.path.join(
                    task.output_files_dir,
                    f"gemini_{task.task_id}.{fmt}"
                )
                shutil.move(found, dest)
                size = os.path.getsize(dest)
                print(f"[Gemini] ✓ File saved → {dest} ({size:,} bytes)")
                return [dest], {}
            msg = (
                f"Code executed (rc=0) but output.{fmt} not found "
                f"in {staging_dir}."
            )
            print(f"[Gemini] ❌ {msg}")
            return [], {fmt: msg}

        output_files: list[str] = []
        output_file_errors: dict = {}

        try:
            for attempt in range(1, MAX_FIX_ATTEMPTS + 1):

                # ── Generation request (no code in report) ────────────
                # When the DRA produced no code block, skip execution on
                # attempt 1 and go straight to requesting code generation
                # using the research report as context.
                if no_initial_code and attempt == 1:
                    print(f"[Gemini] Requesting code generation from "
                          f"research report (attempt 1/{MAX_FIX_ATTEMPTS})...")
                    # Use pass1_report if available (multi-pass), else report_text
                    # (single-pass where report_text IS the research output).
                    # Never rely on previous_interaction_id for context — Gemini
                    # does not carry research context across interactions.
                    source_report = pass1_report if pass1_report else report_text
                    fix_prompt = (
                        f"You produced the following research report but did not "
                        f"include a Python code block to generate the required "
                        f"{fmt.upper()} file.\n\n"
                        f"Research report:\n{source_report}\n\n"
                        f"Using the findings in the report above, write a complete "
                        f"Python code block using {lib_hint} that generates a "
                        f"{fmt.upper()} file. "
                        f"Save the file as exactly: output.{fmt} "
                        f"in the current working directory.\n\n"
                        f"Wrap the code with EXACTLY these sentinel lines:\n\n"
                        f"```python\n"
                        f"# GEMINI_FILE_OUTPUT_START\n"
                        f"# ... your code ...\n"
                        f"# GEMINI_FILE_OUTPUT_END\n"
                        f"```\n\n"
                        f"Do NOT include [cite: N] markers inside the code. "
                        f"Return ONLY the code block, nothing else."
                    )
                    # fall through to the fix API call below
                    stderr_full = ""
                    # jump to fix call by setting a flag
                    goto_fix = True
                else:
                    goto_fix = False

                if not goto_fix:
                    _write_script(code)

                    # ── Execute ───────────────────────────────────────────
                    try:
                        proc = await asyncio.to_thread(
                            subprocess.run,
                            [sys.executable, script_path],
                            cwd=staging_dir,
                            capture_output=True,
                            text=True,
                            timeout=CODE_EXEC_TIMEOUT_SECONDS,
                        )
                    except subprocess.TimeoutExpired:
                        msg = f"Code execution timed out after {CODE_EXEC_TIMEOUT_SECONDS}s"
                        print(f"[Gemini] ❌ {msg}")
                        output_file_errors[fmt] = msg
                        break
                    except Exception as e:
                        msg = f"Code execution error: {e}"
                        print(f"[Gemini] ❌ {msg}")
                        output_file_errors[fmt] = msg
                        break

                    # ── Success ───────────────────────────────────────────
                    if proc.returncode == 0:
                        output_files, output_file_errors = _move_output()
                        break

                    # ── Failure ───────────────────────────────────────────
                    stderr_full = proc.stderr.strip()
                    print(f"[Gemini] ❌ Attempt {attempt}/{MAX_FIX_ATTEMPTS} "
                          f"failed (rc={proc.returncode}): {stderr_full[:200]}")

                    if attempt == MAX_FIX_ATTEMPTS:
                        msg = (
                            f"Code execution failed (rc={proc.returncode}) "
                            f"after {MAX_FIX_ATTEMPTS} attempts: {stderr_full[:600]}"
                        )
                        output_file_errors[fmt] = msg
                        break

                    # ── Ask Gemini to fix ─────────────────────────────────
                    print(f"[Gemini] Requesting code fix "
                          f"(attempt {attempt+1}/{MAX_FIX_ATTEMPTS})...")
                    fix_prompt = (
                        f"The Python code you provided to generate a .{fmt} file "
                        f"failed with this error:\n\n"
                        f"```\n{stderr_full}\n```\n\n"
                        f"Here is the code that failed:\n\n"
                        f"```python\n{code}\n```\n\n"
                        f"Fix the code so it runs without errors and saves "
                        f"output.{fmt} in the current working directory.\n\n"
                        f"Wrap your fixed code with EXACTLY these sentinel lines:\n\n"
                        f"```python\n"
                        f"# GEMINI_FILE_OUTPUT_START\n"
                        f"# ... fixed code ...\n"
                        f"# GEMINI_FILE_OUTPUT_END\n"
                        f"```\n\n"
                        f"Do NOT include [cite: N] markers. "
                        f"Return ONLY the code block, nothing else."
                    )
                try:
                    fix_create_kwargs = dict(
                        input=[{"type": "text", "text": fix_prompt}],
                        agent=self.model,
                        background=True,
                        store=True,
                    )
                    # For execution-failure fix requests, chain off previous_id
                    # so Gemini has the broken code conversation context.
                    # For generation-from-report (no_initial_code, attempt==1),
                    # do NOT chain — research context is pasted inline above.
                    if previous_id is not None and not (no_initial_code and attempt == 1):
                        fix_create_kwargs["previous_interaction_id"] = previous_id
                    fix_interaction = await asyncio.to_thread(
                        self._client.interactions.create,
                        **fix_create_kwargs,
                    )
                    fix_interaction = await self._poll(fix_interaction, task)

                    # Chain the next fix attempt off this one so Gemini has
                    # full context that its previous fix also failed.
                    previous_id = fix_interaction.id

                    # Extract cost of fix call
                    fix_usage = getattr(fix_interaction, "usage", None)
                    if fix_usage:
                        fix_in  = getattr(fix_usage, "total_input_tokens",  0) or 0
                        fix_out = getattr(fix_usage, "total_output_tokens", 0) or 0
                        fix_cached = getattr(fix_usage, "total_cached_tokens", 0) or 0
                        fix_cost_usd += estimate_cost(
                            self.model, fix_in, fix_out, fix_cached
                        )

                    fix_text = ""
                    for output in (getattr(fix_interaction, "outputs", None) or []):
                        fix_text = getattr(output, "text", "") or ""
                        if fix_text:
                            break

                    # Extract fixed code from response
                    fixed_code = _extract_marked_code_block(fix_text, FILE_OUTPUT_MARKER)
                    if fixed_code is None:
                        fixed_code = _extract_last_file_code_block(fix_text, fmt)
                    if fixed_code is None:
                        import re as _re
                        m = _re.search(
                            r'```python\s*\n(.*?)```', fix_text, _re.DOTALL
                        )
                        fixed_code = m.group(1).strip() if m else None

                    if not fixed_code:
                        msg = f"Gemini did not return fixed code on attempt {attempt+1}"
                        print(f"[Gemini] ❌ {msg}")
                        output_file_errors[fmt] = msg
                        break

                    code = _sanitize_code(fixed_code)
                    print(f"[Gemini] Rewrote script with fixed code "
                          f"({len(code)} chars)")

                except Exception as e:
                    msg = f"Gemini fix request failed: {e}"
                    print(f"[Gemini] ❌ {msg}")
                    output_file_errors[fmt] = msg
                    break

        finally:
            if os.path.exists(script_path):
                os.remove(script_path)

        return output_files, output_file_errors, fix_cost_usd

    # ─── Dry run ──────────────────────────────────────────────────────

    def _dry_run_result(self, task: ResearchTask, started_at) -> AgentResult:
        completed_at = datetime.now(timezone.utc)
        iat_status   = ("CLOSED (web search disabled via instruction)"
                        if task.is_closed else "OPEN (web search enabled)")
        file_summary = (", ".join(os.path.basename(f) for f in task.file_paths)
                        or "none")

        file_note = ""
        file_errs = {}
        if task.output_formats:
            fmt = task.output_formats[0]
            file_note = (
                f"\n## File Generation (Dry Run)\n\n"
                f"Requested: {task.output_formats}\n"
                f"Approach: DRA embeds Python code block in report "
                f"(marker: {FILE_OUTPUT_MARKER})\n"
                f"Execution: local subprocess in staging dir\n"
                f"Input files available: yes (on disk)\n"
            )
            file_errs = {fmt: "DRY_RUN"}

        report = (
            f"# [DRY RUN] Gemini Deep Research Report\n\n"
            f"**Task:** {task.task_id}\n"
            f"**Model:** {self.model}\n"
            f"**IAT:** {iat_status}\n"
            f"**Context window:** 1M tokens\n"
            f"**Files:** {file_summary}\n"
            f"{file_note}\n"
            f"## Prompt (first 500 chars)\n\n"
            f"{task.prompt[:500]}{'...' if len(task.prompt) > 500 else ''}\n\n"
            f"## Limitations\n\nDry-run — no actual research performed.\n"
        )

        print(f"[Gemini] Dry run complete for {task.task_id}")
        return AgentResult(
            task_id=task.task_id,
            agent="gemini",
            model=self.model,
            response_text=report,
            citations=[{"url": "https://example.com/mock",
                        "title": "Mock Citation", "snippet": "Dry run"}],
            tool_call_log=[],
            input_tokens=0,
            output_tokens=0,
            total_cost_usd=0.0,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=1,
            completed=True,
            forced_stop=False,
            output_files=[],
            output_file_errors=file_errs,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    def _error_result(self, task: ResearchTask, started_at,
                      error_msg: str) -> AgentResult:
        completed_at = datetime.now(timezone.utc)
        return AgentResult(
            task_id=task.task_id,
            agent="gemini",
            model=self.model,
            response_text="",
            tool_call_log=[],
            input_tokens=0,
            output_tokens=0,
            total_cost_usd=0.0,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=0,
            completed=False,
            forced_stop=True,
            error=error_msg,
            output_files=[],
            output_file_errors={},
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )


# ─── Code block extraction helpers ───────────────────────────────────────────

def _sanitize_code(code: str) -> str:
    """
    Remove Gemini inline citation markers from generated Python code.

    Gemini DRA embeds [cite: N] / [cite: N, M] inside code blocks,
    causing SyntaxError. Three patterns handled in priority order:

    P1  key [cite: N]: value  ->  key: value   (citation between key and colon)
    P2  key: [cite: N]        ->  key: []      (citation as sole value)
    P3  value [cite: N]       ->  value        (inline suffix / comment)
    """
    import re
    # P1: citation annotating the key, sits before the colon
    code = re.sub(r"(['\"][^'\"]*['\"]) *\[cite:\s*[\d,\s]+\](\s*:)",
                  r"\1\2", code)
    # P2: citation is the sole value after a colon — replace with empty list
    code = re.sub(r"(:\s*)\[cite:\s*[\d,\s]+\]", r"\1[]", code)
    # P3: any remaining citation marker — strip entirely (no leading space required)
    code = re.sub(r"\[cite:\s*[\d,\s]+\]", "", code)
    # P4: fix systematic Gemini python-docx API error
    # Gemini always writes table.rows.cells (wrong) — rows is a _Rows container,
    # must index first: table.rows[0].cells, table.rows[1].cells etc.
    code = code.replace(".rows.cells", ".rows[0].cells")
    # P5: fix tuple.text — Gemini treats .cells tuple as single Cell object
    # Pattern: <var>_cells.text = X  →  <var>_cells[0].text = X
    # Also catches row1.text, row_cells.text etc.
    import re as _re
    code = _re.sub(r'(\b\w+_cells)\.text(\s*=)', r'\1[0].text\2', code)
    code = _re.sub(r'(\b\w+_cells)\.text\b', r'\1[0].text', code)
    return code


def _strip_file_output_block(text):
    # type: (str) -> str
    """
    Remove the file-generation code block from a research report.

    When Gemini includes code inline in Pass 1, pass1_report contains
    both the research and the sentinel-wrapped code block. We want only
    the research in response_text — the code is extracted separately.

    Strips the opening fence line (if any), FILE_OUTPUT_START through
    FILE_OUTPUT_END, and the closing fence line (if any).
    Returns text unchanged if sentinels are not found.
    """
    start_idx = text.find(FILE_OUTPUT_START)
    if start_idx == -1:
        return text.strip()
    end_idx = text.find(FILE_OUTPUT_END)
    if end_idx == -1 or end_idx < start_idx:
        return text.strip()

    before = text[:start_idx]
    after  = text[end_idx + len(FILE_OUTPUT_END):]

    # Strip the opening fence line from the tail of before.
    # The fence immediately precedes FILE_OUTPUT_START on its own line.
    lines_before = before.rstrip().split("\n")
    while lines_before and re.match(r"^`{3,}", lines_before[-1].strip()):
        lines_before.pop()
    before = "\n".join(lines_before).rstrip()

    # Strip the closing fence line from the head of after.
    lines_after = after.lstrip("\n").split("\n")
    if lines_after and re.match(r"^`{3,}\s*$", lines_after[0]):
        lines_after.pop(0)
    after = "\n".join(lines_after).lstrip("\n")

    result = before
    if after:
        result = result + "\n\n" + after
    return result.strip()

# Robust fence regex — matches any of:
#   ```python  ```Python  ```py  ```  ````  `````  (3+ backticks, any case)
# Handles the full range of Gemini code fence variations without requiring
# a specific language tag.
_FENCE_RE = re.compile(
    r"`{3,}"             # 3 or more backticks (opening fence)
    r"[ \t]*"            # optional spaces after backticks
    r"(?:python|py)?"    # optional language tag, case-insensitive
    r"[ \t]*\n"          # rest of opening fence line
    r"(.*?)"             # code content — non-greedy
    r"\n?`{3,}",         # closing fence (3+ backticks)
    re.DOTALL | re.IGNORECASE,
)


def _extract_marked_code_block(report_text: str, marker: str) -> Optional[str]:
    """
    Extract the file-generation code block from a DRA response.

    Extraction priority:
    1. Sentinel block — content between FILE_OUTPUT_START and FILE_OUTPUT_END.
       These are explicit markers we ask Gemini to include, making extraction
       immune to fence language-tag variations and surrounding prose.
    2. Fence block containing the marker — scanned with _FENCE_RE which handles
       ```python, ```Python, ```py, plain ```, and 3+ backtick variants.
    """
    # Priority 1: sentinel markers
    start_idx = report_text.find(FILE_OUTPUT_START)
    end_idx   = report_text.find(FILE_OUTPUT_END)
    if start_idx != -1 and end_idx != -1 and end_idx > start_idx:
        raw = report_text[start_idx : end_idx + len(FILE_OUTPUT_END)]
        # Strip any surrounding fence lines (``` wrapping the sentinels)
        lines = [l for l in raw.splitlines()
                 if not re.match(r"^`{3,}", l.strip())]
        return "\n".join(lines).strip()

    # Priority 2: any fenced block containing the marker
    for match in _FENCE_RE.finditer(report_text):
        code = match.group(1)
        if marker in code:
            return code.strip()

    return None


def _extract_last_file_code_block(report_text: str, fmt: str) -> Optional[str]:
    """
    Fallback: find the last fenced code block that references a file-generation
    library appropriate for the requested format.

    Uses _FENCE_RE so it handles ```python, ```Python, ```py, plain ```,
    and 3+ backtick variants — not just ```python.
    """
    library_hints = {
        "xlsx": ["openpyxl", "pandas", "xlsxwriter"],
        "docx": ["docx", "Document"],
        "pptx": ["pptx", "Presentation"],
        "pdf":  ["reportlab", "FPDF"],
    }
    hints = library_hints.get(fmt, [])

    last_match = None
    for match in _FENCE_RE.finditer(report_text):
        code = match.group(1)
        if any(hint in code for hint in hints):
            last_match = code
    return last_match.strip() if last_match else None


def _find_output_file(staging_dir: str, fmt: str,
                      exclude: str) -> Optional[str]:
    """
    Scan staging_dir for any file with the right extension that isn't
    the script itself. Returns the path of the most recently modified one.
    """
    candidates = []
    for fname in os.listdir(staging_dir):
        fpath = os.path.join(staging_dir, fname)
        if (fname.endswith(f".{fmt}")
                and fpath != exclude
                and os.path.isfile(fpath)):
            candidates.append(fpath)
    if not candidates:
        return None
    return max(candidates, key=os.path.getmtime)