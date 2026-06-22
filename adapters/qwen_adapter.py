"""
qwen_adapter.py v2 — Full Amazon GDPVal-aligned Qwen Deep Research Agent.

Architecture
────────────
  Transparent agentic loop: Qwen generates tool_calls → we execute
  locally → feed results back → Qwen continues. Full observability.

8-tool suite (locally executed, matching Amazon GDPVal sandbox exactly)
──────────────────────────────────────────────────────────────────────
  python_execute(code)             — Run Python in sandbox. PRIMARY workhorse.
  bash_execute(command)            — Run bash commands. System ops, piping.
  web_search(query, max_results)   — DuckDuckGo (free default) or Tavily.
  web_fetch(url)                   — Fetch URL content as text.
  read_file(filename)              — Read input file (PDF/XLSX/DOCX/PPTX/CSV/ODS/XLSB).
  write_file(filename, content)    — Write text content to file.
  list_directory(path)             — List files and sizes.
  convert_document(source, format) — LibreOffice headless format conversion.

Pre-installed libraries (matching Amazon GDPVal sandbox):
  Spreadsheets:   openpyxl, xlrd, odfpy
  Documents:      python-docx, python-pptx, reportlab
  PDFs:           pdfplumber
  Images:         Pillow (PIL)
  Data & analysis: pandas, numpy, scipy
  Visualization:  matplotlib, seaborn
  Office conv.:   LibreOffice headless

File handling
─────────────
  Input files:  Extracted to text and injected inline in the system message
                (like Gemini's approach). Originals remain on disk at known
                paths for bash_execute code to reference directly.

  Output files: Model writes code via bash_execute that saves files.
                Adapter detects and captures output files from the staging dir.

Deployment
──────────
  # vLLM (recommended for Qwen 3.6)
  vllm serve Qwen/Qwen3.6-27B \\
      --port 8000 \\
      --reasoning-parser qwen3 \\
      --enable-auto-tool-choice \\
      --tool-call-parser qwen3_coder

  # Then in .env:
  QWEN_BASE_URL=http://localhost:8000/v1
  QWEN_API_KEY=EMPTY
  QWEN_MODEL=Qwen/Qwen3.6-27B

  # Or for DashScope:
  QWEN_BASE_URL=https://dashscope.aliyuncs.com/compatible-mode/v1
  QWEN_API_KEY=sk-...
  QWEN_MODEL=qwen-max-latest

Comparison with other adapters
──────────────────────────────
  Aspect               Claude            o3               Gemini            Qwen (this)
  ─────────────────── ────────────────  ────────────────  ────────────────  ────────────────
  Loop visibility      Transparent       Black box         Black box         Transparent
  Tool execution       Server-side       Server-side       Local subprocess  Local subprocess
  Web search           Server-side       Server-side       Server-side       Local (DDG/Tavily)
  Code execution       Anthropic sandbox OpenAI container  Local subprocess  Local subprocess
  File input method    Files API upload  Files API upload  Inline text       Inline text
  File output capture  Files API DL      Containers API    Sentinel code     Staging dir scan
  IAT-1 enforcement    Hard (tool omit)  Soft              Soft              Hard (tool omit)
  Observability        FULL              LOW               LOW               FULL
  Cost                 API pricing       API pricing       API pricing       Free (local) / API
"""

from __future__ import annotations

import os
import re
import sys
import json
import time
import asyncio
import logging
import subprocess
import tempfile
import shutil
import tempfile
import shutil
from datetime import datetime, timezone
from typing import Optional, Any

from models import ResearchTask, AgentResult, ToolCall

_log = logging.getLogger("dra.qwen")


# ── Pricing ───────────────────────────────────────────────────────────────────
# Local vLLM: free. DashScope: varies by model.
# Only used for budget guards; set to 0 for local serving.

PRICING = {
    # DashScope pricing (USD per 1M tokens, approximate)
    "qwen-max-latest":     {"input_per_mtok": 2.50, "output_per_mtok": 7.50},
    "qwen-plus-latest":    {"input_per_mtok": 0.80, "output_per_mtok": 2.00},
    "qwen-turbo-latest":   {"input_per_mtok": 0.30, "output_per_mtok": 0.60},
    "qwen3.7-max":         {"input_per_mtok": 2.50, "output_per_mtok": 7.50},
    # Local models: free (set to near-zero so budget guard doesn't fire)
    "_local_default":      {"input_per_mtok": 0.001, "output_per_mtok": 0.001},
}


def estimate_cost(model: str, input_tokens: int, output_tokens: int) -> float:
    prices = None
    for key in PRICING:
        if key in model.lower():
            prices = PRICING[key]
            break
    if prices is None:
        prices = PRICING["_local_default"]
    return round(
        (input_tokens / 1_000_000) * prices["input_per_mtok"] +
        (output_tokens / 1_000_000) * prices["output_per_mtok"],
        6,
    )


def _is_local_serving(base_url: str) -> bool:
    """Check if the base URL points to a local server (free inference)."""
    return any(h in base_url for h in ("localhost", "127.0.0.1", "0.0.0.0"))


# ── Code execution constants ─────────────────────────────────────────────────

PYTHON_EXEC_TIMEOUT  = 180   # seconds per python_execute call
BASH_EXEC_TIMEOUT    = 120   # seconds per bash_execute call
CONVERT_TIMEOUT      = 120   # seconds per LibreOffice conversion
MAX_OUTPUT_CHARS     = 50_000
MAX_FILE_READ_CHARS  = 100_000


# ── System prompt ─────────────────────────────────────────────────────────────

RESEARCH_SYSTEM_PROMPT = """\
You are a deep research agent in a Linux sandbox. Produce a comprehensive, \
well-cited research report and any required deliverable artifacts.

## Execution environment

**Python libraries (import directly — no pip install needed):**
- Spreadsheets: `openpyxl`, `xlrd`, `odfpy`
- Documents: `python-docx` (import as `docx`), `python-pptx` (import as `pptx`), `reportlab`
- PDFs: `pdfplumber`
- Images: `Pillow` (import as `PIL`)
- Data & analysis: `pandas`, `numpy`, `scipy`
- Visualization: `matplotlib`, `seaborn`
- Audio: `wave` (stdlib)
- General: `json`, `csv`, `re`, `math`, `statistics`, `collections`, `itertools`

**System tools:** LibreOffice headless (convert_document), bash, grep, awk, sed, sort.

## Tools (11 available)

1. **python_execute** — Run Python code. PRIMARY tool for analysis, calculations, \
visualizations, and generating output files (.xlsx, .docx, .pdf, .pptx, .png, .csv). \
All libraries above are available. Always print() results to verify.
2. **bash_execute** — Run bash commands. File management, piping, text processing.
3. **web_search** — Search the web. Returns titles, URLs, snippets.
4. **web_fetch** — Fetch full text of a URL.
5. **read_file** — Read input file text content. Supports page ranges for PDFs: \
read_file(filename, pages="44-57") reads only those pages. Use page ranges for large docs.
6. **write_file** — Write text content to a file (.txt, .md, .csv, .json, .html).
7. **list_directory** — List files and sizes. Verify output files were created.
8. **convert_document** — Convert formats via LibreOffice (docx→pdf, xlsx→pdf, etc).
9. **read_file_visual** — Read VISUAL content: charts, diagrams, infographics, slide layouts. \
Rasterizes pages to images and uses a vision model to describe every visual element. \
Supports page ranges: read_file_visual(filename, pages="5"). Files marked [VISUAL] in manifest.
10. **search_in_file** — Search for a keyword/phrase within a file. Returns matching page \
numbers and text snippets. Like Ctrl+F. Use to locate specific data before reading full pages.
11. **finish** — (Internal) Signals task completion. Do NOT call this directly.

## MANDATORY reading strategy

You MUST read ALL input files before writing any analysis. A response that skips \
input files will be REJECTED and you will be asked to redo the work.

**For large documents (>20 pages, e.g. annual reports, policy manuals):**
1. First call: read_file(filename, pages="1-5") to get the structure/table of contents.
2. Use search_in_file(filename, query) to locate specific data points.
3. Read ONLY relevant sections using page ranges: read_file(filename, pages="44-57").
4. Use read_file_visual ONLY on specific pages with charts/diagrams you need.
5. NEVER dump an entire large document into context. Be surgical.

**For small documents (<20 pages):**
- Read the full file: read_file(filename).
- Use read_file_visual if the file is marked [VISUAL] in the manifest.

## Methodology

1. READ FILES FIRST: Read every input file (mandatory, enforced). For large files, \
read the TOC/structure first, then target relevant sections.
2. PLAN: Outline 3-5 key questions that need answering.
3. SEARCH: Use web_search for external data (if permitted).
4. ANALYZE: Use python_execute for ALL calculations. Print and verify results.
5. SYNTHESIZE: Cross-reference findings across sources. Note contradictions.
6. GENERATE: If file deliverable required, use python_execute with the right library. \
Verify with list_directory.

## Rules

- Never fabricate data or sources. Say so if you cannot find evidence.
- Always verify calculations — run code, print output, check it.
- Read ALL input files before analysis. Hidden footnotes and corrections are common.
- Cite: [1], [2] for web, [File: name, page X] for input files.
- Output: Markdown headers, executive summary at top, Limitations section at end.
"""

FILE_GEN_ADDENDUM = """\

## FILE GENERATION REQUIREMENT — MANDATORY

You MUST generate the required output file(s) using python_execute.

Library reference for output formats:
- **.xlsx** → `openpyxl` or `pandas` with `engine='openpyxl'`
- **.docx** → `from docx import Document`
- **.pptx** → `from pptx import Presentation`
- **.pdf**  → `from reportlab.lib.pagesizes import letter; from reportlab.pdfgen import canvas` \
OR generate .docx first then convert_document to PDF
- **.csv**  → `csv` module or `pandas.to_csv()`
- **.png/.jpg** → `matplotlib.pyplot.savefig()` or `from PIL import Image`
- **.wav**  → `import wave`

After generating, ALWAYS verify with list_directory. A response missing the file is FAILED.
"""

CLOSED_CORPUS_ADDENDUM = """\

## IMPORTANT: Closed corpus constraint (IAT-1)

This task is STRICTLY CLOSED CORPUS. You must ONLY use the provided input \
files to answer the question. Do NOT use web_search or web_fetch. Do NOT \
use any external knowledge. All conclusions must be derived solely from \
the provided documents.
"""


def _build_system_prompt(task: ResearchTask) -> str:
    prompt = RESEARCH_SYSTEM_PROMPT
    if task.is_closed:
        prompt += CLOSED_CORPUS_ADDENDUM

    # Always include file generation guidance — every task may need deliverables
    prompt += FILE_GEN_ADDENDUM

    # Detector hint as advisor (not instructor)
    if task.output_formats:
        prompt += (
            f"\nDetector hint (may have errors): the prompt likely expects "
            f"**{', '.join(task.output_formats)}** output. Verify this against "
            f"the actual prompt requirements.\n"
        )
    else:
        prompt += (
            "\nDetector hint: no specific output format was auto-detected, but "
            "read the prompt carefully — it may still require a file deliverable "
            "(memo, report, presentation, spreadsheet, etc). When in doubt, "
            "produce a .docx deliverable.\n"
        )
    return prompt


# ── Tool definitions (OpenAI function-calling format) ─────────────────────────

def _build_tool_definitions(task: ResearchTask) -> list[dict]:
    """Build the full 8-tool suite matching Amazon GDPVal sandbox."""
    tools = []

    # 1. python_execute — PRIMARY tool for analysis and file generation
    tools.append({"type": "function", "function": {
        "name": "python_execute",
        "description": (
            "Execute Python code in the sandbox. PRIMARY tool for data analysis, "
            "calculations, visualization, and generating output files. "
            "Pre-installed: openpyxl, xlrd, odfpy, pandas, numpy, scipy, "
            "matplotlib, seaborn, python-docx (import docx), python-pptx "
            "(import pptx), reportlab, pdfplumber, Pillow (import PIL), wave. "
            "Working directory contains all input files. Always print() results."
        ),
        "parameters": {"type": "object", "properties": {
            "code": {"type": "string", "description": "Complete self-contained Python code. Include all imports."},
        }, "required": ["code"]},
    }})

    # 2. bash_execute — system operations
    tools.append({"type": "function", "function": {
        "name": "bash_execute",
        "description": (
            "Execute a bash command. Working directory has all input files. "
            "Use for: file management (cp, mv, ls, cat, head, tail, wc), "
            "text processing (grep, awk, sed, sort), piping. "
            "Prefer python_execute for data analysis and file generation."
        ),
        "parameters": {"type": "object", "properties": {
            "command": {"type": "string", "description": "Bash command to execute."},
        }, "required": ["command"]},
    }})

    # 3. web_search (IAT-1: omitted)
    if task.web_search_enabled:
        tools.append({"type": "function", "function": {
            "name": "web_search",
            "description": (
                "Search the web. Returns titles, URLs, snippets. "
                "Prefer primary sources (official reports, government data, "
                "peer-reviewed papers) over aggregators."
            ),
            "parameters": {"type": "object", "properties": {
                "query": {"type": "string", "description": "Search query."},
                "max_results": {"type": "integer", "description": "Max results (default: 8).", "default": 8},
            }, "required": ["query"]},
        }})

    # 4. web_fetch (IAT-1: omitted)
    if task.web_search_enabled:
        tools.append({"type": "function", "function": {
            "name": "web_fetch",
            "description": "Fetch full text content of a URL. HTML auto-converted to text.",
            "parameters": {"type": "object", "properties": {
                "url": {"type": "string", "description": "URL to fetch."},
            }, "required": ["url"]},
        }})

    # 5. read_file (with page ranges for PDFs)
    tools.append({"type": "function", "function": {
        "name": "read_file",
        "description": (
            "Read input file content. Handles: PDF, XLSX, XLS, ODS, XLSB, "
            "DOCX, PPTX, CSV, TSV, and all text formats. "
            "For PDFs: use pages param to read specific pages (e.g. '1-5', '44-57', '154'). "
            "ALWAYS read input files before starting analysis."
        ),
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string", "description": "Filename from the manifest."},
            "pages": {"type": "string", "description": "Page range for PDFs: '1-5', '44-57', '3'. Omit to read all pages.", "default": ""},
        }, "required": ["filename"]},
    }})

    # 6. write_file
    tools.append({"type": "function", "function": {
        "name": "write_file",
        "description": (
            "Write text content to a file. Use for: .txt, .md, .csv, .json, "
            ".xml, .html, Python scripts, intermediate data. "
            "Do NOT use for binary formats (.xlsx, .docx, .pdf, .png) — "
            "use python_execute with the appropriate library instead."
        ),
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string", "description": "File to create/overwrite."},
            "content": {"type": "string", "description": "Text content."},
        }, "required": ["filename", "content"]},
    }})

    # 7. list_directory
    tools.append({"type": "function", "function": {
        "name": "list_directory",
        "description": (
            "List files and sizes. Use to verify input files exist "
            "and confirm output files were generated."
        ),
        "parameters": {"type": "object", "properties": {
            "path": {"type": "string", "description": "Directory (default: '.').", "default": "."},
        }, "required": []},
    }})

    # 8. convert_document
    tools.append({"type": "function", "function": {
        "name": "convert_document",
        "description": (
            "Convert document between formats via LibreOffice headless. "
            "Supports: docx→pdf, xlsx→pdf, pptx→pdf, odt→pdf, odt→docx, "
            "csv→xlsx, html→pdf, txt→pdf."
        ),
        "parameters": {"type": "object", "properties": {
            "source_file": {"type": "string", "description": "Source file path."},
            "target_format": {"type": "string", "description": "Target extension without dot: pdf, docx, xlsx, etc."},
        }, "required": ["source_file", "target_format"]},
    }})

    # 9. read_file_visual — for charts, diagrams, images in documents
    tools.append({"type": "function", "function": {
        "name": "read_file_visual",
        "description": (
            "Read visual content from a file: charts, diagrams, images, "
            "tables rendered as graphics, slide layouts. "
            "Rasterizes pages and uses a vision model to describe every visual element. "
            "Supports page ranges for PDFs: pages='5' or pages='4-6'. "
            "Use INSTEAD of read_file for files marked [VISUAL] in manifest."
        ),
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string", "description": "Filename from the manifest."},
            "pages": {"type": "string", "description": "Page range to rasterize: '5', '4-6'. Omit for all pages (max 20).", "default": ""},
            "context_hint": {"type": "string", "description": "What to focus on (e.g. 'revenue chart', 'org diagram').", "default": ""},
        }, "required": ["filename"]},
    }})

    # 10. search_in_file — Ctrl+F within a document
    tools.append({"type": "function", "function": {
        "name": "search_in_file",
        "description": (
            "Search for a keyword or phrase within a file. Returns matching page numbers "
            "and surrounding text snippets. Like Ctrl+F in a PDF viewer. "
            "Use this to locate specific data in large documents BEFORE reading full pages. "
            "Works on PDF, DOCX, PPTX, and text files."
        ),
        "parameters": {"type": "object", "properties": {
            "filename": {"type": "string", "description": "Filename to search in."},
            "query": {"type": "string", "description": "Keyword or phrase to search for."},
            "max_results": {"type": "integer", "description": "Max matches to return (default: 10).", "default": 10},
        }, "required": ["filename", "query"]},
    }})

    return tools


# ── Tool execution functions ──────────────────────────────────────────────────

async def _exec_web_search(query: str, max_results: int = 8) -> str:
    """Execute a web search using DuckDuckGo (free, no API key needed)."""
    try:
        try:
            from ddgs import DDGS
        except ImportError:
            from duckduckgo_search import DDGS
        results = []
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                results.append({
                    "title": r.get("title", ""),
                    "url": r.get("href", ""),
                    "snippet": r.get("body", "")[:300],
                })
        if results:
            return json.dumps(results, indent=2, ensure_ascii=False)
        return json.dumps({"message": "No results found", "query": query})
    except ImportError:
        # Fallback: try Tavily
        tavily_key = os.environ.get("TAVILY_API_KEY")
        if tavily_key:
            return await _exec_web_search_tavily(query, max_results, tavily_key)
        return json.dumps({
            "error": "Web search unavailable. Install: pip install duckduckgo-search "
                     "OR set TAVILY_API_KEY for Tavily.",
        })
    except Exception as e:
        return json.dumps({"error": f"Search failed: {str(e)}", "query": query})


async def _exec_web_search_tavily(
    query: str, max_results: int, api_key: str
) -> str:
    """Fallback web search using Tavily API."""
    import httpx
    try:
        async with httpx.AsyncClient(timeout=30) as client:
            resp = await client.post(
                "https://api.tavily.com/search",
                json={
                    "api_key": api_key,
                    "query": query,
                    "max_results": max_results,
                    "include_answer": False,
                },
            )
            resp.raise_for_status()
            data = resp.json()
            results = []
            for r in data.get("results", []):
                results.append({
                    "title": r.get("title", ""),
                    "url": r.get("url", ""),
                    "snippet": r.get("content", "")[:300],
                })
            return json.dumps(results, indent=2, ensure_ascii=False)
    except Exception as e:
        return json.dumps({"error": f"Tavily search failed: {e}", "query": query})


async def _exec_web_fetch(url: str) -> str:
    """Fetch URL content and return as text."""
    import httpx
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/125.0.0.0 Safari/537.36 DRA-Qwen/1.0"
            ),
        }
        async with httpx.AsyncClient(
            timeout=30, follow_redirects=True, max_redirects=5
        ) as client:
            resp = await client.get(url, headers=headers)
            resp.raise_for_status()
            content_type = resp.headers.get("content-type", "")

            if "text/html" in content_type:
                text = _html_to_text(resp.text)
            else:
                text = resp.text

            if len(text) > MAX_OUTPUT_CHARS:
                text = text[:MAX_OUTPUT_CHARS] + (
                    f"\n\n[... TRUNCATED: showing first {MAX_OUTPUT_CHARS} of "
                    f"{len(resp.text)} characters ...]"
                )
            return text

    except Exception as e:
        return f"[Fetch error: {e}]"


def _html_to_text(html: str) -> str:
    """Simple HTML to text conversion. Strips tags, keeps structure."""
    # Remove script and style blocks
    html = re.sub(r'<script[^>]*>.*?</script>', '', html, flags=re.DOTALL | re.IGNORECASE)
    html = re.sub(r'<style[^>]*>.*?</style>', '', html, flags=re.DOTALL | re.IGNORECASE)
    # Convert common block elements to newlines
    html = re.sub(r'<(?:br|p|div|h[1-6]|li|tr)[^>]*>', '\n', html, flags=re.IGNORECASE)
    # Strip remaining tags
    text = re.sub(r'<[^>]+>', '', html)
    # Decode common entities
    for entity, char in [('&amp;', '&'), ('&lt;', '<'), ('&gt;', '>'),
                         ('&quot;', '"'), ('&#39;', "'"), ('&nbsp;', ' ')]:
        text = text.replace(entity, char)
    # Clean up whitespace
    lines = [line.strip() for line in text.splitlines()]
    text = '\n'.join(line for line in lines if line)
    return text


async def _exec_python(code: str, working_dir: str) -> str:
    """Execute Python code via temp file. Clean sandbox execution."""
    script_path = os.path.join(working_dir, "_qwen_exec.py")
    try:
        with open(script_path, "w", encoding="utf-8") as f:
            f.write(code)
        proc = await asyncio.create_subprocess_exec(
            sys.executable, script_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=working_dir,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1", "MPLBACKEND": "Agg"},
        )
        try:
            stdout, stderr = await asyncio.wait_for(
                proc.communicate(), timeout=PYTHON_EXEC_TIMEOUT
            )
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return f"[TIMEOUT: Python execution exceeded {PYTHON_EXEC_TIMEOUT}s and was killed]"
        parts = []
        if stdout:
            out = stdout.decode("utf-8", errors="replace")
            if len(out) > MAX_OUTPUT_CHARS:
                out = out[:MAX_OUTPUT_CHARS] + "\n[... stdout truncated ...]"
            parts.append(out)
        if stderr:
            err = stderr.decode("utf-8", errors="replace")
            if len(err) > MAX_OUTPUT_CHARS:
                err = err[:MAX_OUTPUT_CHARS] + "\n[... stderr truncated ...]"
            parts.append(f"STDERR:\n{err}")
        parts.append(f"EXIT CODE: {proc.returncode}")
        return "\n".join(parts) if parts else "[No output]"
    except Exception as e:
        return f"[Python execution error: {e}]"
    finally:
        if os.path.exists(script_path):
            try: os.remove(script_path)
            except OSError: pass


async def _exec_bash(command: str, working_dir: str) -> str:
    """
    Execute a bash command in the staging directory.

    The staging dir contains all input files on disk, so the model's code
    can reference them by filename directly (same as Amazon GDPVal sandbox).
    """
    try:
        proc = await asyncio.create_subprocess_shell(
            command,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            cwd=working_dir,
            env={**os.environ, "PYTHONDONTWRITEBYTECODE": "1"},
        )
        try:
            stdout, stderr = await asyncio.wait_for(
                proc.communicate(), timeout=BASH_EXEC_TIMEOUT
            )
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return f"[TIMEOUT: Command exceeded {BASH_EXEC_TIMEOUT}s limit and was killed]"

        result_parts = []
        if stdout:
            out = stdout.decode("utf-8", errors="replace")
            if len(out) > MAX_OUTPUT_CHARS:
                out = out[:MAX_OUTPUT_CHARS] + "\n[... stdout truncated ...]"
            result_parts.append(f"STDOUT:\n{out}")
        if stderr:
            err = stderr.decode("utf-8", errors="replace")
            if len(err) > MAX_OUTPUT_CHARS:
                err = err[:MAX_OUTPUT_CHARS] + "\n[... stderr truncated ...]"
            result_parts.append(f"STDERR:\n{err}")

        exit_note = f"EXIT CODE: {proc.returncode}"
        result_parts.append(exit_note)

        return "\n".join(result_parts) if result_parts else "[No output]"

    except Exception as e:
        return f"[Execution error: {e}]"


def _parse_page_range(pages_str: str) -> set[int]:
    """Parse page range string like '1-5', '3', '44-57,154-200' into a set of page numbers."""
    result = set()
    for part in pages_str.replace(" ", "").split(","):
        if "-" in part:
            try:
                start, end = part.split("-", 1)
                result.update(range(int(start), int(end) + 1))
            except ValueError:
                pass
        else:
            try:
                result.add(int(part))
            except ValueError:
                pass
    return result


def _exec_search_in_file(filename: str, query: str, file_paths: list[str], max_results: int = 10) -> str:
    """Search for a keyword/phrase within a file. Returns page numbers and snippets."""
    # Find the file
    target_path = None
    for fpath in file_paths:
        if os.path.basename(fpath).lower() == filename.lower():
            target_path = fpath
            break
    if target_path is None and os.path.exists(filename):
        target_path = filename
    if target_path is None:
        return json.dumps({"error": f"File not found: {filename}"})

    ext = os.path.splitext(target_path)[1].lower()
    query_lower = query.lower()
    matches = []

    try:
        if ext == ".pdf":
            try:
                import pdfplumber
                with pdfplumber.open(target_path) as pdf:
                    for i, page in enumerate(pdf.pages):
                        text = page.extract_text() or ""
                        if query_lower in text.lower():
                            # Extract snippet around the match
                            idx = text.lower().find(query_lower)
                            start = max(0, idx - 100)
                            end = min(len(text), idx + len(query) + 100)
                            snippet = text[start:end].strip()
                            matches.append({"page": i + 1, "snippet": f"...{snippet}..."})
                            if len(matches) >= max_results:
                                break
            except ImportError:
                try:
                    import fitz
                    doc = fitz.open(target_path)
                    for i, page in enumerate(doc):
                        text = page.get_text()
                        if query_lower in text.lower():
                            idx = text.lower().find(query_lower)
                            start = max(0, idx - 100)
                            end = min(len(text), idx + len(query) + 100)
                            matches.append({"page": i + 1, "snippet": f"...{text[start:end].strip()}..."})
                            if len(matches) >= max_results:
                                break
                    doc.close()
                except ImportError:
                    return "[Search unavailable: install pdfplumber or pymupdf]"

        elif ext == ".docx":
            try:
                from docx import Document
                doc = Document(target_path)
                for i, para in enumerate(doc.paragraphs):
                    if query_lower in para.text.lower():
                        matches.append({"paragraph": i + 1, "snippet": para.text[:200]})
                        if len(matches) >= max_results:
                            break
            except ImportError:
                return "[Search unavailable: install python-docx]"

        elif ext in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(target_path, read_only=True, data_only=True)
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    for row_idx, row in enumerate(ws.iter_rows(max_row=10000, values_only=True), 1):
                        row_text = "\t".join(str(c) if c is not None else "" for c in row)
                        if query_lower in row_text.lower():
                            matches.append({"sheet": sheet_name, "row": row_idx, "snippet": row_text[:200]})
                            if len(matches) >= max_results:
                                break
                    if len(matches) >= max_results:
                        break
                wb.close()
            except ImportError:
                return "[Search unavailable: install openpyxl]"

        else:
            # Text-based files
            try:
                with open(target_path, "r", encoding="utf-8", errors="replace") as f:
                    for line_num, line in enumerate(f, 1):
                        if query_lower in line.lower():
                            matches.append({"line": line_num, "snippet": line.strip()[:200]})
                            if len(matches) >= max_results:
                                break
            except Exception as e:
                return f"[Search error: {e}]"

        if matches:
            return json.dumps({"query": query, "matches": len(matches), "results": matches}, indent=2, ensure_ascii=False)
        return json.dumps({"query": query, "matches": 0, "message": f"'{query}' not found in {filename}"})

    except Exception as e:
        return f"[Search error in {filename}: {e}]"


def _exec_read_file(filename: str, file_paths: list[str], pages: str = "") -> str:
    """
    Read an input file by name.

    Searches the task's file_paths for a matching filename.
    Uses text extraction for binary formats (PDF, DOCX, XLSX, PPTX).
    """
    # Find the file
    target_path = None
    for fpath in file_paths:
        if os.path.basename(fpath).lower() == filename.lower():
            target_path = fpath
            break
        if os.path.basename(fpath) == filename:
            target_path = fpath
            break
    if target_path is None:
        # Also try exact path match
        if os.path.exists(filename):
            target_path = filename
    if target_path is None:
        available = [os.path.basename(f) for f in file_paths]
        return json.dumps({
            "error": f"File not found: {filename}",
            "available_files": available,
        })

    ext = os.path.splitext(target_path)[1].lower()

    try:
        # Plain text formats
        if ext in (".txt", ".md", ".csv", ".tsv", ".html", ".htm",
                   ".json", ".xml", ".yaml", ".yml", ".log", ".py"):
            with open(target_path, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            if len(text) > MAX_OUTPUT_CHARS:
                text = text[:MAX_OUTPUT_CHARS] + "\n[... truncated ...]"
            return text

        # PDF (with page range support)
        if ext == ".pdf":
            page_set = _parse_page_range(pages) if pages else None
            try:
                import pdfplumber
                extracted = []
                with pdfplumber.open(target_path) as pdf:
                    total_pages = len(pdf.pages)
                    for i, page in enumerate(pdf.pages):
                        page_num = i + 1
                        if page_set and page_num not in page_set:
                            continue
                        t = page.extract_text() or ""
                        if t.strip():
                            extracted.append(f"[Page {page_num}]\n{t}")
                header = f"[PDF: {total_pages} total pages"
                if page_set:
                    header += f", showing pages {pages}"
                header += "]\n\n"
                text = header + "\n\n".join(extracted)
                if not extracted:
                    text = header + "[No extractable text on requested pages]"
                if len(text) > MAX_FILE_READ_CHARS:
                    text = text[:MAX_FILE_READ_CHARS] + "\n[... truncated ...]"
                return text
            except ImportError:
                try:
                    import fitz
                    doc = fitz.open(target_path)
                    extracted = []
                    for i, p in enumerate(doc):
                        page_num = i + 1
                        if page_set and page_num not in page_set:
                            continue
                        t = p.get_text()
                        if t.strip():
                            extracted.append(f"[Page {page_num}]\n{t}")
                    doc.close()
                    return "\n\n".join(extracted) or "[PDF: no text]"
                except ImportError:
                    return "[PDF: install pdfplumber or pymupdf]"

        # XLSX
        if ext in (".xlsx", ".xls"):
            try:
                from openpyxl import load_workbook
                wb = load_workbook(target_path, read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    parts.append(f"=== Sheet: {sheet_name} ===")
                    for row in ws.iter_rows(max_row=5000, values_only=True):
                        cells = [str(c) if c is not None else "" for c in row]
                        if any(c.strip() for c in cells):
                            parts.append("\t".join(cells))
                wb.close()
                return "\n".join(parts)
            except ImportError:
                return "[XLSX: install openpyxl]"

        # DOCX
        if ext == ".docx":
            try:
                from docx import Document
                doc = Document(target_path)
                paras = [p.text for p in doc.paragraphs if p.text.strip()]
                for table in doc.tables:
                    for row in table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            paras.append("\t".join(cells))
                return "\n".join(paras) or "[Empty document]"
            except ImportError:
                return "[DOCX: install python-docx]"

        # PPTX
        if ext == ".pptx":
            try:
                from pptx import Presentation
                prs = Presentation(target_path)
                slides = []
                for i, slide in enumerate(prs.slides, 1):
                    texts = []
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for para in shape.text_frame.paragraphs:
                                if para.text.strip():
                                    texts.append(para.text)
                        if shape.has_table:
                            for row in shape.table.rows:
                                cells = [c.text.strip() for c in row.cells
                                         if c.text.strip()]
                                if cells:
                                    texts.append("\t".join(cells))
                    if texts:
                        slides.append(f"=== Slide {i} ===\n" + "\n".join(texts))
                return "\n\n".join(slides) or "[Empty presentation]"
            except ImportError:
                return "[PPTX: install python-pptx]"

        # Fallback: try reading as text
        with open(target_path, "r", encoding="utf-8", errors="replace") as f:
            return f.read()[:MAX_OUTPUT_CHARS]

    except Exception as e:
        return f"[Read error for {filename}: {e}]"


# ── File manifest builder ─────────────────────────────────────────────────────

def _build_file_manifest(file_paths: list[str]) -> str:
    """Build a text manifest with [VISUAL] flags for files worth rasterizing."""
    if not file_paths:
        return ""
    try:
        from visual_file_processor import has_visual_content
        can_check_visual = True
    except ImportError:
        can_check_visual = False
    lines = [
        "\n## Input files provided\n",
        "Use `read_file` for text/data. Use `read_file_visual` for files "
        "marked [VISUAL] to see charts, diagrams, images. "
        "Files are also in the working directory for code.\n",
    ]
    for fpath in file_paths:
        name = os.path.basename(fpath)
        ext = os.path.splitext(fpath)[1].lower()
        if os.path.exists(fpath):
            size_kb = os.path.getsize(fpath) / 1024
            visual_tag = ""
            if can_check_visual and has_visual_content(fpath):
                visual_tag = " [VISUAL]"
            lines.append(f"  • **{name}** ({ext}, {size_kb:.1f} KB){visual_tag}")
        else:
            lines.append(f"  ⚠ **{name}** — FILE NOT FOUND")
    lines.append("")
    return "\n".join(lines)


def _exec_write_file(filename: str, content: str, working_dir: str) -> str:
    """Write text content to a file in the working directory."""
    try:
        filepath = os.path.join(working_dir, filename)
        os.makedirs(os.path.dirname(filepath) or working_dir, exist_ok=True)
        with open(filepath, "w", encoding="utf-8") as f:
            f.write(content)
        size = os.path.getsize(filepath)
        return f"File written: {filename} ({size:,} bytes)"
    except Exception as e:
        return f"[Write error for {filename}: {e}]"


def _exec_list_directory(path: str, working_dir: str) -> str:
    """List files and sizes in a directory."""
    try:
        target = path if os.path.isabs(path) else os.path.join(working_dir, path)
        if not os.path.isdir(target):
            return f"[Directory not found: {path}]"
        entries = []
        for entry in sorted(os.listdir(target)):
            if entry.startswith("."):
                continue
            full = os.path.join(target, entry)
            if os.path.isfile(full):
                size = os.path.getsize(full)
                ext = os.path.splitext(entry)[1]
                entries.append(f"  {entry:40s}  {size:>10,} bytes  {ext}")
            elif os.path.isdir(full):
                entries.append(f"  {entry:40s}  [directory]")
        if not entries:
            return f"[Directory is empty: {path}]"
        return f"Contents of {path} ({len(entries)} items):\n" + "\n".join(entries)
    except Exception as e:
        return f"[List directory error: {e}]"


async def _exec_convert_document(
    source_file: str, target_format: str, working_dir: str
) -> str:
    """Convert document formats via LibreOffice headless."""
    source_path = (
        source_file if os.path.isabs(source_file)
        else os.path.join(working_dir, source_file)
    )
    if not os.path.exists(source_path):
        return f"[Source file not found: {source_file}]"
    lo_bin = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo_bin:
        return (
            "[LibreOffice not found. Install: sudo apt install libreoffice. "
            "Alternative: use python_execute with reportlab (PDF) or "
            "openpyxl/python-docx for direct generation.]"
        )
    target_format = target_format.lstrip(".")
    try:
        proc = await asyncio.create_subprocess_exec(
            lo_bin, "--headless", "--convert-to", target_format,
            "--outdir", working_dir, source_path,
            stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=CONVERT_TIMEOUT)
        except asyncio.TimeoutError:
            proc.kill()
            await proc.communicate()
            return f"[TIMEOUT: Conversion exceeded {CONVERT_TIMEOUT}s]"
        if proc.returncode != 0:
            return f"[Conversion failed: {stderr.decode()[:500]}]"
        stem = os.path.splitext(os.path.basename(source_file))[0]
        output_name = f"{stem}.{target_format}"
        output_path = os.path.join(working_dir, output_name)
        if os.path.exists(output_path):
            return f"Converted: {os.path.basename(source_file)} → {output_name} ({os.path.getsize(output_path):,} bytes)"
        return f"Conversion completed but {output_name} not found. Check list_directory."
    except Exception as e:
        return f"[Conversion error: {e}]"


def _snapshot_directory(path: str) -> set[str]:
    """Take a snapshot of filenames in a directory for before/after comparison."""
    if not os.path.isdir(path):
        return set()
    return {e for e in os.listdir(path) if os.path.isfile(os.path.join(path, e)) and not e.startswith("_qwen_")}


# ── Staging directory ─────────────────────────────────────────────────────────

def _prepare_staging_dir(task: ResearchTask) -> str:
    """
    Create a SEPARATE staging directory with symlinks to input files.

    This is the model's working directory. Input symlinks, scratch .py files,
    and intermediate artifacts all live here. ONLY deliverable files get
    copied to task.output_files_dir after the run completes.

    Staging dir is ALWAYS a fresh temp directory, never output_files_dir.
    """
    staging = tempfile.mkdtemp(prefix=f"qwen_staging_{task.task_id}_")

    for fpath in task.file_paths:
        if not os.path.exists(fpath):
            continue
        dest = os.path.join(staging, os.path.basename(fpath))
        if not os.path.exists(dest):
            try:
                os.symlink(os.path.abspath(fpath), dest)
            except (OSError, NotImplementedError):
                shutil.copy2(fpath, dest)

    return staging


# ── Tool fingerprint tracker (deduplication + stagnation detection) ────────────

class _ToolTracker:
    """
    Tracks tool call fingerprints to detect duplicates and stagnation.

    A fingerprint is (tool_name, canonical_args_hash). If the same fingerprint
    appears again, we return the cached result instead of re-executing.

    Stagnation: if N consecutive iterations produce zero new fingerprints,
    the model is going in circles.
    """

    def __init__(self, max_stagnant: int = 3):
        self._seen: dict[str, str] = {}  # fingerprint → cached result
        self._stagnant_count: int = 0
        self._max_stagnant: int = max_stagnant
        self._last_iter_had_new: bool = True

    def fingerprint(self, tool_name: str, tool_args: dict) -> str:
        """Create a canonical fingerprint for a tool call."""
        # Normalize: sort keys, strip whitespace from string values
        def _normalize(obj):
            if isinstance(obj, dict):
                return tuple(sorted((k, _normalize(v)) for k, v in obj.items()))
            if isinstance(obj, str):
                return obj.strip()[:500]  # cap long code blocks
            return obj
        return f"{tool_name}::{hash(_normalize(tool_args))}"

    def check(self, tool_name: str, tool_args: dict) -> tuple[bool, str]:
        """
        Check if this tool call was seen before.

        Returns:
            (is_duplicate, cached_result_or_empty)
        """
        fp = self.fingerprint(tool_name, tool_args)
        if fp in self._seen:
            return True, self._seen[fp]
        return False, ""

    def record(self, tool_name: str, tool_args: dict, result: str):
        """Record a tool call and its result."""
        fp = self.fingerprint(tool_name, tool_args)
        is_new = fp not in self._seen
        self._seen[fp] = result
        return is_new

    def update_stagnation(self, had_new_this_iter: bool):
        """Update stagnation counter. Call once per iteration."""
        if had_new_this_iter:
            self._stagnant_count = 0
        else:
            self._stagnant_count += 1

    @property
    def is_stagnant(self) -> bool:
        return self._stagnant_count >= self._max_stagnant

    @property
    def stagnant_count(self) -> int:
        return self._stagnant_count


# ── The adapter ───────────────────────────────────────────────────────────────

class QwenAdapter:
    """
    Deep Research Agent using Qwen models via OpenAI-compatible API.

    Transparent agentic loop with locally-executed tools.
    Full observability: every tool call, result, and reasoning step is logged.

    Supports:
      - Qwen 3.6-27B (dense), Qwen 3.6-35B-A3B (MoE)
      - Qwen 3.5 family (397B-A17B, 27B, etc.)
      - Qwen 3.7 Max (DashScope API)
      - Any OpenAI-compatible endpoint (vLLM, Ollama, DashScope)
    """

    def __init__(
        self,
        api_key: Optional[str] = None,
        model: Optional[str] = None,
        base_url: Optional[str] = None,
        dry_run: bool = False,
        max_tool_rounds: int = 100,
    ):
        self.api_key = (
            api_key
            or os.environ.get("QWEN_API_KEY")
            or os.environ.get("OPENAI_API_KEY")
            or "EMPTY"
        )
        self.model = (
            model
            or os.environ.get("QWEN_MODEL")
            or "Qwen/Qwen3.6-27B"
        )
        self.base_url = (
            base_url
            or os.environ.get("QWEN_BASE_URL")
            or "http://localhost:8000/v1"
        )
        self.dry_run = dry_run
        self.max_tool_rounds = int(os.environ.get("QWEN_MAX_ROUNDS", str(max_tool_rounds)))

        if not dry_run:
            try:
                from openai import AsyncOpenAI
                self.client = AsyncOpenAI(
                    api_key=self.api_key,
                    base_url=self.base_url,
                    timeout=600,  # 10 min per API call
                )
            except ImportError:
                raise ImportError("pip install openai  (required for Qwen adapter)")

    # ── Main entry point ─────────────────────────────────────────────────

    async def run(self, task: ResearchTask) -> AgentResult:
        started_at = datetime.now(timezone.utc)

        print(f"[Qwen] Starting {task.task_id} "
              f"(model={self.model}, dry_run={self.dry_run})")

        if self.dry_run:
            return self._dry_run_result(task, started_at)

        # Prepare staging directory with input files
        staging_dir = _prepare_staging_dir(task)
        files_before = _snapshot_directory(staging_dir)
        print(f"[Qwen] Staging dir: {staging_dir} ({len(files_before)} input files)")

        # Build conversation
        system_prompt = _build_system_prompt(task)
        file_manifest = _build_file_manifest(task.file_paths)
        user_content = task.prompt + file_manifest

        messages = [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_content},
        ]

        tools = _build_tool_definitions(task)
        print(f"[Qwen] Tools: {[t['function']['name'] for t in tools]}")

        # Agentic loop state
        iterations       = 0
        total_input_tok  = 0
        total_output_tok = 0
        tool_call_log: list[ToolCall] = []
        forced_stop      = False
        final_text       = ""
        tracker          = _ToolTracker(max_stagnant=3)
        malformed_retries = 0

        while True:
            iterations += 1

            # ── Guard: max iterations (absolute safety net) ──────
            if iterations > self.max_tool_rounds:
                print(f"  [GUARD] Max tool rounds ({self.max_tool_rounds}) hit.")
                forced_stop = True
                break

            # ── Guard: stagnation ────────────────────────────────
            if tracker.is_stagnant:
                print(f"  [STAGNATION] {tracker.stagnant_count} consecutive "
                      f"unproductive iterations. Forcing conclusion.")
                messages.append({"role": "user", "content": (
                    "You are repeating previous actions without making progress. "
                    "Stop using tools and produce your final research report now."
                )})
                try:
                    stag_resp = await self.client.chat.completions.create(
                        model=self.model, messages=messages,
                        max_tokens=16384, temperature=0.7,
                    )
                    if stag_resp.usage:
                        total_input_tok += stag_resp.usage.prompt_tokens or 0
                        total_output_tok += stag_resp.usage.completion_tokens or 0
                    final_text = stag_resp.choices[0].message.content or ""
                except Exception:
                    pass
                break

            # ── Guard: budget ────────────────────────────────────
            current_cost = estimate_cost(
                self.model, total_input_tok, total_output_tok
            )
            if current_cost > task.max_cost_usd:
                print(f"  [GUARD] Budget ${task.max_cost_usd} exceeded "
                      f"(${current_cost:.2f}).")
                forced_stop = True
                break

            # ── Guard: timeout ───────────────────────────────────
            elapsed = (datetime.now(timezone.utc) - started_at).total_seconds()
            if elapsed > task.timeout_seconds:
                print(f"  [GUARD] Timeout {task.timeout_seconds}s exceeded.")
                forced_stop = True
                break

            print(f"  [ITER {iterations}] Calling Qwen... "
                  f"msgs={len(messages)}  cost=${current_cost:.4f}  "
                  f"elapsed={elapsed:.0f}s  stagnant={tracker.stagnant_count}")

            # ── API call ─────────────────────────────────────────
            try:
                create_kwargs = dict(
                    model=self.model,
                    messages=messages,
                    max_tokens=16384,
                    temperature=0.7,
                    top_p=0.8,
                    extra_body={
                        "reasoning": {"enabled": True},
                    },
                )
                if tools:
                    create_kwargs["tools"] = tools
                    create_kwargs["tool_choice"] = "auto"

                response = await self.client.chat.completions.create(
                    **create_kwargs
                )

            except Exception as e:
                error_str = str(e)
                # Bug B: Model generated malformed tool call arguments.
                # Retry with a hint instead of killing the task.
                if "function.arguments" in error_str and "JSON" in error_str and malformed_retries < 2:
                    malformed_retries += 1
                    print(f"  [RETRY] Malformed tool args (attempt {malformed_retries}/2), nudging model...")
                    messages.append({
                        "role": "user",
                        "content": (
                            "Your last tool call had malformed arguments (not valid JSON). "
                            "Please try again, ensuring all tool call arguments are "
                            "properly JSON-encoded. For python_execute, make sure the "
                            "code string has escaped quotes and newlines."
                        ),
                    })
                    continue  # retry this iteration
                print(f"  [ERROR] API call failed: {e}")
                return self._error_result(
                    task, started_at, error_str,
                    total_input_tok, total_output_tok,
                    tool_call_log, iterations,
                )

            # ── Extract usage ────────────────────────────────────
            if response.usage:
                total_input_tok += response.usage.prompt_tokens or 0
                total_output_tok += response.usage.completion_tokens or 0

            choice = response.choices[0]
            finish_reason = choice.finish_reason
            message = choice.message

            print(f"  [ITER {iterations}] finish={finish_reason}  "
                  f"in={response.usage.prompt_tokens if response.usage else '?'}  "
                  f"out={response.usage.completion_tokens if response.usage else '?'}")

            # ── finish_reason == None (provider didn't set it) ────
            if finish_reason is None:
                if message.tool_calls and len(message.tool_calls) > 0:
                    finish_reason = "tool_calls"
                    print(f"  [FIX] finish_reason=None with tool_calls → treating as 'tool_calls'")
                elif message.content:
                    finish_reason = "stop"
                    print(f"  [FIX] finish_reason=None with content → treating as 'stop'")
                else:
                    print(f"  [WARN] finish_reason=None with no content or tools, skipping")
                    continue

            # ── finish_reason == "stop" ──────────────────────────
            if finish_reason == "stop":
                final_text = message.content or ""
                break

            # ── finish_reason == "tool_calls" ────────────────────
            if finish_reason == "tool_calls" or (
                message.tool_calls and len(message.tool_calls) > 0
            ):
                # Append assistant message with tool calls + reasoning preservation
                assistant_msg = {"role": "assistant", "content": message.content or ""}
                # Preserve reasoning_details for thinking continuity across turns
                reasoning_details = getattr(message, "reasoning_details", None)
                if reasoning_details:
                    assistant_msg["reasoning_details"] = reasoning_details
                if message.tool_calls:
                    assistant_msg["tool_calls"] = [
                        {
                            "id": tc.id,
                            "type": "function",
                            "function": {
                                "name": tc.function.name,
                                "arguments": tc.function.arguments,
                            },
                        }
                        for tc in message.tool_calls
                    ]
                messages.append(assistant_msg)

                # Execute each tool call (with dedup)
                iter_had_new = False
                for tc in message.tool_calls:
                    tool_name = tc.function.name
                    try:
                        tool_args = json.loads(tc.function.arguments)
                    except json.JSONDecodeError:
                        tool_args = {"_raw": tc.function.arguments}

                    arg_preview = json.dumps(tool_args, ensure_ascii=False)[:100]

                    # Dedup check
                    is_dup, cached = tracker.check(tool_name, tool_args)
                    if is_dup:
                        result_text = (
                            f"[CACHED: You already called {tool_name} with these "
                            f"arguments. The result is in your context above. "
                            f"Here it is again:]\n\n{cached[:2000]}"
                        )
                        print(f"  [DEDUP] {tool_name}({arg_preview}) → cached")
                        tool_duration_ms = 0
                    else:
                        print(f"  [TOOL] {tool_name}({arg_preview})")
                        tool_start = time.monotonic()
                        result_text = await self._execute_tool(
                            tool_name, tool_args, task, staging_dir
                        )
                        tool_duration_ms = int((time.monotonic() - tool_start) * 1000)
                        is_new = tracker.record(tool_name, tool_args, result_text)
                        if is_new:
                            iter_had_new = True
                        print(f"  [TOOL] → {len(result_text)} chars, "
                              f"{tool_duration_ms}ms {'(new)' if is_new else '(seen)'}")

                    tool_call_log.append(ToolCall(
                        iteration=iterations,
                        tool_name=tool_name,
                        tool_input=tool_args,
                        result_preview=result_text[:500],
                        result_tokens=len(result_text) // 4,
                        timestamp=datetime.now(timezone.utc).isoformat(),
                        duration_ms=tool_duration_ms,
                    ))

                    messages.append({
                        "role": "tool",
                        "tool_call_id": tc.id,
                        "content": result_text,
                    })

                tracker.update_stagnation(iter_had_new)
                continue

            # ── finish_reason == "length" ────────────────────────
            if finish_reason == "length":
                print(f"  [WARN] max_tokens hit, asking to continue...")
                messages.append({
                    "role": "assistant",
                    "content": message.content or "",
                })
                messages.append({
                    "role": "user",
                    "content": (
                        "Your response was cut off at the token limit. "
                        "Please continue from where you left off."
                    ),
                })
                continue

            # ── Unknown finish reason ────────────────────────────
            print(f"  [WARN] Unexpected finish_reason={finish_reason}")
            if message.content:
                final_text = message.content
            forced_stop = True
            break

        # ── Forced-stop synthesis ─────────────────────────────────────────
        # If final_text contains raw <tool_call> XML, the model was mid-tool-call
        # when the iteration limit hit. This is not a real response — clear it
        # so synthesis kicks in.
        if final_text and "<tool_call>" in final_text:
            print(f"  [CLEANUP] final_text contains raw <tool_call> XML ({len(final_text)} chars). Clearing for synthesis.")
            final_text = ""
        if forced_stop and not final_text:
            print(f"  [SYNTH] Forced stop — synthesizing from context...")
            messages.append({"role": "user", "content": (
                "You have reached the research limit. Synthesize all findings "
                "into a final well-cited research report. Note incomplete areas."
            )})
            try:
                synth = await self.client.chat.completions.create(
                    model=self.model, messages=messages, max_tokens=8192, temperature=0.7,
                )
                if synth.usage:
                    total_input_tok += synth.usage.prompt_tokens or 0
                    total_output_tok += synth.usage.completion_tokens or 0
                final_text = synth.choices[0].message.content or ""
            except Exception as e:
                final_text = f"[Synthesis error: {e}]"

        # ── Bug A fix: Reconstruct if final_text is empty but work was done ──
        if not final_text and tool_call_log:
            print(f"  [RECOVER] final_text empty despite {len(tool_call_log)} tool calls. "
                  f"Reconstructing from conversation...")
            # Try to get the model to summarize its own work
            messages.append({
                "role": "user",
                "content": (
                    "You completed your analysis but did not produce a final response. "
                    "Please write your complete research report now, synthesizing all "
                    "the findings from your tool calls above."
                ),
            })
            try:
                recover = await self.client.chat.completions.create(
                    model=self.model, messages=messages,
                    max_tokens=16384, temperature=0.7,
                )
                if recover.usage:
                    total_input_tok += recover.usage.prompt_tokens or 0
                    total_output_tok += recover.usage.completion_tokens or 0
                final_text = recover.choices[0].message.content or ""
                if final_text:
                    print(f"  [RECOVER] Reconstructed {len(final_text)} chars")
            except Exception as e:
                print(f"  [RECOVER] Failed: {e}")

            # Last resort: extract from last assistant messages
            if not final_text:
                for msg in reversed(messages):
                    if msg.get("role") == "assistant" and msg.get("content", "").strip():
                        final_text = msg["content"]
                        print(f"  [RECOVER] Used last assistant message ({len(final_text)} chars)")
                        break

        # ── Layer 2: Dynamic minimum tool-call gate ──────────────────────
        if final_text and not forced_stop:
            files_read = sum(1 for tc in tool_call_log if tc.tool_name in ("read_file", "read_file_visual", "search_in_file"))
            code_calls = sum(1 for tc in tool_call_log if tc.tool_name == "python_execute")
            search_calls = sum(1 for tc in tool_call_log if tc.tool_name == "web_search")
            min_reads = len(task.file_paths)
            min_code = 1 if task.output_formats else 0
            min_search = 1 if (task.web_search_enabled and task.iat_type == "IAT-3") else 0
            gate_failures = []
            if files_read < min_reads:
                gate_failures.append(f"read {files_read}/{min_reads} files")
            if code_calls < min_code:
                gate_failures.append(f"code_calls={code_calls}, need >={min_code}")
            if search_calls < min_search:
                gate_failures.append(f"web_searches={search_calls}, need >={min_search}")
            if gate_failures:
                print(f"  [LAYER2] Gate FAILED: {gate_failures}")
                messages.append({"role": "assistant", "content": final_text})
                messages.append({"role": "user", "content": (
                    "Your response is INCOMPLETE:\n"
                    + "\n".join(f"  - {f}" for f in gate_failures)
                    + "\n\nRedo the missing steps, then provide your complete response."
                )})
                final_text = ""
                gate_tracker = _ToolTracker(max_stagnant=3)
                while not gate_tracker.is_stagnant:
                    iterations += 1
                    try:
                        gate_resp = await self.client.chat.completions.create(
                            model=self.model, messages=messages,
                            max_tokens=16384, temperature=0.7, tools=tools, tool_choice="auto",
                        )
                        if gate_resp.usage:
                            total_input_tok += gate_resp.usage.prompt_tokens or 0
                            total_output_tok += gate_resp.usage.completion_tokens or 0
                        choice = gate_resp.choices[0]
                        if choice.finish_reason == "stop":
                            final_text = choice.message.content or ""
                            break
                        elif choice.message.tool_calls:
                            asst = {"role": "assistant", "content": choice.message.content or ""}
                            _rd = getattr(choice.message, "reasoning_details", None)
                            if _rd: asst["reasoning_details"] = _rd
                            asst["tool_calls"] = [
                                {"id": tc.id, "type": "function",
                                 "function": {"name": tc.function.name, "arguments": tc.function.arguments}}
                                for tc in choice.message.tool_calls
                            ]
                            messages.append(asst)
                            gate_new = False
                            for tc in choice.message.tool_calls:
                                ta = json.loads(tc.function.arguments) if tc.function.arguments else {}
                                rt = await self._execute_tool(tc.function.name, ta, task, staging_dir)
                                if gate_tracker.record(tc.function.name, ta, rt):
                                    gate_new = True
                                tool_call_log.append(ToolCall(
                                    iteration=iterations, tool_name=tc.function.name,
                                    tool_input=ta, result_preview=rt[:500],
                                    result_tokens=len(rt)//4,
                                    timestamp=datetime.now(timezone.utc).isoformat(), duration_ms=0,
                                ))
                                messages.append({"role": "tool", "tool_call_id": tc.id, "content": rt})
                            gate_tracker.update_stagnation(gate_new)
                    except Exception as e:
                        print(f"  [LAYER2] Remediation error: {e}")
                        break
                print(f"  [LAYER2] Done, final_text={'YES' if final_text else 'NO'}")

        # ── Layer 4: Self-verification ───────────────────────────────────
        if final_text and not forced_stop:
            print(f"  [LAYER4] Running self-verification...")
            verify_prompt = (
                f"## SELF-VERIFICATION\n\n"
                f"**ORIGINAL PROMPT:**\n{task.prompt}\n\n"
                f"**YOUR RESPONSE:**\n{final_text}\n\n"
                f"**INPUT FILES READ:** {sum(1 for tc in tool_call_log if tc.tool_name in ('read_file', 'read_file_visual'))}/{len(task.file_paths)}\n\n"
                f"**CHECKLIST:**\n"
                f"1. Did you read ALL {len(task.file_paths)} input files?\n"
                f"2. Does your response directly answer what was asked?\n"
                f"3. Are calculations verified with python_execute?\n"
                f"4. Are claims cited?\n"
                f"5. Did you generate required output files?\n"
                f"6. Any charts/diagrams to examine with read_file_visual?\n\n"
                f"If ANY check fails, fix it now. If ALL pass: VERIFICATION: PASS"
            )
            messages.append({"role": "assistant", "content": final_text})
            messages.append({"role": "user", "content": verify_prompt})
            try:
                v_tracker = _ToolTracker(max_stagnant=3)
                verify_resp = await self.client.chat.completions.create(
                    model=self.model, messages=messages,
                    max_tokens=16384, temperature=0.5, tools=tools, tool_choice="auto",
                )
                if verify_resp.usage:
                    total_input_tok += verify_resp.usage.prompt_tokens or 0
                    total_output_tok += verify_resp.usage.completion_tokens or 0
                v_choice = verify_resp.choices[0]
                if v_choice.message.tool_calls:
                    print(f"  [LAYER4] Model found issues, running fixes...")
                    asst = {"role": "assistant", "content": v_choice.message.content or ""}
                    _rd = getattr(v_choice.message, "reasoning_details", None)
                    if _rd: asst["reasoning_details"] = _rd
                    asst["tool_calls"] = [
                        {"id": tc.id, "type": "function",
                         "function": {"name": tc.function.name, "arguments": tc.function.arguments}}
                        for tc in v_choice.message.tool_calls
                    ]
                    messages.append(asst)
                    for tc in v_choice.message.tool_calls:
                        ta = json.loads(tc.function.arguments) if tc.function.arguments else {}
                        rt = await self._execute_tool(tc.function.name, ta, task, staging_dir)
                        v_tracker.record(tc.function.name, ta, rt)
                        tool_call_log.append(ToolCall(
                            iteration=iterations+1, tool_name=tc.function.name,
                            tool_input=ta, result_preview=rt[:500],
                            result_tokens=len(rt)//4,
                            timestamp=datetime.now(timezone.utc).isoformat(), duration_ms=0,
                        ))
                        messages.append({"role": "tool", "tool_call_id": tc.id, "content": rt})
                    while not v_tracker.is_stagnant:
                        vr = await self.client.chat.completions.create(
                            model=self.model, messages=messages,
                            max_tokens=16384, temperature=0.5, tools=tools, tool_choice="auto",
                        )
                        if vr.usage:
                            total_input_tok += vr.usage.prompt_tokens or 0
                            total_output_tok += vr.usage.completion_tokens or 0
                        vc = vr.choices[0]
                        if vc.finish_reason == "stop" or not vc.message.tool_calls:
                            if vc.message.content:
                                final_text = vc.message.content
                            break
                        asst2 = {"role": "assistant", "content": vc.message.content or ""}
                        _rd = getattr(vc.message, "reasoning_details", None)
                        if _rd: asst2["reasoning_details"] = _rd
                        asst2["tool_calls"] = [
                            {"id": tc.id, "type": "function",
                             "function": {"name": tc.function.name, "arguments": tc.function.arguments}}
                            for tc in vc.message.tool_calls
                        ]
                        messages.append(asst2)
                        v_new = False
                        for tc in vc.message.tool_calls:
                            ta = json.loads(tc.function.arguments) if tc.function.arguments else {}
                            rt = await self._execute_tool(tc.function.name, ta, task, staging_dir)
                            if v_tracker.record(tc.function.name, ta, rt):
                                v_new = True
                            messages.append({"role": "tool", "tool_call_id": tc.id, "content": rt})
                        v_tracker.update_stagnation(v_new)
                    messages.append({"role": "user", "content": "Provide your complete corrected final response."})
                    fr = await self.client.chat.completions.create(
                        model=self.model, messages=messages, max_tokens=16384, temperature=0.7,
                    )
                    if fr.usage:
                        total_input_tok += fr.usage.prompt_tokens or 0
                        total_output_tok += fr.usage.completion_tokens or 0
                    revised = fr.choices[0].message.content or ""
                    if revised and len(revised) > len(final_text) * 0.5:
                        final_text = revised
                    print(f"  [LAYER4] Response revised ({len(final_text)} chars)")
                elif "VERIFICATION: PASS" in (v_choice.message.content or ""):
                    print(f"  [LAYER4] Verification PASSED")
                else:
                    revised = v_choice.message.content or ""
                    if revised and "VERIFICATION: PASS" not in revised and len(revised) > 500:
                        final_text = revised
                        print(f"  [LAYER4] Response revised inline ({len(revised)} chars)")
                    else:
                        print(f"  [LAYER4] Verification passed (implicit)")
            except Exception as e:
                print(f"  [LAYER4] Verification error (non-fatal): {e}")

        # ── Conclusion step: model classifies files ──────────────────────
        deliverable_files = []
        if final_text:
            print(f"  [CONCLUSION] Classifying output files...")
            files_after = _snapshot_directory(staging_dir)
            new_files = sorted(files_after - files_before)
            new_files = [f for f in new_files if not f.endswith(".pyc") and not f.startswith(".")]
            if new_files:
                file_listing = "\n".join(
                    f"  {f:40s}  {os.path.getsize(os.path.join(staging_dir, f)):>10,} bytes"
                    for f in new_files if os.path.isfile(os.path.join(staging_dir, f))
                )
                conclusion_prompt = (
                    f"ORIGINAL PROMPT:\n{task.prompt}\n\n"
                    f"ALL FILES IN YOUR WORKING DIRECTORY:\n{file_listing}\n\n"
                    f"DETECTOR HINT: {task.output_formats or '(none detected)'}\n\n"
                    "Respond in EXACTLY this JSON format, nothing else:\n"
                    "{\n"
                    '  "required_formats": ["docx"],\n'
                    '  "deliverable_files": ["Report.docx"],\n'
                    '  "intermediate_files": ["script.py", "data.json"],\n'
                    '  "missing_formats": []\n'
                    "}"
                )
                try:
                    conc_resp = await self.client.chat.completions.create(
                        model=self.model,
                        messages=[{"role": "user", "content": conclusion_prompt}],
                        max_tokens=1024, temperature=0.1,
                    )
                    if conc_resp.usage:
                        total_input_tok += conc_resp.usage.prompt_tokens or 0
                        total_output_tok += conc_resp.usage.completion_tokens or 0
                    raw = conc_resp.choices[0].message.content or ""
                    raw = re.sub(r'^```json\s*|```\s*$', '', raw.strip())
                    conclusion = json.loads(raw)
                    deliverable_files = conclusion.get("deliverable_files", [])
                    intermediate = conclusion.get("intermediate_files", [])
                    missing = conclusion.get("missing_formats", [])
                    model_formats = conclusion.get("required_formats", [])
                    print(f"  [CONCLUSION] Deliverables: {deliverable_files}")
                    print(f"  [CONCLUSION] Intermediate (filtered): {intermediate}")
                    if missing:
                        print(f"  [CONCLUSION] Missing formats: {missing} — generating...")
                        messages.append({"role": "user", "content": (
                            f"You identified missing deliverable formats: {missing}. "
                            f"Generate them now using python_execute."
                        )})
                        miss_tracker = _ToolTracker(max_stagnant=3)
                        while not miss_tracker.is_stagnant:
                            mr = await self.client.chat.completions.create(
                                model=self.model, messages=messages,
                                max_tokens=16384, temperature=0.7, tools=tools, tool_choice="auto",
                            )
                            if mr.usage:
                                total_input_tok += mr.usage.prompt_tokens or 0
                                total_output_tok += mr.usage.completion_tokens or 0
                            mc = mr.choices[0]
                            if mc.finish_reason == "stop" or not mc.message.tool_calls:
                                break
                            asst = {"role": "assistant", "content": mc.message.content or ""}
                            _rd = getattr(mc.message, "reasoning_details", None)
                            if _rd: asst["reasoning_details"] = _rd
                            asst["tool_calls"] = [
                                {"id": tc.id, "type": "function",
                                 "function": {"name": tc.function.name, "arguments": tc.function.arguments}}
                                for tc in mc.message.tool_calls
                            ]
                            messages.append(asst)
                            mn = False
                            for tc in mc.message.tool_calls:
                                ta = json.loads(tc.function.arguments) if tc.function.arguments else {}
                                rt = await self._execute_tool(tc.function.name, ta, task, staging_dir)
                                if miss_tracker.record(tc.function.name, ta, rt):
                                    mn = True
                                messages.append({"role": "tool", "tool_call_id": tc.id, "content": rt})
                            miss_tracker.update_stagnation(mn)
                        files_now = _snapshot_directory(staging_dir)
                        for f in sorted(files_now - files_before):
                            if f not in deliverable_files and not f.endswith((".py", ".pyc")):
                                ext = os.path.splitext(f)[1].lstrip(".")
                                if ext in [m.lstrip(".") for m in missing]:
                                    deliverable_files.append(f)
                    if task.output_formats and model_formats and set(task.output_formats) != set(model_formats):
                        print(f"  [CONCLUSION] Format mismatch — detector: {task.output_formats}, model: {model_formats}")
                except (json.JSONDecodeError, Exception) as e:
                    print(f"  [CONCLUSION] Parse error ({e}), falling back to snapshot")
                    deliverable_files = [f for f in new_files if not f.endswith((".py", ".pyc", ".json"))]

            # Safety net: scan disk for deliverable files the model missed
            DELIVERABLE_EXTS = {".docx", ".pptx", ".xlsx", ".pdf", ".csv", ".png", ".jpg", ".wav"}
            response_docx_pattern = f"{task.task_id}"
            for f in new_files:
                if f in deliverable_files:
                    continue
                ext = os.path.splitext(f)[1].lower()
                if ext in DELIVERABLE_EXTS and response_docx_pattern not in f:
                    src = os.path.join(staging_dir, f)
                    if os.path.isfile(src) and os.path.getsize(src) > 100:
                        deliverable_files.append(f)
                        print(f"  [SAFETY] Recovered missed deliverable: {f}")

        # ── Ensure output_files_dir exists ───────────────────────────────
        if not task.output_files_dir:
            task.output_files_dir = os.path.join(
                tempfile.gettempdir(), "dra_output_files", task.task_id
            )

        # ── Always generate response docx ────────────────────────────────
        response_docx = None
        response_docx_name = ""
        if final_text:
            response_docx = self._generate_response_docx(final_text, task.task_id, staging_dir)
            if response_docx:
                response_docx_name = os.path.basename(response_docx)
                print(f"  [DOCX] Response saved: {response_docx_name}")

        # ── Copy ONLY deliverables + response docx to output dir ─────────
        output_files = []
        output_file_errors: dict = {}
        os.makedirs(task.output_files_dir, exist_ok=True)

        for fname in deliverable_files:
            src = os.path.join(staging_dir, fname)
            if os.path.isfile(src):
                dest = os.path.join(task.output_files_dir, fname)
                shutil.copy2(src, dest)
                output_files.append(dest)
                print(f"  [FILE] Deliverable: {fname}")

        if response_docx:
            dest = os.path.join(task.output_files_dir, response_docx_name)
            if not os.path.exists(dest):
                shutil.copy2(response_docx, dest)
            output_files.append(dest)

        # ── Build result ──────────────────────────────────────────────────
        completed_at = datetime.now(timezone.utc)
        total_cost = estimate_cost(self.model, total_input_tok, total_output_tok)
        completed = bool(final_text) and not bool(output_file_errors)

        print(f"  [DONE] iters={iterations}  cost=${total_cost:.4f}  "
              f"tools={len(tool_call_log)}  files={len(output_files)}  "
              f"forced_stop={forced_stop}")

        return AgentResult(
            task_id=task.task_id,
            agent="qwen",
            model=self.model,
            response_text=final_text,
            citations=self._extract_citations(final_text),
            tool_call_log=tool_call_log,
            input_tokens=total_input_tok,
            output_tokens=total_output_tok,
            total_cost_usd=total_cost,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=iterations,
            completed=completed,
            forced_stop=forced_stop,
            error=(
                f"File generation failed: {list(output_file_errors.keys())}"
                if output_file_errors else None
            ),
            output_files=output_files,
            output_file_errors=output_file_errors,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    # ── Tool dispatch ────────────────────────────────────────────────────

    async def _execute_tool(
        self, tool_name: str, tool_args: dict,
        task: ResearchTask, staging_dir: str,
    ) -> str:
        """Route a tool call to the appropriate executor."""
        if tool_name == "python_execute":
            return await _exec_python(tool_args.get("code", ""), staging_dir)
        elif tool_name == "bash_execute":
            return await _exec_bash(tool_args.get("command", ""), staging_dir)
        elif tool_name == "web_search":
            return await _exec_web_search(tool_args.get("query", ""), tool_args.get("max_results", 8))
        elif tool_name == "web_fetch":
            return await _exec_web_fetch(tool_args.get("url", ""))
        elif tool_name == "read_file":
            return _exec_read_file(tool_args.get("filename", ""), task.file_paths, tool_args.get("pages", ""))
        elif tool_name == "write_file":
            return _exec_write_file(tool_args.get("filename", ""), tool_args.get("content", ""), staging_dir)
        elif tool_name == "list_directory":
            return _exec_list_directory(tool_args.get("path", "."), staging_dir)
        elif tool_name == "convert_document":
            return await _exec_convert_document(tool_args.get("source_file", ""), tool_args.get("target_format", "pdf"), staging_dir)
        elif tool_name == "read_file_visual":
            return await self._exec_read_file_visual(tool_args.get("filename", ""), tool_args.get("context_hint", ""), task, tool_args.get("pages", ""))
        elif tool_name == "search_in_file":
            return _exec_search_in_file(tool_args.get("filename", ""), tool_args.get("query", ""), task.file_paths, tool_args.get("max_results", 10))
        else:
            return json.dumps({"error": f"Unknown tool: {tool_name}"})

    # ── Visual file reading ──────────────────────────────────────────────

    async def _exec_read_file_visual(
        self, filename: str, context_hint: str, task: ResearchTask, pages: str = ""
    ) -> str:
        """Read visual content from a file using the vision model pipeline."""
        # Find the file
        target_path = None
        for fpath in task.file_paths:
            if os.path.basename(fpath).lower() == filename.lower():
                target_path = fpath
                break
        if target_path is None and os.path.exists(filename):
            target_path = filename
        if target_path is None:
            available = [os.path.basename(f) for f in task.file_paths]
            return json.dumps({"error": f"File not found: {filename}", "available_files": available})

        try:
            from visual_file_processor import rasterize_file, describe_visual_content
        except ImportError:
            return (
                "[Visual file reading unavailable: visual_file_processor.py not found. "
                "Falling back to text extraction.]"
                "\n\n" + _exec_read_file(filename, task.file_paths)
            )

        print(f"  [VISUAL] Rasterizing {filename}...")
        pages = rasterize_file(target_path)
        if not pages:
            return (
                f"[Could not rasterize {filename} — no visual content or "
                f"missing dependencies (pymupdf/pdf2image/Pillow). "
                f"Falling back to text extraction.]\n\n"
                + _exec_read_file(filename, task.file_paths)
            )

        print(f"  [VISUAL] {len(pages)} pages → calling vision model...")
        description = await describe_visual_content(
            pages, filename, context_hint=context_hint,
        )
        print(f"  [VISUAL] Got {len(description)} chars from vision model")

        # Combine text extraction + visual description
        text_content = _exec_read_file(filename, task.file_paths)
        return (
            f"{text_content}\n\n"
            f"══ VISUAL CONTENT (from vision model) ══\n"
            f"{description}"
        )

    # ── Output file collection ───────────────────────────────────────────

    def _collect_output_files(
        self, staging_dir: str, task: ResearchTask, files_before: set[str]
    ) -> list[str]:
        """Detect new files created during the agentic loop via before/after snapshot."""
        if not os.path.isdir(staging_dir):
            return []

        files_after = _snapshot_directory(staging_dir)
        new_files = files_after - files_before
        new_files = {f for f in new_files if not f.endswith(".pyc") and not f.startswith(".")}

        if not new_files:
            return []

        found = []
        for entry in sorted(new_files):
            src = os.path.join(staging_dir, entry)
            if not os.path.isfile(src):
                continue
            if task.output_files_dir and task.output_files_dir != staging_dir:
                os.makedirs(task.output_files_dir, exist_ok=True)
                dest = os.path.join(task.output_files_dir, entry)
                shutil.copy2(src, dest)
                found.append(dest)
                print(f"  [FILE] Captured: {entry} → {dest} ({os.path.getsize(dest):,} bytes)")
            else:
                found.append(src)
                print(f"  [FILE] Found: {entry} ({os.path.getsize(src):,} bytes)")
        return found

    # ── Parsing helpers ──────────────────────────────────────────────────

    def _extract_citations(self, text: str) -> list[dict]:
        """Extract numbered citations from the final report."""
        pattern = r'\[(\d+)\]\s+(.+?)\s*[—\-]+\s*(https?://\S+)'
        return [
            {
                "index": int(m.group(1)),
                "title": m.group(2).strip(),
                "url": m.group(3).strip(),
            }
            for m in re.finditer(pattern, text)
        ]

    # ── Response docx generation ─────────────────────────────────────────

    def _generate_response_docx(
        self, final_text: str, task_id: str, staging_dir: str
    ) -> Optional[str]:
        """
        Always generate a .docx file containing the response text.

        This is an infrastructure-level step (not model-generated).
        Every task produces at least this docx, regardless of whether
        the prompt asked for specific output files.

        Converts markdown-style text to a structured Word document with
        headings, paragraphs, and basic formatting.
        """
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
        except ImportError:
            print("  [DOCX] python-docx not installed, skipping response docx")
            return None

        try:
            doc = Document()

            # Title
            title_para = doc.add_heading(f"Research Report — {task_id}", level=0)

            # Parse markdown-ish text into docx structure
            lines = final_text.split("\n")
            for line in lines:
                stripped = line.strip()
                if not stripped:
                    continue

                # Markdown headings
                if stripped.startswith("# "):
                    doc.add_heading(stripped[2:].strip(), level=1)
                elif stripped.startswith("## "):
                    doc.add_heading(stripped[3:].strip(), level=2)
                elif stripped.startswith("### "):
                    doc.add_heading(stripped[4:].strip(), level=3)
                elif stripped.startswith("#### "):
                    doc.add_heading(stripped[5:].strip(), level=4)
                elif stripped.startswith("- ") or stripped.startswith("* "):
                    # Bullet points
                    doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
                elif re.match(r'^\d+\.\s', stripped):
                    # Numbered lists
                    doc.add_paragraph(
                        re.sub(r'^\d+\.\s', '', stripped).strip(),
                        style="List Number"
                    )
                elif stripped.startswith("|") and stripped.endswith("|"):
                    # Skip markdown table separator lines
                    if all(c in "-| :" for c in stripped):
                        continue
                    doc.add_paragraph(stripped, style="No Spacing")
                elif stripped.startswith("**") and stripped.endswith("**"):
                    # Bold standalone lines
                    p = doc.add_paragraph()
                    run = p.add_run(stripped.strip("*").strip())
                    run.bold = True
                else:
                    doc.add_paragraph(stripped)

            # Footer with metadata
            doc.add_paragraph("")  # spacer
            footer = doc.add_paragraph()
            footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
            run = footer.add_run(
                f"Generated by Qwen DRA | Model: {self.model} | Task: {task_id}"
            )
            run.font.size = Pt(8)
            run.italic = True

            # Save
            filename = f"{task_id}_response.docx"
            filepath = os.path.join(staging_dir, filename)
            doc.save(filepath)
            return filepath

        except Exception as e:
            print(f"  [DOCX] Response docx generation failed: {e}")
            return None

    # ── Dry run / error results ──────────────────────────────────────────

    def _dry_run_result(self, task: ResearchTask, started_at) -> AgentResult:
        completed_at = datetime.now(timezone.utc)
        iat = "CLOSED (web tools omitted)" if task.is_closed else "OPEN"
        fmts = ", ".join(task.output_formats) if task.output_formats else "none"
        files = ", ".join(os.path.basename(f) for f in task.file_paths) or "none"

        report = (
            f"# [DRY RUN] Qwen Deep Research Agent\n\n"
            f"**Task:** {task.task_id}\n"
            f"**Model:** {self.model}\n"
            f"**Base URL:** {self.base_url}\n"
            f"**IAT:** {iat}\n"
            f"**Output formats:** {fmts}\n"
            f"**Files:** {files}\n\n"
            f"## Prompt\n\n{task.prompt[:500]}"
            f"{'...' if len(task.prompt) > 500 else ''}\n"
        )
        print(f"[Qwen] Dry run: {task.task_id} | IAT={iat} | formats={fmts}")

        return AgentResult(
            task_id=task.task_id,
            agent="qwen",
            model=self.model,
            response_text=report,
            citations=[],
            tool_call_log=[],
            input_tokens=0,
            output_tokens=0,
            total_cost_usd=0.0,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=0,
            completed=True,
            forced_stop=False,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    def _error_result(
        self, task, started_at, error_msg,
        input_tokens=0, output_tokens=0,
        tool_call_log=None, iterations=0,
    ) -> AgentResult:
        completed_at = datetime.now(timezone.utc)
        return AgentResult(
            task_id=task.task_id,
            agent="qwen",
            model=self.model,
            response_text="",
            tool_call_log=tool_call_log or [],
            input_tokens=input_tokens,
            output_tokens=output_tokens,
            total_cost_usd=estimate_cost(self.model, input_tokens, output_tokens),
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=iterations,
            completed=False,
            forced_stop=True,
            error=error_msg,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )