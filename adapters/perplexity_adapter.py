"""
perplexity_adapter.py — Black-box search-native adapter for Perplexity Sonar Deep Research.

Perplexity is architecturally different from the other providers: it's
search-native. Every query automatically triggers web search — there's
no "file analysis only" mode. This creates a fundamental tension with
IAT-1 (Closed) tasks where the agent should NOT use external data.

Architecture:
    ┌─────────────────────────────────────────────────┐
    │          PerplexityAdapter                       │
    │                                                  │
    │  1. Extract file text (no native file upload)    │
    │  2. Build prompt with inline file content        │
    │  3. POST /chat/completions (OpenAI-compatible)   │
    │     → Synchronous, blocks 1-3 minutes            │
    │  4. Extract report text + citations array         │
    │  5. Check IAT-1 compliance (unauthorized cites)  │
    │  6. Return AgentResult                           │
    │                                                  │
    │  Structural limitations:                         │
    │  - No file upload API                            │
    │  - No MCP support                                │
    │  - No tool configuration                         │
    │  - Search is ALWAYS on (cannot be disabled)      │
    │  - 128K context window (smallest of 4 providers) │
    │                                                  │
    │  Observability: MINIMAL                          │
    │  - Final report text                             │
    │  - Citation URLs (top-level array)               │
    │  - Token counts                                  │
    │  - No search queries, no reasoning trace         │
    └─────────────────────────────────────────────────┘

IAT-1 (Closed) handling:
    Since Perplexity cannot disable web search, we:
    1. Prepend an explicit "DO NOT use web search" instruction
    2. After receiving the response, check citations for external URLs
       not present in the provided files
    3. Flag any unauthorized external citations in the result metadata

    This is a known structural compromise — documented in the study
    methodology. Perplexity results on IAT-1 tasks should be interpreted
    with this caveat.
"""

from __future__ import annotations

import os
import re
import time
import json
import asyncio
from datetime import datetime, timezone
from typing import Optional
from pathlib import Path

from models import ResearchTask, AgentResult, ToolCall


# ─── Pricing constants ──────────────────────────────────────────────────
# Source: https://docs.perplexity.ai/guides/pricing
# Sonar Deep Research pricing (as of early 2026)

PRICING = {
    "sonar-deep-research": {
        "input_per_mtok": 2.00,
        "output_per_mtok": 8.00,
        # Perplexity also charges per search query (~$5/1000 queries)
        # but this isn't reported in the API response, so we estimate.
        "search_surcharge_per_request": 0.03,  # rough estimate
    },
}


def estimate_cost(model: str, input_tokens: int, output_tokens: int) -> float:
    """
    Estimate USD cost from token counts.

    Note: Perplexity's actual cost includes a per-search-query charge
    that is NOT reflected in token counts. The estimate here covers
    only the token-based portion. Actual cost may be 10-30% higher.
    """
    prices = PRICING.get(model, PRICING["sonar-deep-research"])
    input_cost = (input_tokens / 1_000_000) * prices["input_per_mtok"]
    output_cost = (output_tokens / 1_000_000) * prices["output_per_mtok"]
    surcharge = prices.get("search_surcharge_per_request", 0)
    return round(input_cost + output_cost + surcharge, 6)


# ─── File text extraction (lightweight) ──────────────────────────────────
# Perplexity has no file upload API, so we must extract text and stuff
# it into the prompt. This is the binding constraint — large corpora
# will not fit in Perplexity's 128K context window.

def extract_file_text(fpath: str, max_chars: int = 40_000) -> str:
    """
    Best-effort text extraction for inline inclusion.

    For evaluation fairness, we use the same extraction logic as the
    shared file_processor but without the Claude-specific content block
    formatting. The result is plain text that goes into the prompt.

    Handles: TXT, MD, CSV, HTML, JSON, XML natively.
    For binary formats (PDF, DOCX, XLSX, PPTX), falls back to
    description-only if extraction libraries aren't available.
    """
    ext = Path(fpath).suffix.lower()

    # Plain text formats — read directly
    if ext in (".txt", ".md", ".html", ".htm", ".json", ".xml", ".csv", ".tsv"):
        try:
            with open(fpath, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[... content truncated ...]"
            return text
        except Exception as e:
            return f"[Error reading {os.path.basename(fpath)}: {e}]"

    # DOCX — extract with python-docx if available
    if ext == ".docx":
        try:
            from docx import Document
            doc = Document(fpath)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[... content truncated ...]"
            return text
        except ImportError:
            return f"[File: {os.path.basename(fpath)} — .docx extraction requires python-docx]"
        except Exception as e:
            return f"[Error extracting {os.path.basename(fpath)}: {e}]"

    # XLSX/XLS — extract with openpyxl if available
    if ext in (".xlsx", ".xls"):
        try:
            from openpyxl import load_workbook
            wb = load_workbook(fpath, read_only=True, data_only=True)
            parts = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                parts.append(f"=== Sheet: {sheet_name} ===")
                for row in ws.iter_rows(max_row=500, values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    parts.append("\t".join(cells))
            wb.close()
            text = "\n".join(parts)
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[... content truncated ...]"
            return text
        except ImportError:
            return f"[File: {os.path.basename(fpath)} — .xlsx extraction requires openpyxl]"
        except Exception as e:
            return f"[Error extracting {os.path.basename(fpath)}: {e}]"

    # PDF — extract with pdfplumber if available
    if ext == ".pdf":
        try:
            import pdfplumber
            with pdfplumber.open(fpath) as pdf:
                pages = []
                for i, page in enumerate(pdf.pages):
                    text = page.extract_text() or ""
                    if text.strip():
                        pages.append(f"--- Page {i+1} ---\n{text}")
                    if sum(len(p) for p in pages) > max_chars:
                        pages.append("\n[... remaining pages truncated ...]")
                        break
            return "\n\n".join(pages)
        except ImportError:
            return f"[File: {os.path.basename(fpath)} — .pdf extraction requires pdfplumber]"
        except Exception as e:
            return f"[Error extracting {os.path.basename(fpath)}: {e}]"

    # PPTX — extract with python-pptx if available
    if ext == ".pptx":
        try:
            from pptx import Presentation
            prs = Presentation(fpath)
            parts = []
            for i, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            if para.text.strip():
                                slide_texts.append(para.text)
                if slide_texts:
                    parts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))
            text = "\n\n".join(parts)
            if len(text) > max_chars:
                text = text[:max_chars] + "\n\n[... content truncated ...]"
            return text
        except ImportError:
            return f"[File: {os.path.basename(fpath)} — .pptx extraction requires python-pptx]"
        except Exception as e:
            return f"[Error extracting {os.path.basename(fpath)}: {e}]"

    # Images — cannot inline image data into Perplexity
    if ext in (".jpg", ".jpeg", ".png", ".gif", ".webp"):
        return (
            f"[File: {os.path.basename(fpath)} — Image file. "
            f"Perplexity does not support image input. "
            f"Content not available for this agent.]"
        )

    return f"[File: {os.path.basename(fpath)} — Unsupported format: {ext}]"


# ─── IAT-1 compliance checking ──────────────────────────────────────────

# Prefix added to IAT-1 (Closed) task prompts
IAT1_INSTRUCTION = """
CRITICAL CONSTRAINT — CLOSED INFORMATION TASK:
You MUST answer this question using ONLY the information provided in the
files below. Do NOT use web search results, external databases, or any
information source other than the provided documents.

If you cannot answer a question from the provided files alone, explicitly
state "Insufficient information in provided documents" rather than
searching externally.

Any claims you make MUST be traceable to specific content in the provided
files. Do NOT supplement with external knowledge.

"""


def check_iat1_compliance(
    response_text: str,
    citations: list[dict],
    file_names: list[str],
) -> dict:
    """
    Check whether a Perplexity response on an IAT-1 (Closed) task
    contains unauthorized external citations.

    Returns a compliance report:
      {
        "compliant": bool,
        "total_citations": int,
        "external_citations": int,
        "external_urls": [str],
        "compliance_ratio": float,  # 1.0 = fully compliant
      }

    This is imperfect — Perplexity may use external knowledge without
    citing it, which we cannot detect. But unauthorized citations are
    a clear signal of non-compliance.
    """
    external_urls = []
    file_names_lower = [fn.lower() for fn in file_names]

    for cite in citations:
        url = cite.get("url", "")
        title = cite.get("title", "").lower()

        # Check if the citation references one of the provided files
        is_file_ref = any(fn in title or fn in url.lower() for fn in file_names_lower)

        if not is_file_ref and url:
            external_urls.append(url)

    total = len(citations)
    external = len(external_urls)

    return {
        "compliant": external == 0,
        "total_citations": total,
        "external_citations": external,
        "external_urls": external_urls,
        "compliance_ratio": 1.0 if total == 0 else (total - external) / total,
    }


# ─── The adapter ──────────────────────────────────────────────────────────

class PerplexityAdapter:
    """
    Runs a deep research task using Perplexity's Sonar Deep Research API.

    This adapter handles:
      1. File text extraction and prompt construction
      2. IAT-1 instruction injection for Closed tasks
      3. Synchronous API call (OpenAI-compatible endpoint)
      4. Citation extraction from the response
      5. IAT-1 compliance checking

    Usage:
        adapter = PerplexityAdapter(api_key="pplx-...")
        result = await adapter.run(task)

    Dry-run mode:
        adapter = PerplexityAdapter(dry_run=True)
        result = await adapter.run(task)
    """

    def __init__(
        self,
        api_key: Optional[str] = None,
        model: str = "sonar-deep-research",
        dry_run: bool = False,
    ):
        self.api_key = api_key or os.environ.get("PERPLEXITY_API_KEY")
        self.model = model
        self.dry_run = dry_run

        if not dry_run:
            try:
                from openai import AsyncOpenAI
                self.client = AsyncOpenAI(
                    api_key=self.api_key,
                    base_url="https://api.perplexity.ai",
                    timeout=300,  # 5 min — Perplexity DR can be slow
                )
            except ImportError:
                raise ImportError(
                    "pip install openai  "
                    "(Perplexity uses OpenAI-compatible API)"
                )

    async def run(self, task: ResearchTask) -> AgentResult:
        """Execute a deep research task via Perplexity's API."""
        started_at = datetime.now(timezone.utc)

        print(f"[Perplexity] Starting task {task.task_id} "
              f"(model={self.model}, dry_run={self.dry_run})")

        if self.dry_run:
            return self._dry_run_result(task, started_at)

        try:
            # ── Step 1: Build prompt with inline file content ───────
            messages = self._build_messages(task)
            prompt_tokens_est = sum(len(m["content"]) // 4 for m in messages)
            print(f"[Perplexity] Prompt built: ~{prompt_tokens_est} tokens est.")

            if prompt_tokens_est > 120_000:
                print(f"[Perplexity] WARNING: Prompt likely exceeds 128K context. "
                      f"Response quality may degrade.")

            # ── Step 2: Call Perplexity API ─────────────────────────
            # Synchronous call — blocks until complete (1-3 min typical)
            print(f"[Perplexity] Calling API (synchronous, may take 1-3 min)...")

            response = await self.client.chat.completions.create(
                model=self.model,
                messages=messages,
            )

            # ── Step 3: Extract results ────────────────────────────
            return self._extract_result(response, task, started_at)

        except Exception as e:
            print(f"[Perplexity] Error: {e}")
            return self._error_result(task, started_at, str(e))

    # ─── Message construction ─────────────────────────────────────────

    def _build_messages(self, task: ResearchTask) -> list[dict]:
        """
        Build the messages payload for Perplexity's Chat API.

        Since Perplexity has no file upload, we extract text from all
        files and inline them into the user message. This is the
        binding constraint — 128K context means large corpora will
        be truncated.

        For IAT-1 (Closed) tasks, we prepend an explicit instruction
        not to use web search.
        """
        # System message: research methodology
        system_msg = (
            "You are a deep research analyst. Produce a comprehensive, "
            "well-structured report based on the user's research question. "
            "Cite all claims with numbered references. Include an executive "
            "summary and limitations section."
        )

        # Build user message with inline file content
        user_parts = []

        # IAT-1 injection
        if task.is_closed:
            user_parts.append(IAT1_INSTRUCTION)
            print(f"[Perplexity] IAT-1 (Closed): injecting DO NOT SEARCH instruction")

        # Inline file content
        if task.file_paths:
            user_parts.append("=" * 60)
            user_parts.append("PROVIDED REFERENCE DOCUMENTS")
            user_parts.append("=" * 60)

            # Budget: distribute 128K context across files + prompt
            # Reserve ~30K tokens for the prompt + response
            remaining_budget = 90_000 * 4  # ~90K tokens in chars
            per_file_budget = remaining_budget // max(len(task.file_paths), 1)

            for fpath in task.file_paths:
                fname = os.path.basename(fpath)
                user_parts.append(f"\n{'─' * 40}")
                user_parts.append(f"FILE: {fname}")
                user_parts.append(f"{'─' * 40}")

                file_text = extract_file_text(fpath, max_chars=per_file_budget)
                user_parts.append(file_text)

            user_parts.append("\n" + "=" * 60)
            user_parts.append("END OF REFERENCE DOCUMENTS")
            user_parts.append("=" * 60 + "\n")

        # The actual research prompt
        user_parts.append(task.prompt)

        return [
            {"role": "system", "content": system_msg},
            {"role": "user", "content": "\n".join(user_parts)},
        ]

    # ─── Result extraction ────────────────────────────────────────────

    def _extract_result(self, response, task, started_at) -> AgentResult:
        """
        Extract results from Perplexity's Chat Completions response.

        Perplexity returns:
          - choices[0].message.content: the report text
          - citations: top-level array of URLs (unique to Perplexity)
          - usage: {prompt_tokens, completion_tokens, total_tokens}
        """
        completed_at = datetime.now(timezone.utc)

        # Extract report text
        report_text = ""
        if response.choices:
            report_text = response.choices[0].message.content or ""

        # Extract citations — Perplexity returns these as a top-level array
        raw_citations = getattr(response, 'citations', []) or []
        citations = []
        for cite in raw_citations:
            if isinstance(cite, str):
                citations.append({"url": cite, "title": "", "snippet": ""})
            elif isinstance(cite, dict):
                citations.append({
                    "url": cite.get("url", ""),
                    "title": cite.get("title", ""),
                    "snippet": cite.get("snippet", "")[:200],
                })

        # Token usage
        input_tokens = 0
        output_tokens = 0
        if hasattr(response, 'usage') and response.usage:
            input_tokens = getattr(response.usage, 'prompt_tokens', 0)
            output_tokens = getattr(response.usage, 'completion_tokens', 0)

        total_cost = estimate_cost(self.model, input_tokens, output_tokens)

        # ── IAT-1 compliance check ──────────────────────────────────
        iat1_report = None
        if task.is_closed:
            file_names = [os.path.basename(f) for f in task.file_paths]
            iat1_report = check_iat1_compliance(
                report_text, citations, file_names
            )
            if iat1_report["compliant"]:
                print(f"[Perplexity] IAT-1 compliance: PASSED "
                      f"(no external citations detected)")
            else:
                print(f"[Perplexity] IAT-1 compliance: FAILED "
                      f"({iat1_report['external_citations']} external citations)")
                for url in iat1_report["external_urls"][:5]:
                    print(f"  → {url}")

        # Build metadata string for compliance info
        error_msg = None
        if iat1_report and not iat1_report["compliant"]:
            error_msg = (
                f"IAT-1 VIOLATION: {iat1_report['external_citations']} "
                f"unauthorized external citations detected. "
                f"Compliance ratio: {iat1_report['compliance_ratio']:.1%}"
            )

        print(f"[Perplexity] Done. "
              f"tokens=({input_tokens}in, {output_tokens}out), "
              f"cost=${total_cost:.2f}, "
              f"citations={len(citations)}")

        return AgentResult(
            task_id=task.task_id,
            agent="perplexity",
            model=self.model,
            response_text=report_text,
            citations=citations,
            tool_call_log=[],  # No observability into search behavior
            input_tokens=input_tokens,
            output_tokens=output_tokens,
            total_cost_usd=total_cost,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=1,  # Single synchronous call
            completed=True,
            forced_stop=False,
            error=error_msg,  # IAT-1 violation info if applicable
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    # ─── Dry run ──────────────────────────────────────────────────────

    def _dry_run_result(self, task: ResearchTask, started_at) -> AgentResult:
        """Return a mock result without making any API calls."""
        completed_at = datetime.now(timezone.utc)

        is_closed = task.is_closed
        iat_status = "CLOSED (instruction-only enforcement)" if is_closed else "OPEN"
        file_summary = ", ".join(os.path.basename(f) for f in task.file_paths) or "none"

        # Simulate IAT-1 compliance check
        mock_compliance = ""
        mock_error = None
        if is_closed:
            mock_compliance = (
                "\n## IAT-1 Compliance Check (Simulated)\n\n"
                "In production, the adapter checks all citations against\n"
                "the provided file names. Any URL not referencing the source\n"
                "documents is flagged as an unauthorized external citation.\n\n"
                "**Structural caveat:** Perplexity CANNOT disable web search.\n"
                "IAT-1 enforcement relies on prompt instruction + post-hoc\n"
                "compliance checking. Results should be interpreted with this\n"
                "limitation in mind.\n"
            )
            mock_error = (
                "IAT-1 STRUCTURAL WARNING: Perplexity cannot disable web search. "
                "Compliance enforced via instruction injection + citation audit only."
            )

        # Estimate context budget usage
        total_file_chars = 0
        for fpath in task.file_paths:
            if os.path.exists(fpath):
                total_file_chars += os.path.getsize(fpath)
        est_tokens = total_file_chars // 4
        context_pct = (est_tokens / 128_000) * 100

        report = (
            f"# [DRY RUN] Perplexity Deep Research Report\n\n"
            f"**Task:** {task.task_id}\n"
            f"**Model:** {self.model}\n"
            f"**IAT:** {iat_status}\n"
            f"**Research Type:** {task.research_type or 'unspecified'}\n"
            f"**Context Window:** 128K tokens (smallest of 4 providers)\n"
            f"**Files:** {file_summary}\n"
            f"**Est. File Context Usage:** ~{est_tokens:,} tokens ({context_pct:.0f}% of 128K)\n\n"
            f"## Executive Summary\n\n"
            f"This is a dry-run mock response. In production, Perplexity's "
            f"sonar-deep-research model would:\n"
            f"1. Receive file text inlined in the prompt (no file upload API)\n"
            f"2. Run its search-native research loop (always searches web)\n"
            f"3. Return a report with a top-level citations array\n\n"
            f"## Structural Limitations\n\n"
            f"Perplexity is the most constrained agent for this evaluation:\n"
            f"- No file upload API (text must be inlined in prompt)\n"
            f"- No MCP support (no Tier 2 large-corpus access)\n"
            f"- 128K context window (vs. 200K Claude, 200K OpenAI, 1M Gemini)\n"
            f"- Web search CANNOT be disabled (IAT-1 structurally compromised)\n"
            f"- No image/vision support (image files are skipped)\n\n"
            f"Perplexity's strength is on IAT-3 (Open) tasks requiring\n"
            f"comprehensive web research with real-time data.\n"
            f"{mock_compliance}\n"
            f"## Prompt Received\n\n"
            f"{task.prompt[:500]}{'...' if len(task.prompt) > 500 else ''}\n\n"
            f"## Limitations\n\n"
            f"Dry-run mode — no actual research performed.\n"
        )

        print(f"[Perplexity] Dry run complete for {task.task_id}")
        print(f"[Perplexity]   IAT: {iat_status}")
        print(f"[Perplexity]   Files: {len(task.file_paths)} "
              f"(~{est_tokens:,} tokens, {context_pct:.0f}% of 128K)")
        print(f"[Perplexity]   Search: ALWAYS ON (cannot be disabled)")

        return AgentResult(
            task_id=task.task_id,
            agent="perplexity",
            model=self.model,
            response_text=report,
            citations=[{"url": "https://example.com/mock", "title": "Mock Citation", "snippet": "Dry run"}],
            tool_call_log=[],
            input_tokens=0,
            output_tokens=0,
            total_cost_usd=0.0,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=1,
            completed=True,
            forced_stop=False,
            error=mock_error,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )

    def _error_result(self, task, started_at, error_msg) -> AgentResult:
        """Construct an error result when the API call fails."""
        completed_at = datetime.now(timezone.utc)
        return AgentResult(
            task_id=task.task_id,
            agent="perplexity",
            model=self.model,
            response_text="",
            tool_call_log=[],
            input_tokens=0,
            output_tokens=0,
            total_cost_usd=0.0,
            total_duration_sec=(completed_at - started_at).total_seconds(),
            iterations=0,
            completed=False,
            forced_stop=True,
            error=error_msg,
            started_at=started_at.isoformat(),
            completed_at=completed_at.isoformat(),
        )
